"""
Daniel Kordmodanlou - CV PowerPoint Builder
Creates a polished, modern CV presentation in TWO visual themes (same content).

  python build_cv.py            # midnight  (dark "Midnight Executive")
  python build_cv.py midnight   #  → Daniel's CV (1).pptx
  python build_cv.py memphis    # bright playful "Memphis" (Slidesgo-style)
  python build_cv.py both       #  → also Daniel's CV - Memphis.pptx

Every slide is drawn through semantic theme tokens (apply_theme), so the same
builders render either look. "midnight" = deep navy + amber/gold, ivory cards.
"memphis" = white background, coral/cyan/yellow/purple, Lexend Deca + Questrial
fonts, big corner circles + dot clusters, outlined cards. There are NO
animations in this deck.
"""
import os
import sys
import struct
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# --- Raw palette: Midnight Executive ---
NAVY     = RGBColor(0x0F, 0x17, 0x2A)   # deep navy (main bg accents)
INDIGO   = RGBColor(0x1E, 0x2A, 0x55)   # mid navy / cards
ROYAL    = RGBColor(0x2E, 0x42, 0x82)   # royal blue (secondary)
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)   # accent
GOLD     = RGBColor(0xFB, 0xBF, 0x24)   # accent light
IVORY    = RGBColor(0xF8, 0xFA, 0xFC)   # near-white
SLATE    = RGBColor(0x94, 0xA3, 0xB8)   # body muted
DARK_TEXT = RGBColor(0x1F, 0x29, 0x37)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xE2, 0xE8, 0xF0)
TEAL     = RGBColor(0x14, 0xB8, 0xA6)
MAGENTA  = RGBColor(0xEC, 0x4D, 0x88)

# --- Raw palette: MEMPHIS (bright / playful, from "Daniel's CV.pptx") ---
M_CORAL  = RGBColor(0xEC, 0x66, 0x3D)   # dominant accent
M_CYAN   = RGBColor(0x78, 0xDD, 0xF4)   # light cyan — underline / pills / dots
M_TEAL   = RGBColor(0x1C, 0x92, 0xBD)   # deeper cyan — readable accent text
M_YELLOW = RGBColor(0xF7, 0xB8, 0x45)   # decorative only
M_PURPLE = RGBColor(0x67, 0x58, 0x9D)
M_GREEN  = RGBColor(0x69, 0xD9, 0x71)   # decorative only
M_PINK   = RGBColor(0xEF, 0xA1, 0xBD)   # decorative only
M_INK    = RGBColor(0x1C, 0x1C, 0x1C)
M_GREY_D = RGBColor(0x4A, 0x4A, 0x4A)
M_GREY_M = RGBColor(0x80, 0x80, 0x80)
M_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
M_CARD   = RGBColor(0xF1, 0xF3, 0xF5)   # near-white card surface
M_LAV    = RGBColor(0xCF, 0xC7, 0xE6)   # muted text on purple panel

MEDIA = "extracted/ppt/media"

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Decorative shape colours used by the Memphis corner/edge ornaments.
_DECO_PAL = [M_CORAL, M_CYAN, M_YELLOW, M_PURPLE, M_GREEN, M_PINK]

# Output filename per theme.
OUT_FILE = {
    "midnight": "Daniel's CV (1).pptx",
    "memphis":  "Daniel's CV - Memphis.pptx",
}


# ====================================================================
# THEME  ·  semantic tokens, rebound per theme by apply_theme()
# ====================================================================
# Builders never name a raw colour; they use these tokens. apply_theme()
# swaps them so one set of builders renders either visual style.
def apply_theme(name):
    midnight = dict(
        H_FONT="Calibri", B_FONT="Calibri Light",
        PAGE=IVORY, FEAT=NAVY, PANEL=NAVY, FRAME=NAVY,
        INK=NAVY, INK_F=IVORY, ON_DARK=IVORY,
        SUB=ROYAL, SUB_F=GOLD, BODY=DARK_TEXT, MUTE=SLATE, MUTE_D=SLATE,
        CARD=NAVY, CARD2=INDIGO, CHEAD=IVORY, CBODY=LIGHT_GREY,
        ACC1=AMBER, ACC2=GOLD, ACC3=TEAL, EYE=AMBER, PANEL_ACC=GOLD,
        CHIP_BG=INDIGO, CHIP_FG=GOLD,
        STACKHDR=INDIGO, STACKCHIP_BG=INDIGO, STACKCHIP_FG=IVORY,
        PILL_BG=INDIGO, PILL_FG=GOLD, PILL2_BG=ROYAL, PILL2_FG=IVORY,
        FOOT=SLATE, UNDER=None,
        BAR_TRACK=RGBColor(0x3B, 0x4E, 0x84), WMARK=RGBColor(0x21, 0x2E, 0x55),
        DIVLINE=LIGHT_GREY,
        CARD_LINE="none", MARKER="bar", DECO="none",
    )
    memphis = dict(
        H_FONT="Lexend Deca", B_FONT="Questrial",
        PAGE=M_WHITE, FEAT=M_WHITE, PANEL=M_PURPLE, FRAME=M_CARD,
        INK=M_INK, INK_F=M_INK, ON_DARK=M_WHITE,
        SUB=M_GREY_D, SUB_F=M_GREY_D, BODY=M_INK, MUTE=M_GREY_M, MUTE_D=M_LAV,
        CARD=M_CARD, CARD2=M_CARD, CHEAD=M_INK, CBODY=M_GREY_D,
        ACC1=M_CORAL, ACC2=M_PURPLE, ACC3=M_TEAL, EYE=M_CORAL, PANEL_ACC=M_CYAN,
        CHIP_BG=M_PURPLE, CHIP_FG=M_WHITE,
        STACKHDR=M_CORAL, STACKCHIP_BG=M_PURPLE, STACKCHIP_FG=M_WHITE,
        PILL_BG=M_CYAN, PILL_FG=M_INK, PILL2_BG=M_PURPLE, PILL2_FG=M_WHITE,
        FOOT=M_GREY_M, UNDER=M_CYAN,
        BAR_TRACK=RGBColor(0xE6, 0xE8, 0xEC), WMARK=RGBColor(0xF0, 0xEC, 0xF8),
        DIVLINE=RGBColor(0xE5, 0xE5, 0xE5),
        CARD_LINE="accent", MARKER="dot", DECO="memphis",
    )
    chosen = memphis if name == "memphis" else midnight
    globals().update(chosen)
    globals()["THEME_NAME"] = name


# Initialise to midnight so module-level defaults resolve before build().
apply_theme("midnight")


# ====================================================================
# ROLE MODES  ·  tailor the CV's wording to a target role
# ====================================================================
# `python build_cv.py [midnight|memphis|both] [general|backend]`. Default
# role=general keeps the content written throughout this file. A role profile
# overrides the positioning, the per-role descriptions, the skills and which
# project slides show — everything else falls back via rp().
ROLE = "general"

ROLE_PROFILES = {
    "general": {},
    "backend": {
        "suffix": "Backend",
        "eyebrow": "PORTFOLIO  ·  BACKEND SOFTWARE ENGINEER  ·  2026",
        "title_role": "Backend Software Engineer  ·  APIs · Distributed Systems · Cloud",
        "about_sub": "I build reliable, scalable backend systems — APIs, data and cloud infrastructure.",
        "about_bullets": [
            "[b]Backend & distributed-systems[/b] engineer — REST / gRPC services in [b]FastAPI, Spring Boot, Django[/b] over [b]PostgreSQL / MySQL / Redis[/b], containerized on [b]Docker / Kubernetes / AWS[/b].",
            "Shipped a [b]URL-shortener SaaS[/b] (Java microservices on K8s + AWS with a Hadoop redundancy layer) and a [b]from-scratch search engine[/b] (inverted index, tf-idf, KNN).",
            "Built the [b]FastAPI + LangGraph[/b] backends for Guardian (multi-agent, idempotent 911 tooling) and CaseLogic (hybrid Chroma + SQLite FTS5 retrieval).",
            "Strong CS fundamentals: [b]XV6 kernel[/b] syscalls + scheduling, [b]multi-core C[/b] (OpenMP / CUDA), and CI/CD on .NET / Azure / GitHub Actions.",
        ],
        "exp_sub": "From IoT networking to cloud microservices and ML platforms — a track of shipping reliable backends.",
        "exp_overview": {
            "Machine Learning Associate": "Built the FastAPI + Ollama/LangChain services and RAG pipeline (FAISS) behind a real-time VR coaching system.",
            "Senior Unity Developer": "Built CI/CD delivery, Docker workflows and Claude-assisted tooling that auto-investigates exceptions and commits fixes.",
            "Research & Teaching Assistant": "Built reproducible test tooling and motion-capture data pipelines; root-caused an SDK bug. TA for OOP (Java).",
            "Software Developer": "Built simulation backends with fault-tolerant pathfinding and real-time monitoring dashboards (DOTS/ECS).",
            "Lead Unity3D Developer": "Led the online-multiplayer backend: Photon + PlayFab (Azure) matchmaking, Firebase analytics, DQN + MCTS AI.",
            "Volunteer Full-Stack Developer": "Maintained & optimized a Django backend (+ Vue.js) for a 100+ country student-exchange platform.",
            "Software Engineer Intern": "Built a real-time IoT multiplayer backend — WebSockets / gRPC / Protobuf / MQTT.",
            "Technical Staff & TA": "TA for Operating Systems, C, Advanced Java and AI. GameCraft events.",
        },
        "exp": {
            "vector": {
                "summary": "Built the backend services and ML data pipeline for a real-time VR firefighter training system — FastAPI, Ollama/LangChain, RAG over FAISS.",
                "bullets": [
                    "Engineered the [b]LLM Coaching Layer[/b] as a service: [b]FastAPI[/b] + [b]Ollama / LangChain[/b] with a [b]RAG[/b] retrieval pipeline (FAISS similarity) over training manuals.",
                    "Architected a [b]Deviation Engine[/b] over sliding windows of 3D skeletal telemetry — MPJPE and quaternion distances scored against training standards.",
                    "Built a full-body [b]data-collection pipeline[/b] mapping headset telemetry to the [b]NTU-25[/b] skeleton with cross-platform frame normalization.",
                    "Collaborated on structured evaluation of generative-AI pipelines.",
                ],
                "tags": ["FastAPI", "RAG", "FAISS", "Ollama", "LangChain", "NTU-25", "Python", "Data Pipeline"],
            },
            "dreamforge": {
                "summary": "Core contributor to a C# + Python engine. Built CI/CD delivery, Docker workflows and Claude-assisted tooling that auto-resolves build failures.",
                "bullets": [
                    "Designed [b]CI/CD delivery components[/b] and [b]Docker-based workflows[/b] in [b].NET / Azure / GitHub Actions[/b] for automated deployment and exception analysis.",
                    "Built [b]AI-assisted tooling (Claude)[/b] that auto-intercepts exceptions, investigates root causes across C# / Python / Docker layers, and commits permanent fixes.",
                    "Performed deep code reviews; led peers in resolving structural inefficiencies and maintaining robust debugging workflows.",
                    "Focus on reliable tooling, procedural generation and continuous delivery.",
                ],
                "tags": ["C#", ".NET", "Python", "Docker", "Azure", "GitHub Actions", "CI/CD", "Claude"],
            },
            "biomotion": {
                "summary": "Built reproducible test tooling and motion-capture data pipelines; root-caused a critical SDK bug. TA for OOP (Java), 5 semesters.",
                "bullets": [
                    "Investigated and reported a critical [b]stereo-geometry bug[/b] in Meta XR SDK Unity; shipped a [b]reproducible test project[/b] documenting the root cause.",
                    "Built [b]data pipelines[/b] for VR + motion-capture experiments on Meta Quest / ARKit / Unity / Unreal.",
                    "Presented research at the [b]47th European Conference on Visual Perception (ECVP 2025)[/b], Mainz.",
                    "Taught Object-Oriented Programming ([b]Java[/b]) for 5 semesters — design patterns and systems thinking.",
                ],
                "tags": ["Java", "Test Harness", "Data Pipeline", "Meta Quest", "ARKit", "ECVP 2025"],
            },
            "tectotrack": {
                "summary": "Built simulation backends for high-traffic environments with fault-tolerant pathfinding and real-time monitoring dashboards.",
                "bullets": [
                    "Built [b]simulation frameworks[/b] for high-traffic environments with [b]fault-tolerant pathfinding[/b] for adaptive crowd movement.",
                    "Developed [b]real-time monitoring dashboards[/b] for performance, agent state and logs.",
                    "Built internal tooling for debugging, visualization and flow adjustment.",
                    "Optimized high-density scenes with [b]ECS[/b] (Unity DOTS).",
                ],
                "tags": ["C#", "DOTS / ECS", "Pathfinding", "Dashboards", "Simulation", "Monitoring"],
            },
            "techu": {
                "summary": "Led the online-multiplayer backend: Photon + PlayFab (Azure) matchmaking, Firebase analytics, and AI opponents (DQN + MCTS).",
                "bullets": [
                    "Engineered [b]online multiplayer[/b] via [b]Photon[/b] + [b]PlayFab (Azure)[/b] for matchmaking and session management.",
                    "Integrated [b]Firebase[/b] for real-time crash analytics and diagnostics.",
                    "Built AI opponents with [b]DQN[/b] + [b]Monte Carlo Tree Search[/b].",
                    "Owned the full pipeline: gameplay, backend, UI and store releases.",
                ],
                "tags": ["Photon", "PlayFab", "Azure", "Firebase", "DQN", "MCTS", "C#"],
            },
            "iaeste": {
                "summary": "Maintained and optimized a Django backend (with a Vue.js frontend) for a student-exchange platform serving 100+ countries.",
                "bullets": [
                    "Maintained the platform backend in [b]Django[/b] (with a [b]Vue.js[/b] frontend).",
                    "Optimized backend routes and reduced latency on key endpoints.",
                    "Improved UI/UX for student-exchange applications.",
                    "Contributed long-term maintainability through refactoring and tests.",
                ],
                "tags": ["Django", "Python", "Vue.js", "REST", "Optimization"],
            },
            "sepantab": {
                "summary": "Built a real-time IoT multiplayer backend with resilient distributed networking — gRPC, Protocol Buffers, MQTT, WebSockets.",
                "bullets": [
                    "Implemented real-time networked messaging using [b]WebSockets[/b].",
                    "Built resilient distributed communication with [b]gRPC[/b] + [b]Protocol Buffers[/b] + [b]MQTT[/b].",
                    "Designed low-latency state-sync networking infrastructure in [b]C#[/b].",
                    "Delivered an IoT-based WebGL multiplayer experience.",
                ],
                "tags": ["gRPC", "Protobuf", "MQTT", "WebSockets", "C#", "IoT", "Real-time"],
            },
        },
        "skills": [
            ("Core",
             ["Java", "Python", "C#", "C++", "Go", "C", "SQL"]),
            ("Backend & Data",
             ["FastAPI", "Spring Boot", "Django", "gRPC", "REST", "PostgreSQL", "MySQL", "Redis",
              "pgvector", "Chroma", "SQLite FTS5", "MQTT", "Protobuf"]),
            ("Infra & Tooling",
             ["Docker", "Kubernetes", "AWS (ECR)", "Azure", ".NET", "GitHub Actions", "Nginx",
              "Hadoop", "OpenMP", "CUDA", "Firebase", "PyTorch"]),
        ],
        "proj_div_sub": "Backend services, data systems & low-level work — and open source.",
        "hide_proj": ["minigames", "hypervigilance", "gaugan"],
    },
}


def apply_role(name):
    global ROLE
    ROLE = name if name in ROLE_PROFILES else "general"


def rp(key, default=None):
    """Return the active role's override for `key`, or `default`."""
    return ROLE_PROFILES.get(ROLE, {}).get(key, default)


# Deck state — (re)created per build by _reset_deck().
prs = None
blank_layout = None
_deco_i = [0]


def _reset_deck():
    """Fresh Presentation + per-build state (lets us render both themes)."""
    global prs, blank_layout, _deco_i
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]
    _deco_i = [0]


# --------------------- helpers ---------------------

def add_rect(slide, x, y, w, h, fill, line=None, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    if not shadow:
        shp.shadow.inherit = False
        # disable shadow via xml
        spPr = shp.fill._xPr
        # remove existing effectLst if any
        for el in spPr.findall(qn('a:effectLst')):
            spPr.remove(el)
        effectLst = etree.SubElement(spPr, qn('a:effectLst'))
    return shp


def add_round_rect(slide, x, y, w, h, fill, radius=0.08, line=None, line_w=1.25):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    # set corner radius via adj
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    # remove shadow
    spPr = shp.fill._xPr
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))
    return shp


def add_oval(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    spPr = shp.fill._xPr
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))
    return shp


def add_text(slide, x, y, w, h, text, *,
             font=None, size=14, bold=False, color=None,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, line_spacing=None):
    if font is None:
        font = H_FONT
    if color is None:
        color = BODY
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, bullets, *,
                font=None, size=14, color=None, line_spacing=1.2,
                bullet_char="•", spacing_after=4):
    if font is None:
        font = B_FONT
    if color is None:
        color = BODY
    tb = slide.shapes.add_textbox(x, y, w, h)
    try:
        tb.name = "BULLETS"          # tag so the entrance reveals it per bullet
    except Exception:
        pass
    tf = tb.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(spacing_after)
        # handle bold tokens [b]text[/b]
        parts = _split_bold(b)
        first = True
        if bullet_char:
            run = p.add_run()
            run.text = bullet_char + "  "
            run.font.name = font
            run.font.size = Pt(size)
            run.font.color.rgb = color
            first = False
        for txt, is_bold in parts:
            run = p.add_run()
            run.text = txt
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = is_bold
            run.font.color.rgb = color
    return tb


def _split_bold(s):
    """Parse markers [b]...[/b] inline into list of (text, bold)."""
    out = []
    i = 0
    while i < len(s):
        start = s.find("[b]", i)
        if start == -1:
            out.append((s[i:], False))
            break
        if start > i:
            out.append((s[i:start], False))
        end = s.find("[/b]", start)
        if end == -1:
            out.append((s[start:], False))
            break
        out.append((s[start+3:end], True))
        i = end + 4
    return out


def fill_slide(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    spPr = bg.fill._xPr
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))
    return bg


_deco_i  # noqa: F821 — module global, (re)set by _reset_deck()


def page_bg(s, feature=False):
    """Fill the slide background for the active theme and (memphis) scatter
    playful corner ornaments that stay clear of the content area."""
    fill_slide(s, FEAT if feature else PAGE)
    if DECO == "memphis":
        idx = _deco_i[0]; _deco_i[0] += 1
        a = _DECO_PAL[idx % len(_DECO_PAL)]
        b = _DECO_PAL[(idx + 3) % len(_DECO_PAL)]
        # big quarter-circle hugging the (empty) top-right corner
        add_oval(s, Inches(11.85), Inches(-1.05), Inches(2.6), Inches(2.6), a)
        add_oval(s, Inches(11.35), Inches(1.78), Inches(0.26), Inches(0.26), b)
        add_oval(s, Inches(12.98), Inches(1.46), Inches(0.16), Inches(0.16), M_CYAN)


def cline(accent):
    """Outline colour for a card: the accent (memphis) or none (midnight)."""
    return accent if CARD_LINE == "accent" else None


def add_section_marker(slide, num_text, color=AMBER):
    """Section eyebrow marker: amber bar (midnight) or coral dot (memphis),
    followed by the label. Theme-aware: ``color`` is ignored in favour of EYE
    so every call site updates with the active theme."""
    if MARKER == "dot":
        add_oval(slide, Inches(0.5), Inches(0.54), Inches(0.22), Inches(0.22), EYE)
        lx = Inches(0.82); lw = Inches(8.5)
    else:
        # Midnight: keep the original geometry (label box width 3in, x=0.7).
        add_rect(slide, Inches(0.5), Inches(0.5), Inches(0.06), Inches(0.5), EYE)
        lx = Inches(0.7); lw = Inches(3)
    add_text(slide, lx, Inches(0.5), lw, Inches(0.5),
             num_text, font=H_FONT, size=11, bold=True, color=EYE,
             anchor=MSO_ANCHOR.MIDDLE)


def _png_ratio(path):
    """Return height/width of a PNG from its IHDR (no Pillow needed)."""
    with open(path, "rb") as f:
        w, h = struct.unpack(">II", f.read(24)[16:24])
    return h / float(w)


def add_image(slide, path, x, y, w=None, h=None):
    if not os.path.exists(path):
        print(f"WARNING: image not found {path}")
        return None
    if w and h:
        return slide.shapes.add_picture(path, x, y, w, h)
    elif w:
        return slide.shapes.add_picture(path, x, y, width=w)
    elif h:
        return slide.shapes.add_picture(path, x, y, height=h)
    return slide.shapes.add_picture(path, x, y)


def page_footer(slide, label, num):
    # subtle footer with name and slide number (themed text colour)
    add_text(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.4),
             "Danial Kordmodanlou  ·  CV 2026", font=B_FONT, size=9,
             color=FOOT)
    add_text(slide, Inches(11.8), Inches(7.05), Inches(1.1), Inches(0.4),
             f"{num:02d}  /  {label}", font=B_FONT, size=9, color=FOOT,
             align=PP_ALIGN.RIGHT).name = "FOOTNUM"


def light_footer(slide, label, num):
    add_text(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.4),
             "Danial Kordmodanlou  ·  CV 2026", font=B_FONT, size=9,
             color=FOOT)
    add_text(slide, Inches(11.8), Inches(7.05), Inches(1.1), Inches(0.4),
             f"{num:02d}  /  {label}", font=B_FONT, size=9,
             color=FOOT,
             align=PP_ALIGN.RIGHT).name = "FOOTNUM"


# ====================================================================
# SLIDE 1 — TITLE
# ====================================================================
def slide_title():
    s = prs.slides.add_slide(blank_layout)
    if DECO == "memphis":
        fill_slide(s, M_WHITE)
        # playful geometric composition
        add_oval(s, Inches(-1.2), Inches(-1.2), Inches(3.0), Inches(3.0), M_CORAL)
        add_oval(s, Inches(1.75), Inches(0.35), Inches(0.42), Inches(0.42), M_CYAN)
        add_oval(s, Inches(11.3), Inches(0.5), Inches(1.5), Inches(1.5), M_YELLOW)
        add_oval(s, Inches(12.75), Inches(2.25), Inches(0.3), Inches(0.3), M_CYAN)
        add_oval(s, Inches(11.5), Inches(5.9), Inches(2.0), Inches(2.0), M_PURPLE)
        add_oval(s, Inches(10.8), Inches(6.7), Inches(0.3), Inches(0.3), M_CORAL)
        add_oval(s, Inches(0.6), Inches(6.75), Inches(0.34), Inches(0.34), M_PINK)
        # content sits below the top-left corner circle so nothing is covered
        add_text(s, Inches(1.0), Inches(2.05), Inches(11), Inches(0.4),
                 rp("eyebrow", "PORTFOLIO  ·  CURRICULUM VITAE  ·  2026"),
                 font=H_FONT, size=12, bold=True, color=M_CORAL)
        add_text(s, Inches(1.0), Inches(2.5), Inches(12), Inches(1.5),
                 "Danial Kordmodanlou", font=H_FONT, size=60, bold=True, color=M_INK)
        add_rect(s, Inches(1.05), Inches(3.95), Inches(2.6), Inches(0.09), M_CYAN)
        add_text(s, Inches(1.05), Inches(4.2), Inches(11), Inches(0.6),
                 rp("title_role", "Computer Vision  ·  AI / ML  ·  XR & Game Engineering"),
                 font=B_FONT, size=19, color=M_GREY_D)
        add_text(s, Inches(1.05), Inches(5.1), Inches(11), Inches(0.4),
                 "M.Sc. Computer Science  ·  York University  ·  Graduating June 2026",
                 font=B_FONT, size=14, color=M_INK)
        add_text(s, Inches(1.05), Inches(5.5), Inches(11), Inches(0.4),
                 "Machine Learning Associate, Vector Institute",
                 font=B_FONT, size=13, italic=True, color=M_GREY_M)
        add_text(s, Inches(1.05), Inches(6.5), Inches(11.5), Inches(0.4),
                 "danielkordm@gmail.com   ·   +1 (437) 559 3462   ·   linkedin.com/in/danial-kord   ·   github.com/Danial-Kord",
                 font=B_FONT, size=12, color=M_GREY_M)
        return

    fill_slide(s, NAVY)
    # decorative amber vertical bar on left
    add_rect(s, Inches(0.6), Inches(1.0), Inches(0.12), Inches(5.5), AMBER)
    # accent dot
    add_oval(s, Inches(11.6), Inches(0.7), Inches(0.5), Inches(0.5), AMBER)
    add_oval(s, Inches(11.85), Inches(0.95), Inches(0.5), Inches(0.5), INDIGO)
    # tag
    add_text(s, Inches(1.0), Inches(1.05), Inches(6), Inches(0.4),
             rp("eyebrow", "PORTFOLIO  ·  CURRICULUM VITAE  ·  2026"),
             font=H_FONT, size=12, bold=True, color=AMBER)
    # main title
    add_text(s, Inches(1.0), Inches(1.7), Inches(11), Inches(1.6),
             "Danial Kordmodanlou",
             font=H_FONT, size=72, bold=True, color=IVORY)
    # subtitle / role
    add_text(s, Inches(1.0), Inches(3.5), Inches(11), Inches(0.6),
             rp("title_role", "Computer Vision  ·  AI / ML  ·  XR & Game Engineering"),
             font=B_FONT, size=22, color=GOLD)
    # divider
    add_rect(s, Inches(1.0), Inches(4.4), Inches(2.5), Inches(0.04), AMBER)
    # contact info
    add_text(s, Inches(1.0), Inches(4.6), Inches(11), Inches(0.4),
             "M.Sc. Candidate, Computer Science  ·  York University, Toronto",
             font=B_FONT, size=16, color=IVORY)
    add_text(s, Inches(1.0), Inches(5.05), Inches(11), Inches(0.4),
             "Machine Learning Associate, Vector Institute",
             font=B_FONT, size=14, italic=True, color=SLATE)
    # contact strip at bottom
    add_text(s, Inches(1.0), Inches(6.2), Inches(11.5), Inches(0.4),
             "danielkordm@gmail.com   ·   +1 (437) 559 3462   ·   linkedin.com/in/danial-kord   ·   github.com/Danial-Kord",
             font=B_FONT, size=12, color=SLATE)


# ====================================================================
# SLIDE 2 — ABOUT ME
# ====================================================================
def slide_about():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s)
    # left identity panel
    add_rect(s, 0, 0, Inches(5.2), SLIDE_H, PANEL)
    # photo on left panel
    add_image(s, f"{MEDIA}/image39.jpg", Inches(0.9), Inches(1.6), w=Inches(3.4))
    # small accent bar under photo
    add_rect(s, Inches(0.9), Inches(5.3), Inches(0.6), Inches(0.06), ACC1)
    add_text(s, Inches(0.9), Inches(5.4), Inches(4), Inches(0.4),
             "DANIAL KORDMODANLOU",
             font=H_FONT, size=14, bold=True, color=ON_DARK)
    add_text(s, Inches(0.9), Inches(5.75), Inches(4), Inches(0.4),
             "M.Sc. Computer Science", font=B_FONT, size=12, color=PANEL_ACC)
    add_text(s, Inches(0.9), Inches(6.05), Inches(4), Inches(0.4),
             "Toronto, Canada", font=B_FONT, size=11, color=MUTE_D)

    # right content
    add_text(s, Inches(5.7), Inches(0.7), Inches(7), Inches(0.5),
             "ABOUT ME", font=H_FONT, size=12, bold=True, color=EYE)
    add_rect(s, Inches(5.7), Inches(1.05), Inches(0.6), Inches(0.06), ACC1)

    add_text(s, Inches(5.7), Inches(1.25), Inches(7), Inches(1.0),
             "Hello, I'm Danial.",
             font=H_FONT, size=38, bold=True, color=INK)
    add_text(s, Inches(5.7), Inches(2.05), Inches(7), Inches(0.6),
             rp("about_sub", "I build at the intersection of vision, AI and immersive systems."),
             font=B_FONT, size=18, italic=True, color=SUB)

    bullets = rp("about_bullets", [
        "Recently completed a [b]Machine Learning Associate[/b] term at the Vector Institute, building real-time VR firefighter training with skeletal-telemetry deviation detection and LLM-driven coaching.",
        "M.Sc. researcher in the [b]BioMotion Lab[/b] at York University ([b]graduating June 2026[/b]) — stereopsis, motion parallax and immersive telepresence on Meta Quest / Unity / Unreal.",
        "5+ years shipping production [b]Unity3D, C#, Python, ML[/b] systems for games, simulation, training and digital humans.",
        "Open-source author of [b]DigiHuman[/b] (500⭐ GitHub), an automatic 3D character animation pipeline from monocular video.",
    ])
    add_bullets(s, Inches(5.7), Inches(2.85), Inches(7.2), Inches(4.0),
                bullets, size=14, color=BODY, line_spacing=1.25,
                bullet_char="●")
    # stats strip
    stat_y = Inches(6.3)
    stats = [
        ("9+", "yrs coding"),
        ("4.0/4.0", "MSc GPA"),
        ("500+", "GitHub ⭐"),
        ("Top 0.5%", "AI exam"),
    ]
    for i, (n, lbl) in enumerate(stats):
        x = Inches(5.7 + i*1.85)
        add_text(s, x, stat_y, Inches(1.7), Inches(0.45),
                 n, font=H_FONT, size=22, bold=True, color=ACC1)
        add_text(s, x, Inches(6.75), Inches(1.7), Inches(0.3),
                 lbl, font=B_FONT, size=10, color=BODY)
    light_footer(s, "ABOUT", 2)


# ====================================================================
# SLIDE 3 — QUOTE
# ====================================================================
def slide_quote():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s, feature=True)
    # giant quote glyph
    add_text(s, Inches(0.6), Inches(0.5), Inches(2), Inches(2.2),
             "“", font=H_FONT, size=220, bold=True, color=ACC1)
    # text
    add_text(s, Inches(1.4), Inches(2.4), Inches(10.5), Inches(2.5),
             "Yesterday is history, tomorrow is a mystery,",
             font=H_FONT, size=34, bold=True, color=INK_F, align=PP_ALIGN.LEFT)
    add_text(s, Inches(1.4), Inches(3.1), Inches(10.5), Inches(0.7),
             "but today is a gift — that's why it's called present.",
             font=H_FONT, size=34, bold=True, color=INK_F, align=PP_ALIGN.LEFT)
    # attribution
    add_rect(s, Inches(1.4), Inches(4.4), Inches(0.4), Inches(0.04), ACC1)
    add_text(s, Inches(2.0), Inches(4.25), Inches(8), Inches(0.5),
             "Master Oogway — Kung Fu Panda",
             font=B_FONT, size=16, italic=True, color=SUB_F)
    page_footer(s, "QUOTE", 3)


# ====================================================================
# SLIDE 4 — TABLE OF CONTENTS
# ====================================================================
def slide_toc():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s)
    # header
    add_text(s, Inches(0.7), Inches(0.7), Inches(8), Inches(0.5),
             "TABLE OF CONTENTS", font=H_FONT, size=12, bold=True, color=EYE)
    add_rect(s, Inches(0.7), Inches(1.05), Inches(0.6), Inches(0.06), ACC1)
    add_text(s, Inches(0.7), Inches(1.25), Inches(11), Inches(1.0),
             "What's inside.", font=H_FONT, size=44, bold=True, color=INK)

    sections = [
        ("01", "Education",   "Academic foundations — NODET, AUT, York University"),
        ("02", "Experience",  "Industry & research — Vector, DreamForge, BioMotion Lab, TectoTrack, Techu, IAESTE"),
        ("03", "Skills",      "Languages, frameworks, certifications, honors & awards"),
        ("04", "Projects",    "DigiHuman, SAFEZone AI, CaseLogic, LaTeX CV Builder … and many more"),
        ("05", "Hobbies",     "Movies, games, animes, board games, sports"),
    ]
    y = Inches(2.85)
    for num, title, desc in sections:
        # numbered card
        add_text(s, Inches(0.7), y, Inches(1.2), Inches(0.7),
                 num, font=H_FONT, size=42, bold=True, color=ACC1)
        add_text(s, Inches(2.0), y, Inches(4), Inches(0.45),
                 title, font=H_FONT, size=20, bold=True, color=INK)
        add_text(s, Inches(2.0), y + Inches(0.45), Inches(10), Inches(0.4),
                 desc, font=B_FONT, size=12, color=BODY)
        # thin divider
        add_rect(s, Inches(0.7), y + Inches(0.85), Inches(12), Inches(0.012), DIVLINE)
        y += Inches(0.82)
    light_footer(s, "TOC", 4)


# ====================================================================
# SECTION DIVIDERS
# ====================================================================
def section_divider(num, label, sub, slide_num):
    s = prs.slides.add_slide(blank_layout)
    page_bg(s, feature=True)
    # large numeric watermark
    add_text(s, Inches(0.6), Inches(0.6), Inches(6), Inches(5.5),
             num, font=H_FONT, size=380, bold=True,
             color=WMARK)  # very pale watermark number
    # accent bar
    add_rect(s, Inches(6.0), Inches(3.0), Inches(0.12), Inches(1.6), ACC1)
    # label
    add_text(s, Inches(6.3), Inches(2.9), Inches(7), Inches(0.5),
             f"SECTION {num}", font=H_FONT, size=14, bold=True, color=EYE)
    add_text(s, Inches(6.3), Inches(3.35), Inches(7), Inches(1.2),
             label, font=H_FONT, size=68, bold=True, color=INK_F)
    add_text(s, Inches(6.3), Inches(4.65), Inches(7), Inches(0.6),
             sub, font=B_FONT, size=16, italic=True, color=SUB_F)
    page_footer(s, label.upper(), slide_num)
    return s


# ====================================================================
# SLIDE - EDUCATION TIMELINE
# ====================================================================
def slide_education():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s)
    add_section_marker(s, "01 / EDUCATION")
    add_text(s, Inches(0.7), Inches(1.05), Inches(11), Inches(0.9),
             "Education", font=H_FONT, size=44, bold=True, color=INK)
    add_text(s, Inches(0.7), Inches(1.95), Inches(11), Inches(0.4),
             "Three institutions · three milestones",
             font=B_FONT, size=15, italic=True, color=SUB)

    # horizontal timeline
    line_y = Inches(3.5)
    add_rect(s, Inches(0.9), line_y, Inches(11.5), Inches(0.04), INK)
    # three nodes
    xs = [Inches(1.5), Inches(6.3), Inches(11.1)]
    years = ["2013 — 2017", "2017 — 2022", "2024 — 2026"]
    titles = ["High School Diploma", "Bachelor of CE", "Master of CS"]
    schools = ["NODET (Allameh Helli)\nMathematics & Physics", "Amirkabir University of Technology\nAI Concentration", "York University, Toronto\nVR / Computer Vision"]
    gpas = ["CGPA 19.56 / 20  (4.0 equiv.)", "CGPA 3.86 / 4.0", "CGPA 4.0 / 4.0  — Fully Funded"]
    for x, year, t, sch, gpa in zip(xs, years, titles, schools, gpas):
        # node circle
        add_oval(s, x - Inches(0.18), line_y - Inches(0.15), Inches(0.36), Inches(0.36),
                 ACC1)
        add_oval(s, x - Inches(0.08), line_y - Inches(0.05), Inches(0.18), Inches(0.18),
                 PAGE)
        # year above
        add_text(s, x - Inches(1.5), Inches(2.7), Inches(3), Inches(0.4),
                 year, font=H_FONT, size=12, bold=True, color=ACC1,
                 align=PP_ALIGN.CENTER)
        # title and details below
        add_text(s, x - Inches(1.8), Inches(3.95), Inches(3.6), Inches(0.5),
                 t, font=H_FONT, size=16, bold=True, color=INK,
                 align=PP_ALIGN.CENTER)
        add_text(s, x - Inches(1.8), Inches(4.4), Inches(3.6), Inches(1.0),
                 sch, font=B_FONT, size=11, color=BODY,
                 align=PP_ALIGN.CENTER)
        add_text(s, x - Inches(1.8), Inches(5.3), Inches(3.6), Inches(0.4),
                 gpa, font=B_FONT, size=10, italic=True, color=SUB,
                 align=PP_ALIGN.CENTER)
    # highlights box at the bottom
    add_round_rect(s, Inches(0.7), Inches(6.05), Inches(12), Inches(0.85), CARD,
                   line=cline(ACC1))
    add_text(s, Inches(0.95), Inches(6.18), Inches(11.5), Inches(0.4),
             "Highlights", font=H_FONT, size=11, bold=True, color=ACC1)
    add_text(s, Inches(0.95), Inches(6.45), Inches(11.7), Inches(0.45),
             "Top 1% in Iranian university entrance exam  ·  Top 0.5% AI Graduate National Exam  ·  Fully-funded MSc + Mitacs $10K + L2M $10K",
             font=B_FONT, size=11, color=CBODY)
    light_footer(s, "EDUCATION", 6)


# ====================================================================
# SLIDE - EDUCATION DETAIL: York
# ====================================================================
def slide_edu_york():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s)
    add_section_marker(s, "01 / EDUCATION  ·  M.SC.")
    # right photo strip
    add_rect(s, Inches(8.6), 0, Inches(4.733), SLIDE_H, FRAME)
    add_image(s, f"{MEDIA}/image41.jpg", Inches(8.85), Inches(2.0), w=Inches(4.2))
    add_text(s, Inches(8.85), Inches(5.6), Inches(4.2), Inches(0.4),
             "York University", font=H_FONT, size=13, bold=True, color=CHEAD)
    add_text(s, Inches(8.85), Inches(5.9), Inches(4.2), Inches(0.3),
             "Toronto, Ontario, Canada", font=B_FONT, size=10, color=ACC2)

    add_text(s, Inches(0.7), Inches(1.05), Inches(8), Inches(0.45),
             "2024 — JUNE 2026  ·  M.SC.", font=H_FONT, size=12,
             bold=True, color=EYE)
    add_text(s, Inches(0.7), Inches(1.45), Inches(8), Inches(1.0),
             "Master of Computer Science", font=H_FONT, size=34, bold=True, color=INK)
    add_text(s, Inches(0.7), Inches(2.4), Inches(8), Inches(0.5),
             "York University, Toronto", font=B_FONT, size=16, italic=True, color=SUB)
    # stat row
    stat_y = Inches(3.2)
    cards = [
        ("4.0 / 4.0", "Final GPA"),
        ("100%", "Fully-funded RA + TA"),
        ("$10K", "Mitacs Funding"),
    ]
    for i, (n, lbl) in enumerate(cards):
        x = Inches(0.7 + i*2.6)
        add_round_rect(s, x, stat_y, Inches(2.35), Inches(0.95), CARD,
                       line=cline(ACC1))
        add_text(s, x + Inches(0.15), stat_y + Inches(0.1), Inches(2.1), Inches(0.45),
                 n, font=H_FONT, size=20, bold=True, color=ACC1)
        add_text(s, x + Inches(0.15), stat_y + Inches(0.55), Inches(2.1), Inches(0.35),
                 lbl, font=B_FONT, size=10, color=CBODY)
    # focus bullets
    add_text(s, Inches(0.7), Inches(4.45), Inches(8), Inches(0.4),
             "FOCUS AREAS", font=H_FONT, size=11, bold=True, color=EYE)
    bullets = [
        "VR & motion capture experiments on [b]Meta Quest, ArKit, Unity3D, Unreal[/b]",
        "Depth perception research: isolating [b]stereopsis[/b] and [b]motion parallax[/b]",
        "Discovered & published a stereo-geometry IPD bug in [b]Meta XR SDK Unity[/b]",
        "Presented [b]illusory parallax in VR[/b] at ECVP 2025 (Mainz, Germany)",
        "TA for OOP / Java for [b]5 semesters[/b] — design patterns, optimization",
    ]
    add_bullets(s, Inches(0.7), Inches(4.85), Inches(8.0), Inches(2.4),
                bullets, size=12, line_spacing=1.2, bullet_char="▸",
                color=BODY)
    light_footer(s, "EDUCATION", 7)


# ====================================================================
# SLIDE - EDUCATION DETAIL: AUT
# ====================================================================
def slide_edu_aut():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s, feature=True)
    add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
             "01 / EDUCATION  ·  B.SC.", font=H_FONT, size=11, bold=True, color=EYE)
    add_rect(s, Inches(0.7), Inches(0.85), Inches(0.4), Inches(0.04), ACC1)
    # left photo (use rotated version)
    add_image(s, f"{MEDIA}/image32_fixed.jpg", Inches(0.7), Inches(1.3), h=Inches(5.7))
    # right text
    add_text(s, Inches(5.5), Inches(1.3), Inches(7), Inches(0.45),
             "2017 — 2022", font=H_FONT, size=12, bold=True, color=ACC2)
    add_text(s, Inches(5.5), Inches(1.7), Inches(7), Inches(1.0),
             "Bachelor of Computer Engineering",
             font=H_FONT, size=28, bold=True, color=INK_F)
    add_text(s, Inches(5.5), Inches(2.75), Inches(7), Inches(0.5),
             "Amirkabir University of Technology (AUT)", font=B_FONT, size=15,
             italic=True, color=SUB_F)
    # GPA card
    add_round_rect(s, Inches(5.5), Inches(3.4), Inches(3.4), Inches(0.95), CARD2,
                   line=cline(ACC1))
    add_text(s, Inches(5.7), Inches(3.5), Inches(3.2), Inches(0.4),
             "CGPA  3.86 / 4.0", font=H_FONT, size=16, bold=True, color=ACC1)
    add_text(s, Inches(5.7), Inches(3.95), Inches(3.2), Inches(0.3),
             "140 / 140 credits  ·  AI focus", font=B_FONT, size=10, color=CBODY)
    add_round_rect(s, Inches(9.05), Inches(3.4), Inches(3.4), Inches(0.95), CARD2,
                   line=cline(ACC1))
    add_text(s, Inches(9.25), Inches(3.5), Inches(3.2), Inches(0.4),
             "Top 0.5%", font=H_FONT, size=16, bold=True, color=ACC1)
    add_text(s, Inches(9.25), Inches(3.95), Inches(3.2), Inches(0.3),
             "AI Graduate National Exam", font=B_FONT, size=10, color=CBODY)

    add_text(s, Inches(5.5), Inches(4.55), Inches(7), Inches(0.4),
             "FOCUS AREAS", font=H_FONT, size=11, bold=True, color=EYE)
    bullets = [
        "Concentrated on [b]Artificial Intelligence[/b] and Computer Vision",
        "Teaching Assistant for AI, Advanced Java, C, and Operating Systems",
        "Technical Staff member for [b]GameCraft[/b] events at AUT",
        "Capstone thesis: [b]DigiHuman[/b] — 500⭐ open-source release",
    ]
    add_bullets(s, Inches(5.5), Inches(4.95), Inches(7.4), Inches(2.0),
                bullets, size=12, line_spacing=1.25, bullet_char="▸",
                color=CHEAD)
    page_footer(s, "EDUCATION", 8)


# ====================================================================
# SLIDE - EXPERIENCE OVERVIEW
# ====================================================================
def slide_exp_overview():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s)
    add_section_marker(s, "02 / EXPERIENCE")
    add_text(s, Inches(0.7), Inches(1.05), Inches(11), Inches(0.9),
             "Career timeline", font=H_FONT, size=40, bold=True, color=INK)
    add_text(s, Inches(0.7), Inches(1.85), Inches(11), Inches(0.4),
             rp("exp_sub", "From hyper-casual games to VR research and AI tooling — a track of shipping production systems."),
             font=B_FONT, size=14, italic=True, color=SUB)

    # Vertical timeline on left
    items = [
        ("2026", "Machine Learning Associate",     "Vector Institute · Toronto",        "Real-time VR firefighter training: skeletal-telemetry deviation engine + LLM coaching layer (RAG, FAISS, Ollama).", ACC1),
        ("2025–26", "Senior Unity Developer", "DreamForge · Miami (Remote)",        "AI-driven engine in C# / Python. Built Claude-assisted tools that auto-investigate exceptions and commit fixes.", ACC1),
        ("2024–26", "Research & Teaching Assistant", "BioMotion Lab · York University", "VR + motion-capture experiments on stereopsis / motion parallax. Presented at ECVP 2025.", ACC1),
        ("2023–25", "Software Developer",     "TectoTrack · Toronto",               "Large-scale crowd simulations for airports & malls with fault-tolerant pathfinding and DOTS ECS.", ACC2),
        ("2023–24", "Lead Unity3D Developer", "Techu · Japan (Remote)",             "Strategic cross-platform card game with DQN + MCTS AI opponents, Photon/PlayFab multiplayer.", ACC2),
        ("2022",       "Volunteer Full-Stack Developer", "IAESTE · Tehran",            "Maintained Vue.js + Django student-exchange platform.", MUTE),
        ("2021",       "Software Engineer Intern", "Sepantab (Hojres) · Tehran",        "Built Stone Thrower — WebGL multiplayer with WebSockets / gRPC / MQTT.", MUTE),
        ("2017–22","Technical Staff & TA",    "AUT · Tehran",                       "GameCraft events. TA for AI, Adv. Java, C, OS.", MUTE),
    ]
    # 2-column grid
    col_w = Inches(5.95)
    col_h = Inches(1.07)
    gx = [Inches(0.7), Inches(6.85)]
    gy = [Inches(2.6), Inches(3.7), Inches(4.8), Inches(5.9)]
    for i, item in enumerate(items):
        c = i % 2
        r = i // 2
        x = gx[c]
        y = gy[r]
        year, role, org, summary, dot_color = item
        # card
        add_round_rect(s, x, y, col_w, col_h, CARD, radius=0.1, line=cline(dot_color))
        # date pill on left
        add_rect(s, x + Inches(0.0), y, Inches(0.08), col_h, dot_color)
        add_text(s, x + Inches(0.2), y + Inches(0.1), Inches(1.4), Inches(0.4),
                 year, font=H_FONT, size=11, bold=True, color=dot_color)
        add_text(s, x + Inches(0.2), y + Inches(0.43), Inches(5.8), Inches(0.4),
                 role, font=H_FONT, size=13, bold=True, color=CHEAD)
        add_text(s, x + Inches(0.2), y + Inches(0.72), Inches(5.8), Inches(0.3),
                 org, font=B_FONT, size=10, italic=True, color=ACC2)
    light_footer(s, "EXPERIENCE", 10)


# ====================================================================
# EXPERIENCE DETAIL TEMPLATE
# ====================================================================
def exp_detail(s, *, period, company, location, title, summary, bullets,
               tags, image_path=None, image_pos="right", slide_num=0,
               accent=None, key=None):
    if accent is None:
        accent = ACC1
    # role-mode may re-frame this position toward the target role
    ov = rp("exp", {}).get(key, {})
    title = ov.get("title", title)
    summary = ov.get("summary", summary)
    bullets = ov.get("bullets", bullets)
    tags = ov.get("tags", tags)
    page_bg(s)
    add_section_marker(s, "02 / EXPERIENCE", accent)
    # header
    add_text(s, Inches(0.7), Inches(1.05), Inches(8), Inches(0.4),
             period.upper(), font=H_FONT, size=12, bold=True, color=accent)
    add_text(s, Inches(0.7), Inches(1.45), Inches(11), Inches(0.9),
             title, font=H_FONT, size=34, bold=True, color=INK)
    add_text(s, Inches(0.7), Inches(2.35), Inches(11), Inches(0.4),
             f"{company}  ·  {location}",
             font=B_FONT, size=14, italic=True, color=SUB)
    add_text(s, Inches(0.7), Inches(2.8), Inches(11), Inches(0.4),
             summary, font=B_FONT, size=13, color=BODY)

    content_x = Inches(0.7)
    content_w = Inches(7.5)
    if image_path:
        if image_pos == "right":
            add_rect(s, Inches(8.6), Inches(3.4), Inches(4.0), Inches(3.4), FRAME)
            add_image(s, image_path, Inches(8.7), Inches(3.5), w=Inches(3.8))
        else:
            content_x = Inches(5.3)
            content_w = Inches(7.5)
            add_rect(s, Inches(0.7), Inches(3.4), Inches(4.0), Inches(3.4), FRAME)
            add_image(s, image_path, Inches(0.8), Inches(3.5), w=Inches(3.8))

    # bullets
    add_text(s, content_x, Inches(3.45), Inches(6), Inches(0.35),
             "WHAT I DID", font=H_FONT, size=10, bold=True, color=accent)
    add_bullets(s, content_x, Inches(3.8), content_w, Inches(2.8), bullets,
                size=12, line_spacing=1.2, bullet_char="▸",
                color=BODY)
    # tag chips at bottom
    tag_y = Inches(6.6)
    tx = content_x
    for tag in tags:
        chip_w = Inches(0.18 + 0.085 * len(tag))
        add_round_rect(s, tx, tag_y, chip_w, Inches(0.32), CHIP_BG, radius=0.4)
        add_text(s, tx, tag_y + Inches(0.04), chip_w, Inches(0.28),
                 tag, font=B_FONT, size=9, bold=True, color=CHIP_FG,
                 align=PP_ALIGN.CENTER)
        tx += chip_w + Inches(0.08)
    light_footer(s, "EXPERIENCE", slide_num)


def slide_exp_mla():
    s = prs.slides.add_slide(blank_layout)
    exp_detail(s,
               period="January 2026 — May 2026",
               company="Vector Institute",
               location="Toronto, Canada",
               title="Machine Learning Associate (MLA)",
               summary="Real-time VR firefighter training in partnership with DXTR — deviation detection from headset telemetry + LLM-driven coaching.",
               bullets=[
                   "Architected a [b]Deviation Engine[/b] over sliding windows of 3D skeletal telemetry; computed [b]MPJPE[/b] and quaternion angular distances to score deviations from training standards.",
                   "Built an [b]LLM Coaching Layer[/b] using [b]RAG[/b] (FAISS similarity search) and [b]Ollama / LangChain[/b] to generate natural-language feedback grounded in training manuals.",
                   "Engineered a full-body [b]pose collection pipeline[/b] mapping headset telemetry to the [b]NTU-25[/b] skeleton with cross-platform frame normalization.",
                   "Collaborated on structured evaluation of generative-AI pipelines for high-fidelity training scenarios.",
               ],
               tags=["RAG", "FAISS", "Ollama", "LangChain", "NTU-25", "VR Telemetry", "Quaternion", "MPJPE"],
               slide_num=11, key="vector")


def slide_exp_dreamforge():
    s = prs.slides.add_slide(blank_layout)
    exp_detail(s,
               period="December 2025 — February 2026",
               company="DreamForge",
               location="Miami, Florida (Remote)",
               title="Senior Unity Developer",
               summary="Core contributor to a C# + Python AI-driven engine. Built Claude-assisted developer tools that automatically intercept exceptions and commit codebase resolutions.",
               bullets=[
                   "Designed [b]CI/CD delivery components[/b] and [b]Docker-based workflows[/b] in [b].NET / Azure / GitHub Actions[/b] for automated deployment and exception analysis.",
                   "Built [b]AI-assisted tools (Claude)[/b] that auto-intercept exceptions, investigate root causes and commit permanent codebase fixes.",
                   "Performed deep code reviews and led peers in resolving structural inefficiencies; maintained debugging workflows in Visual Studio.",
                   "Focus on robust tooling, procedural generation, and continuous delivery pipelines.",
               ],
               tags=["C#", ".NET", "Python", "Docker", "Azure", "GitHub Actions", "Claude", "Unity"],
               slide_num=12, key="dreamforge")


def slide_dreamforge_cicd():
    """DreamForge — animated crash-to-merged-fix pipeline (DigiHuman-style).
    A runtime exception flows Backtrace -> .NET -> queue -> a triage agent;
    criticals run Claude Code in an Azure container that opens a reviewed PR
    and merges — every step is mirrored in Slack."""
    s = prs.slides.add_slide(blank_layout)
    fill_slide(s, M_WHITE)

    RED   = RGBColor(0xDB, 0x3A, 0x34)   # the crash / critical path
    GREEN = RGBColor(0x2F, 0xA0, 0x55)   # merged
    GREY  = RGBColor(0x8A, 0x94, 0xA6)   # non-critical
    SLACK = RGBColor(0x4A, 0x15, 0x4B)   # Slack aubergine

    # --- decorative blobs (static page frame) ---
    add_oval(s, Inches(11.98), Inches(-0.80), Inches(2.05), Inches(2.05), M_PURPLE)
    add_oval(s, Inches(12.63), Inches(2.92), Inches(0.74), Inches(0.74), M_CYAN)
    add_oval(s, Inches(-0.55), Inches(5.45), Inches(1.25), Inches(1.25), M_YELLOW)
    add_oval(s, Inches(0.45), Inches(0.34), Inches(0.24), Inches(0.24), DH_CORAL)

    # --- title + subtitle (static) ---
    add_text(s, Inches(0.5), Inches(0.10), Inches(12.33), Inches(0.52),
             "DreamForge — Autonomous Exception Resolution", font=H_FONT,
             size=28, bold=True, color=DH_INK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(0.62), Inches(12.33), Inches(0.34),
             "A live crash is captured, triaged, and fixed by Claude Code in "
             "Azure — merged via a reviewed PR, all mirrored in Slack",
             font=B_FONT, size=11.5, italic=True, color=M_TEAL,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ---- nested helpers ----
    def card(x, y, w, h, fill, *, radius=0.10, line=None, line_w=1.25):
        return add_round_rect(s, Inches(x), Inches(y), Inches(w), Inches(h),
                              fill, radius=radius, line=line, line_w=line_w)

    def t(x, y, w, h, text, *, size=10, bold=False, color=DH_INK,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, font=None, ls=None):
        add_text(s, Inches(x), Inches(y), Inches(w), Inches(h), text,
                 font=font or H_FONT, size=size, bold=bold, color=color,
                 align=align, anchor=anchor, line_spacing=ls)

    def hchev(x, y, fill, left=False, w=0.42, h=0.5):
        add_chevron(s, Inches(x), Inches(y), Inches(w), Inches(h), fill, left=left)

    def rail(x, y, w, h, fill):
        add_rect(s, Inches(x), Inches(y), Inches(w), Inches(h), fill)

    def vchev(x, y, fill, up=False, w=0.44, h=0.42):
        c = add_chevron(s, Inches(x), Inches(y), Inches(w), Inches(h), fill)
        c.rotation = 270 if up else 90

    def stage(x, y, w, h, num, title, sub, *, num_fill=DH_CORAL, num_color=M_WHITE,
              fill=DH_MINT, tcolor=DH_INK, scolor=DH_INK, line=None, tsize=13):
        card(x, y, w, h, fill, line=line, line_w=2.25 if line else 1.25)
        cd = 0.4
        add_oval(s, Inches(x + w / 2 - cd / 2), Inches(y + 0.11), Inches(cd),
                 Inches(cd), num_fill)
        t(x + w / 2 - cd / 2, y + 0.11, cd, cd, num, size=14, bold=True,
          color=num_color, anchor=MSO_ANCHOR.MIDDLE)
        t(x + 0.06, y + 0.55, w - 0.12, 0.4, title, size=tsize, bold=True,
          color=tcolor)
        t(x + 0.06, y + 0.95, w - 0.12, h - 0.99, sub, size=8.7, color=scolor,
          ls=1.05)

    TOP_Y, TH = 1.55, 1.5
    BOT_Y, BH = 4.35, 1.5
    TOP_BOT = TOP_Y + TH           # 3.05 — bottom edge of the top row

    # ===== TOP ROW (left -> right): Exception -> Backtrace -> .NET -> Queue -> Classify
    TW = 2.14
    stage(0.40, TOP_Y, TW, TH, "1", "Exception", "thrown in the live app",
          num_fill=RED)
    hchev(2.56, TOP_Y + 0.50, DH_ARROW)
    stage(2.99, TOP_Y, TW, TH, "2", "Backtrace", "cloud crash monitoring")
    hchev(5.15, TOP_Y + 0.50, DH_ARROW)
    stage(5.58, TOP_Y, TW, TH, "3", ".NET server", "receives the webhook")
    hchev(7.74, TOP_Y + 0.50, DH_ARROW)
    stage(8.17, TOP_Y, TW, TH, "4", "Queue", "enqueued for triage")
    hchev(10.33, TOP_Y + 0.50, DH_ARROW)
    stage(10.76, TOP_Y, TW, TH, "5", "Classify", "critical or non-critical?",
          num_fill=M_PURPLE)

    # ===== BOTTOM ROW (right -> left): Azure -> Claude Code -> Open PR -> Merge
    BW = 2.7
    # 6 - Azure container (+ the red "critical" drop from the classifier)
    stage(10.00, BOT_Y, BW, BH, "6", "Azure container",
          "isolated container spins up")
    rail(11.33, TOP_BOT, 0.05, BOT_Y - TOP_BOT, RED)
    vchev(11.13, BOT_Y - 0.40, RED)
    t(11.42, TOP_BOT + 0.10, 1.6, 0.26, "critical", size=9.5, bold=True,
      color=RED, align=PP_ALIGN.LEFT)
    hchev(9.53, BOT_Y + 0.50, DH_ARROW, left=True)
    # 7 - Claude Code (the star, highlighted)
    stage(6.80, BOT_Y, BW, BH, "7", "Claude Code",
          "fixes on a fresh git pull with the exception", line=DH_CORAL)
    hchev(6.33, BOT_Y + 0.50, DH_ARROW, left=True)
    # 8 - Open PR
    stage(3.60, BOT_Y, BW, BH, "8", "Open PR",
          "commit · push · assign reviewers")
    hchev(3.13, BOT_Y + 0.50, DH_ARROW, left=True)
    # 9 - Merge
    stage(0.40, BOT_Y, BW, BH, "9", "Merge", "reviewers approve → merged",
          fill=GREEN, tcolor=M_WHITE, scolor=M_WHITE, num_fill=M_WHITE,
          num_color=GREEN)

    # ===== non-critical branch (an aside off the classifier) =====
    card(4.70, TOP_BOT + 0.37, 3.2, 0.5, GREY, radius=0.16)
    t(4.70, TOP_BOT + 0.37, 3.2, 0.5, "non-critical  ·  logged & deprioritized",
      size=9, bold=True, color=M_WHITE, anchor=MSO_ANCHOR.MIDDLE)
    rail(11.00, TOP_BOT, 0.05, 0.66, GREY)
    rail(7.90, TOP_BOT + 0.61, 11.05 - 7.90, 0.05, GREY)
    hchev(7.62, TOP_BOT + 0.46, GREY, left=True, w=0.3, h=0.34)

    # ===== Slack observability rail (mirrors every step) =====
    card(0.40, 6.15, 12.53, 0.58, SLACK, radius=0.14)
    for i, dc in enumerate([RGBColor(0xE0, 0x1E, 0x5A), RGBColor(0xEC, 0xB2, 0x2E),
                            RGBColor(0x2E, 0xB6, 0x7D), RGBColor(0x36, 0xC5, 0xF0)]):
        add_oval(s, Inches(0.62 + i * 0.17), Inches(6.35), Inches(0.13),
                 Inches(0.13), dc)
    t(1.50, 6.15, 1.3, 0.58, "Slack", size=14, bold=True, color=M_WHITE,
      align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    t(2.80, 6.15, 9.9, 0.58,
      "one live thread mirrors the whole interaction — crash · triage · "
      "Claude's fix · PR · merge", size=10.5, color=M_WHITE,
      align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    _dh_fade_transition(s)
    _dh_build_anim(s)


def slide_exp_biomotion():
    s = prs.slides.add_slide(blank_layout)
    exp_detail(s,
               period="January 2024 — April 2026",
               company="BioMotion Lab, York University",
               location="Toronto, Canada",
               title="Research & Teaching Assistant",
               summary="Designed VR + motion-capture experiments on depth perception, isolating stereopsis from motion parallax.",
               bullets=[
                   "Conducted VR and motion-capture experiments using [b]Meta Quest, ArKit, Unity3D, Unreal Engine[/b].",
                   "Investigated and reported a critical stereo-geometry issue in [b]Meta XR SDK Unity[/b] (OVRCameraRig per-eye mode) — effective ~2× IPD scaling bug; published a reproducible test project.",
                   "Presented research on [b]illusory parallax in VR[/b] at the [b]47th European Conference on Visual Perception (ECVP 2025)[/b] in Mainz, Germany.",
                   "Taught Object-Oriented Programming (Java) for [b]5 semesters[/b], emphasizing design patterns and optimization.",
               ],
               tags=["Meta Quest", "ARKit", "Unity", "Unreal", "ECVP 2025", "Stereopsis", "Java OOP"],
               image_path=f"{MEDIA}/image4.png",
               image_pos="right",
               slide_num=13, key="biomotion")


def slide_exp_tectotrack():
    s = prs.slides.add_slide(blank_layout)
    exp_detail(s,
               period="August 2023 — February 2025",
               company="TectoTrack",
               location="Toronto, Canada",
               title="Software Developer",
               summary="Built large-scale simulations for airports and malls. Optimized high-density scenes with Unity DOTS / ECS.",
               bullets=[
                   "Built [b]simulation frameworks[/b] for high-traffic environments with fault-tolerant pathfinding for adaptive crowd movement.",
                   "Improved realism of agents with [b]adaptive behavior modeling[/b].",
                   "Developed [b]real-time monitoring dashboards[/b] for performance, agent states and logs.",
                   "Built internal tools for debugging, visualization and flow adjustment.",
                   "Delivered optimized scenes for high-density regions with [b]ECS[/b] (Unity DOTS).",
               ],
               tags=["Unity", "DOTS / ECS", "Pathfinding", "C#", "Dashboards", "Simulation"],
               image_path=f"{MEDIA}/image11.png",
               image_pos="right",
               slide_num=14, key="tectotrack")


def slide_exp_techu():
    s = prs.slides.add_slide(blank_layout)
    exp_detail(s,
               period="January 2023 — April 2024",
               company="Techu",
               location="Japan (Remote)",
               title="Lead Unity3D Game Developer",
               summary="Led development of Techu — a cross-platform strategic card game. Built AI agents and online multiplayer end-to-end.",
               bullets=[
                   "[b]Led[/b] gameplay and client engineering for a cross-platform multiplayer game.",
                   "Engineered AI opponents with [b]DQN[/b] + [b]Monte Carlo Tree Search[/b] for dynamic decision making.",
                   "Online multiplayer via [b]Photon[/b] + [b]PlayFab (Azure)[/b] for matchmaking.",
                   "Integrated [b]Firebase[/b] for real-time crash analytics and diagnostics.",
                   "Managed the full development pipeline: gameplay, UI, backend, store releases.",
               ],
               tags=["Unity", "C#", "DQN", "MCTS", "PlayFab", "Photon", "Firebase"],
               image_path=f"{MEDIA}/image18.png",
               image_pos="right",
               slide_num=15, key="techu")


def slide_exp_iaeste():
    s = prs.slides.add_slide(blank_layout)
    exp_detail(s,
               period="January 2022 — December 2022",
               company="IAESTE  ·  International Association for Exchange of Students",
               location="Tehran, Iran  ·  Volunteer",
               title="Volunteer Full Stack Developer",
               summary="Maintained and enhanced IAESTE's web platform serving 100+ countries.",
               bullets=[
                   "Maintained the platform using [b]Vue.js[/b] (frontend) and [b]Django[/b] (backend).",
                   "Improved UI/UX for student-exchange applications.",
                   "Optimized backend routes and reduced latency on key endpoints.",
                   "Contributed long-term maintainability through refactoring and tests.",
               ],
               tags=["Vue.js", "Django", "Python", "Volunteer"],
               image_path=f"{MEDIA}/image2.png",
               image_pos="right",
               slide_num=16,
               accent=ACC2, key="iaeste")


def slide_exp_sepantab():
    s = prs.slides.add_slide(blank_layout)
    exp_detail(s,
               period="June 2021 — August 2021",
               company="Sepantab (Beh Andishan Noavar)",
               location="Tehran, Iran",
               title="Software Engineer Intern",
               summary="Built Stone Thrower — a local-multiplayer IoT WebGL game with resilient networking.",
               bullets=[
                   "Developed [b]Stone Thrower[/b], a local-multiplayer IoT-based [b]WebGL[/b] game in Unity3D.",
                   "Implemented real-time networked messaging using [b]WebSockets[/b].",
                   "Built resilient distributed communication using [b]gRPC[/b] + [b]Protocol Buffers[/b] + [b]MQTT[/b].",
                   "Designed networking infrastructure in C# for low-latency state sync.",
               ],
               tags=["C#", "Unity WebGL", "WebSockets", "gRPC", "Protobuf", "MQTT", "IoT"],
               slide_num=17,
               accent=ACC2, key="sepantab")


# ====================================================================
# SKILLS section
# ====================================================================
def slide_skills():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s)
    add_section_marker(s, "03 / SKILLS")
    add_text(s, Inches(0.7), Inches(1.05), Inches(11), Inches(0.9),
             "Technical stack", font=H_FONT, size=40, bold=True, color=INK)
    add_text(s, Inches(0.7), Inches(1.85), Inches(11), Inches(0.4),
             "Languages, frameworks, runtimes and tools — grouped by depth.",
             font=B_FONT, size=14, italic=True, color=SUB)

    # 3 columns (role mode may re-group them)
    cols = rp("skills", [
        ("Fluent",       ["Unity3D", "Unreal Engine", "C#", "Java", "Python", "C++", "C"]),
        ("Intermediate", ["PyTorch", "Django", "Vue.js", "React", "Next.js", "OpenGL", "Docker", "Kubernetes", "AWS (ECR)", "OpenMP", "CUDA", "MySQL", "Spring Boot", "Firebase", "PlayFab", "Photon"]),
        ("Familiar",     ["DOTS / ECS", "Go", "Hadoop", "Keras", "Blender", "MediaPipe", "OpenCV", "LangChain", "Ollama", "FAISS"]),
    ])
    _skcolors = [ACC1, ACC2, ACC3]
    col_x = [Inches(0.7), Inches(5.0), Inches(9.3)]
    col_w = Inches(3.9)
    for i, (label, items) in enumerate(cols):
        color = _skcolors[i % 3]
        x = col_x[i]
        add_round_rect(s, x, Inches(2.7), col_w, Inches(4.0), CARD, radius=0.05,
                       line=cline(color))
        add_rect(s, x, Inches(2.7), col_w, Inches(0.5), STACKHDR)
        lbl_color = ON_DARK if DECO == "memphis" else color
        add_text(s, x + Inches(0.25), Inches(2.78), col_w - Inches(0.5), Inches(0.4),
                 label.upper(), font=H_FONT, size=14, bold=True, color=lbl_color,
                 anchor=MSO_ANCHOR.MIDDLE)
        # items as chips
        cx = x + Inches(0.22)
        cy = Inches(3.4)
        max_x = x + col_w - Inches(0.22)
        for item in items:
            chip_w = Inches(0.25 + 0.09 * len(item))
            if cx + chip_w > max_x:
                cx = x + Inches(0.22)
                cy += Inches(0.42)
            add_round_rect(s, cx, cy, chip_w, Inches(0.32), STACKCHIP_BG, radius=0.4)
            add_text(s, cx, cy + Inches(0.04), chip_w, Inches(0.28),
                     item, font=B_FONT, size=9, bold=True, color=STACKCHIP_FG,
                     align=PP_ALIGN.CENTER)
            cx += chip_w + Inches(0.08)
    light_footer(s, "SKILLS", 18)


def slide_languages():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s, feature=True)
    add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
             "03 / SKILLS  ·  COMMUNICATION", font=H_FONT, size=11, bold=True, color=EYE)
    add_rect(s, Inches(0.7), Inches(0.85), Inches(0.4), Inches(0.04), ACC1)
    add_text(s, Inches(0.7), Inches(1.0), Inches(11), Inches(0.9),
             "Languages & Communication", font=H_FONT, size=36, bold=True, color=INK_F)

    # English card
    add_round_rect(s, Inches(0.7), Inches(2.4), Inches(5.95), Inches(3.6), CARD2,
                   line=cline(ACC1))
    add_text(s, Inches(1.0), Inches(2.6), Inches(5.0), Inches(0.5),
             "English", font=H_FONT, size=24, bold=True, color=ACC1)
    add_text(s, Inches(1.0), Inches(3.1), Inches(5.0), Inches(0.4),
             "IELTS Certified", font=B_FONT, size=12, italic=True, color=ACC2)
    en_rows = [("Speaking", 7.0), ("Listening", 7.0), ("Reading", 6.5), ("Writing", 6.5)]
    by = Inches(3.7)
    for label, score in en_rows:
        add_text(s, Inches(1.0), by, Inches(2.0), Inches(0.35),
                 label, font=B_FONT, size=12, color=CBODY)
        # bar
        full_bar_w = Inches(2.6)
        add_rect(s, Inches(3.1), by + Inches(0.1), full_bar_w, Inches(0.18),
                 BAR_TRACK)
        filled = Inches(2.6 * (score / 9.0))
        add_rect(s, Inches(3.1), by + Inches(0.1), filled, Inches(0.18), ACC1)
        add_text(s, Inches(5.75), by, Inches(0.8), Inches(0.35),
                 f"{score}", font=H_FONT, size=12, bold=True, color=ACC2,
                 align=PP_ALIGN.RIGHT)
        by += Inches(0.45)

    # Persian card
    add_round_rect(s, Inches(6.85), Inches(2.4), Inches(5.95), Inches(3.6), CARD2,
                   line=cline(ACC1))
    add_text(s, Inches(7.15), Inches(2.6), Inches(5.0), Inches(0.5),
             "Persian", font=H_FONT, size=24, bold=True, color=ACC1)
    add_text(s, Inches(7.15), Inches(3.1), Inches(5.0), Inches(0.4),
             "Native speaker", font=B_FONT, size=12, italic=True, color=ACC2)
    fa_rows = [("Speaking", "Native"), ("Listening", "Native"),
               ("Reading", "Native"), ("Writing", "Native")]
    by = Inches(3.7)
    for label, score in fa_rows:
        add_text(s, Inches(7.15), by, Inches(2.0), Inches(0.35),
                 label, font=B_FONT, size=12, color=CBODY)
        add_rect(s, Inches(9.25), by + Inches(0.1), Inches(2.6), Inches(0.18), ACC1)
        add_text(s, Inches(11.9), by, Inches(0.8), Inches(0.35),
                 score, font=H_FONT, size=11, bold=True, color=ACC2,
                 align=PP_ALIGN.RIGHT)
        by += Inches(0.45)
    page_footer(s, "SKILLS", 19)


def slide_certifications():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s)
    add_section_marker(s, "03 / SKILLS  ·  CERTIFICATIONS & HACKATHONS")
    add_text(s, Inches(0.7), Inches(1.05), Inches(11), Inches(0.9),
             "Certifications & hackathons", font=H_FONT, size=36, bold=True, color=INK)
    add_text(s, Inches(0.7), Inches(1.9), Inches(11), Inches(0.4),
             "Selected coursework, hands-on bootcamps and competitive events.",
             font=B_FONT, size=13, italic=True, color=SUB)

    cards = [
        ("2026",  "L2M Lab to Market",
         "$10,000 funding accepted into the Lab-to-Market entrepreneurship program.",
         ACC1, "Funded"),
        ("2026",  "Legal-Tech Hackathon — 3rd Place",
         "Built CaseLogic in 24h: hybrid retrieval (Chroma + SQLite FTS5), per-claim verification, multi-agent planning.",
         ACC1, "3rd / Toronto"),
        ("2025",  "Boson AI × MScAC Hackathon",
         "Built SAFEZone AI — generative-AI therapist with voice cloning and lip-sync. Advanced to the semi-finals.",
         ACC2, "Semi-Finalist"),
        ("2025",  "Build Multi-Agent Applications Bootcamp",
         "Maven cohort. Hands-on building multi-agent AI systems with LangGraph, Cursor IDE, and n8n.",
         ACC2, "Certified"),
        ("2021",  "GANs Specialization",
         "DeepLearning.AI — image synthesis, pose editing, conditional & decoupled GANs.",
         MUTE, "Coursera"),
        ("2017",  "Top 1% Iranian Konkur",
         "Ranked in the top 1% of the Iranian National University Entrance Exam (Konkur).",
         MUTE, "National"),
    ]
    cols = 3
    cw = Inches(4.1)
    ch = Inches(2.2)
    gap_x = Inches(0.15)
    gap_y = Inches(0.25)
    sx = Inches(0.7)
    sy = Inches(2.45)
    for i, (year, title, body, color, badge) in enumerate(cards):
        r = i // cols
        c = i % cols
        x = sx + c * (cw + gap_x)
        y = sy + r * (ch + gap_y)
        add_round_rect(s, x, y, cw, ch, CARD, radius=0.06, line=cline(color))
        # left color band
        add_rect(s, x, y, Inches(0.1), ch, color)
        add_text(s, x + Inches(0.3), y + Inches(0.15), Inches(2), Inches(0.3),
                 year, font=H_FONT, size=11, bold=True, color=color)
        # badge pill right
        bw = Inches(0.28 + 0.085 * len(badge))
        add_round_rect(s, x + cw - bw - Inches(0.2), y + Inches(0.15),
                       bw, Inches(0.3), color, radius=0.4)
        badge_fg = ON_DARK if DECO == "memphis" else INK
        add_text(s, x + cw - bw - Inches(0.2), y + Inches(0.18),
                 bw, Inches(0.25),
                 badge, font=B_FONT, size=8, bold=True, color=badge_fg,
                 align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.3), y + Inches(0.5), cw - Inches(0.5), Inches(0.7),
                 title, font=H_FONT, size=14, bold=True, color=CHEAD)
        add_text(s, x + Inches(0.3), y + Inches(1.15), cw - Inches(0.5), Inches(1.0),
                 body, font=B_FONT, size=10, color=CBODY,
                 line_spacing=1.2)
    light_footer(s, "SKILLS", 20)


def slide_honors():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s, feature=True)
    add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
             "03 / SKILLS  ·  HONORS & AWARDS", font=H_FONT, size=11, bold=True, color=EYE)
    add_rect(s, Inches(0.7), Inches(0.85), Inches(0.4), Inches(0.04), ACC1)
    add_text(s, Inches(0.7), Inches(1.0), Inches(11), Inches(0.9),
             "Honors & Awards", font=H_FONT, size=36, bold=True, color=INK_F)
    add_text(s, Inches(0.7), Inches(1.85), Inches(11), Inches(0.4),
             "Selected recognitions across academic, open-source and competitive work.",
             font=B_FONT, size=13, italic=True, color=SUB_F)
    items = [
        ("2025", "Accepted into L2M Lab to Market with $10,000 funding"),
        ("2024", "Received Fully-Funded MSc position at York University"),
        ("2024", "Awarded additional $10,000 Mitacs funding"),
        ("2023", "500+ GitHub stars on DigiHuman open-source project"),
        ("2023", "Top 0.5% in Iran's AI Graduate National Exam"),
        ("2017", "Top 1% in Konkur, the Iranian National University Entrance Exam"),
        ("2015", "4th place at the National RoboCup competition (out of 32 teams)"),
    ]
    y = Inches(2.6)
    yr_fg = ON_DARK if DECO == "memphis" else INK
    for year, text in items:
        add_round_rect(s, Inches(0.7), y, Inches(1.5), Inches(0.55), ACC1, radius=0.3)
        add_text(s, Inches(0.7), y + Inches(0.1), Inches(1.5), Inches(0.4),
                 year, font=H_FONT, size=14, bold=True, color=yr_fg,
                 align=PP_ALIGN.CENTER)
        add_text(s, Inches(2.4), y + Inches(0.1), Inches(10.3), Inches(0.45),
                 text, font=B_FONT, size=14, color=CHEAD)
        y += Inches(0.62)
    page_footer(s, "SKILLS", 21)


# ====================================================================
# PROJECTS
# ====================================================================
def project_cover(slide, *, name, tagline, year, tags, accent=None, slide_num=0):
    if accent is None:
        accent = ACC1
    add_section_marker(slide, "04 / PROJECTS", accent)
    add_text(slide, Inches(0.7), Inches(1.05), Inches(8), Inches(0.4),
             f"PROJECT  ·  {year}", font=H_FONT, size=12, bold=True, color=accent)
    add_text(slide, Inches(0.7), Inches(1.45), Inches(11), Inches(1.1),
             name, font=H_FONT, size=44, bold=True, color=INK)
    add_text(slide, Inches(0.7), Inches(2.5), Inches(11), Inches(0.5),
             tagline, font=B_FONT, size=16, italic=True, color=SUB)


def slide_proj_digihuman_cover():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s)
    project_cover(s, name="DigiHuman",
                  tagline="Real-time 3D character animation from monocular video — open-source thesis project, 500⭐ on GitHub.",
                  year="2022 — PRESENT",
                  tags=[], slide_num=23)
    # show GitHub card image as hero
    add_rect(s, Inches(0.7), Inches(3.2), Inches(7.5), Inches(3.7), FRAME)
    add_image(s, f"{MEDIA}/image9.png", Inches(0.85), Inches(3.7), w=Inches(7.2))
    # bullets on right
    add_text(s, Inches(8.6), Inches(3.25), Inches(4), Inches(0.4),
             "STACK", font=H_FONT, size=11, bold=True, color=EYE)
    add_bullets(s, Inches(8.6), Inches(3.65), Inches(4.5), Inches(3.2),
                [
                    "[b]Unity3D[/b] high-level animation",
                    "[b]Python + Blender[/b] pipeline",
                    "[b]BlazePose[/b], [b]FaceMesh[/b], [b]MediaPipe[/b]",
                    "[b]Kalman[/b] + low-pass smoothing",
                    "30 FPS real-time throughput",
                    "C# / Python / Open Source",
                ],
                size=12, line_spacing=1.25, bullet_char="▸")
    light_footer(s, "PROJECTS", 23)


# --- DigiHuman pipeline assets + a minimal fade-build animation engine ------
DH_MEDIA = "digihuman_assets"
DH_MINT  = RGBColor(0xD4, 0xF0, 0xD5)
DH_CORAL = RGBColor(0xEE, 0x6B, 0x6B)
DH_ARROW = RGBColor(0xFC, 0xA0, 0x8E)
DH_INK   = RGBColor(0x2B, 0x2B, 0x2B)
GUARDIAN_MEDIA = "guardian_assets"
GUARD_BG = RGBColor(0x0C, 0x10, 0x19)   # dark navy matching the Guardian diagram/video
_PNS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def add_chevron(slide, x, y, w, h, fill, left=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    if left:
        shp.rotation = 180
    spPr = shp.fill._xPr
    for el in spPr.findall(qn('a:effectLst')): spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))
    return shp


def _dh_fade_transition(slide, speed="med"):
    sld = slide.element
    for el in sld.findall(qn('p:transition')): sld.remove(el)
    tr = etree.fromstring(f'<p:transition xmlns:p="{_PNS}" spd="{speed}"><p:fade/></p:transition>')
    csld = sld.find(qn('p:cSld'))
    idx = list(sld).index(csld) + 1
    while idx < len(sld) and sld[idx].tag == qn('p:clrMapOvr'): idx += 1
    sld.insert(idx, tr)


def _dh_anim_par(spid, dur_ms, offset_ms, node_type, nid, para=None):
    ctn = nid(); setid = nid(); animid = nid()
    if para is None:
        tgt = f'<p:spTgt spid="{spid}"/>'
    else:
        tgt = (f'<p:spTgt spid="{spid}"><p:txEl>'
               f'<p:pRg st="{para}" end="{para}"/></p:txEl></p:spTgt>')
    return (f'<p:par><p:cTn id="{ctn}" presetID="10" presetClass="entr" presetSubtype="0" '
            f'fill="hold" grpId="0" nodeType="{node_type}">'
            f'<p:stCondLst><p:cond delay="{offset_ms}"/></p:stCondLst><p:childTnLst>'
            f'<p:set><p:cBhvr><p:cTn id="{setid}" dur="1" fill="hold">'
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
            f'<p:tgtEl>{tgt}</p:tgtEl>'
            f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst></p:cBhvr>'
            f'<p:to><p:strVal val="visible"/></p:to></p:set>'
            f'<p:anim calcmode="lin" valueType="num"><p:cBhvr additive="base">'
            f'<p:cTn id="{animid}" dur="{dur_ms}" fill="hold"/>'
            f'<p:tgtEl>{tgt}</p:tgtEl>'
            f'<p:attrNameLst><p:attrName>style.opacity</p:attrName></p:attrNameLst></p:cBhvr>'
            f'<p:tavLst><p:tav tm="0"><p:val><p:fltVal val="0"/></p:val></p:tav>'
            f'<p:tav tm="100000"><p:val><p:fltVal val="1"/></p:val></p:tav></p:tavLst></p:anim>'
            f'</p:childTnLst></p:cTn></p:par>')


def _dh_build_anim(slide, *, dur_ms=400):
    """Grouped click reveals — the page frame (title, decorations) is static;
    each box (tile / pill / bar) fades in TOGETHER with its text & image on a
    single click. Shapes before the first box stay static (present on enter)."""
    def container(sp):
        try:
            return (sp.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE
                    and int(sp.width or 0) >= int(Inches(1.5)))
        except Exception:
            return False
    def center_in(sp, c):
        try:
            cx = int(sp.left or 0) + int(sp.width or 0) // 2
            cy = int(sp.top or 0) + int(sp.height or 0) // 2
            return (int(c.left) <= cx <= int(c.left) + int(c.width)
                    and int(c.top) <= cy <= int(c.top) + int(c.height))
        except Exception:
            return False
    shapes = list(slide.shapes)
    first = next((i for i, sp in enumerate(shapes) if container(sp)), None)
    if first is None:
        return
    groups, cur, cont = [], [], None
    for sp in shapes[first:]:
        try:
            if (int(sp.left or 0) == 0 and int(sp.top or 0) == 0
                    and int(sp.width or 0) == int(SLIDE_W)
                    and int(sp.height or 0) == int(SLIDE_H)):
                continue
        except Exception:
            pass
        if container(sp) and not (cont is not None and center_in(sp, cont)):
            if cur:
                groups.append(cur)
            cur = [sp.shape_id]; cont = sp
        else:
            cur.append(sp.shape_id)
    if cur:
        groups.append(cur)
    if not groups:
        return
    counter = [3]
    def nid():
        v = counter[0]; counter[0] += 1; return v
    steps = []
    builds_list = []
    for g in groups:
        click = nid(); inner = nid()
        anims = []
        for j, spid in enumerate(g):
            nt = "clickEffect" if j == 0 else "withEffect"
            anims.append(_dh_anim_par(spid, dur_ms, 0, nt, nid))
            builds_list.append(f'<p:bldP spid="{spid}" grpId="0"/>')
        steps.append(
            f'<p:par><p:cTn id="{click}" fill="hold">'
            f'<p:stCondLst><p:cond delay="indefinite"/></p:stCondLst><p:childTnLst>'
            f'<p:par><p:cTn id="{inner}" fill="hold">'
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst>{"".join(anims)}</p:childTnLst></p:cTn></p:par>'
            f'</p:childTnLst></p:cTn></p:par>')
    builds = "".join(builds_list)
    xml = (f'<p:timing xmlns:p="{_PNS}"><p:tnLst><p:par>'
           f'<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>'
           f'<p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>'
           f'{"".join(steps)}</p:childTnLst></p:cTn>'
           f'<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
           f'<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
           f'</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst>'
           f'<p:bldLst>{builds}</p:bldLst></p:timing>')
    sld = slide.element
    for el in sld.findall(qn('p:timing')): sld.remove(el)
    sld.append(etree.fromstring(xml))


def _apply_auto_entrance(slide, *, dur_ms=450):
    """Non-bullet content fades in automatically when the slide opens; a bullet
    list (text box tagged 'BULLETS') is NOT auto-chained — each bullet reveals
    on its OWN click. Skips the full-slide background and the footer strip."""
    shapes = []
    for sp in slide.shapes:
        try:
            l = int(sp.left or 0); t = int(sp.top or 0)
            w = int(sp.width or 0); h = int(sp.height or 0)
        except Exception:
            l = t = w = h = 0
        if l == 0 and t == 0 and w == int(SLIDE_W) and h == int(SLIDE_H):
            continue
        if t >= int(Inches(7.0)):
            continue
        shapes.append(sp)
    if not shapes:
        return

    def is_bullet(sp):
        try:
            return sp.name == "BULLETS"
        except Exception:
            return False

    def npar(sp):
        try:
            return max(1, len(list(sp.text_frame.paragraphs)))
        except Exception:
            return 1

    counter = [3]
    def nid():
        v = counter[0]; counter[0] += 1; return v

    def step_par(anims, on_click):
        cid = nid(); iid = nid()
        delay = "indefinite" if on_click else "0"
        return (f'<p:par><p:cTn id="{cid}" fill="hold">'
                f'<p:stCondLst><p:cond delay="{delay}"/></p:stCondLst><p:childTnLst>'
                f'<p:par><p:cTn id="{iid}" fill="hold">'
                f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
                f'<p:childTnLst>{"".join(anims)}</p:childTnLst></p:cTn></p:par>'
                f'</p:childTnLst></p:cTn></p:par>')

    steps_xml, whole_ids, para_ids = [], [], []
    # 1) auto step — all non-bullet content fades in on slide enter
    auto_anims, first = [], True
    for sp in shapes:
        if not is_bullet(sp):
            nt = "afterEffect" if first else "withEffect"
            first = False
            auto_anims.append(_dh_anim_par(sp.shape_id, dur_ms, 0, nt, nid))
            whole_ids.append(sp.shape_id)
    if auto_anims:
        steps_xml.append(step_par(auto_anims, on_click=False))
    # 2) one On-Click step per bullet (paragraph)
    for sp in shapes:
        if is_bullet(sp):
            para_ids.append(sp.shape_id)
            for k in range(npar(sp)):
                anim = _dh_anim_par(sp.shape_id, dur_ms, 0, "clickEffect", nid, para=k)
                steps_xml.append(step_par([anim], on_click=True))
    if not steps_xml:
        return
    builds = ("".join(f'<p:bldP spid="{s}" grpId="0" build="p"/>' for s in para_ids)
              + "".join(f'<p:bldP spid="{s}" grpId="0"/>' for s in whole_ids))
    xml = (f'<p:timing xmlns:p="{_PNS}"><p:tnLst><p:par>'
           f'<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>'
           f'<p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>'
           f'{"".join(steps_xml)}</p:childTnLst></p:cTn>'
           f'<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
           f'<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
           f'</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst>'
           f'<p:bldLst>{builds}</p:bldLst></p:timing>')
    sld = slide.element
    for el in sld.findall(qn('p:timing')):
        sld.remove(el)
    sld.append(etree.fromstring(xml))


def slide_proj_digihuman_arch():
    """DigiHuman pipeline — recreated from the template's page 34, animated."""
    s = prs.slides.add_slide(blank_layout)
    fill_slide(s, WHITE)

    # decorative template-style blobs (behind content; fade in first)
    add_oval(s, Inches(11.95), Inches(-0.75), Inches(2.1), Inches(2.1), M_PURPLE)
    add_oval(s, Inches(12.55), Inches(3.15), Inches(0.85), Inches(0.85), DH_CORAL)
    add_oval(s, Inches(-0.55), Inches(6.35), Inches(1.5), Inches(1.5), M_CYAN)
    add_oval(s, Inches(0.40), Inches(0.25), Inches(0.26), Inches(0.26), M_YELLOW)

    add_text(s, Inches(0.5), Inches(0.14), Inches(12.33), Inches(0.7),
             "DigiHuman", font=H_FONT, size=34, bold=True, color=DH_INK,
             align=PP_ALIGN.CENTER)

    def tile(x, y, w, h):
        add_round_rect(s, Inches(x), Inches(y), Inches(w), Inches(h), DH_MINT, radius=0.09)

    def pill(x, y, w, h, text, size=12, r=0.5):
        add_round_rect(s, Inches(x), Inches(y), Inches(w), Inches(h), DH_CORAL, radius=r)
        add_text(s, Inches(x), Inches(y), Inches(w), Inches(h), text, font=H_FONT,
                 size=size, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

    def label(x, y, w, text):
        add_text(s, Inches(x), Inches(y), Inches(w), Inches(0.38), text, font=H_FONT,
                 size=13, bold=True, color=DH_INK, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

    # ----- TOP ROW: Video -> Server -----
    tile(0.93, 0.85, 2.30, 2.45)
    add_image(s, f"{DH_MEDIA}/dh_01.jpg", Inches(1.10), Inches(1.02), w=Inches(1.94))
    label(0.93, 2.86, 2.30, "Video")
    tile(3.75, 0.85, 2.30, 2.45)
    add_image(s, f"{DH_MEDIA}/dh_00.png", Inches(3.93), Inches(1.00), w=Inches(1.95))
    label(3.75, 2.86, 2.30, "Server")
    add_chevron(s, Inches(3.20), Inches(1.55), Inches(0.60), Inches(0.48), DH_ARROW)
    pill(2.50, 3.42, 1.95, 0.46, "Upload")

    # ----- Video Processing -> branch to the three landmark bars -----
    pill(5.71, 1.74, 2.34, 0.55, "Video Processing")
    add_rect(s, Inches(8.05), Inches(1.03), Inches(0.035), Inches(2.05), DH_ARROW)
    bars = [("Facial landmarks", 0.68, "dh_04.png", 10.72, 0.50, 0.89),
            ("Body Pose landmarks", 1.73, "dh_03.png", 10.62, 1.52, 0.94),
            ("Hand landmarks", 2.72, "dh_05.png", 10.61, 2.56, 0.89)]
    for txt, by, icon, ix, iy, isz in bars:
        add_rect(s, Inches(8.05), Inches(by + 0.33), Inches(0.30), Inches(0.035), DH_ARROW)
        add_round_rect(s, Inches(8.24), Inches(by), Inches(3.27), Inches(0.70), DH_MINT, radius=0.3)
        add_text(s, Inches(8.45), Inches(by), Inches(2.05), Inches(0.70), txt,
                 font=H_FONT, size=12, bold=True, color=DH_INK, anchor=MSO_ANCHOR.MIDDLE)
        add_image(s, f"{DH_MEDIA}/{icon}", Inches(ix), Inches(iy), w=Inches(isz))
    pill(8.95, 3.50, 1.95, 0.62, "Receive data from server", size=10, r=0.25)

    # ----- BOTTOM ROW (right -> left): Processing -> Smoothing -> Rendering -----
    tile(8.70, 3.86, 2.55, 2.55)
    add_image(s, f"{DH_MEDIA}/dh_02.png", Inches(9.10), Inches(4.05), w=Inches(1.80))
    pill(8.50, 6.42, 2.75, 0.95,
         "Processing 3D landmarks, calculating rotation & position of body joints",
         size=8.5, r=0.12)
    add_chevron(s, Inches(8.15), Inches(4.92), Inches(0.55), Inches(0.50), DH_ARROW, left=True)
    tile(5.53, 3.86, 2.60, 2.55)
    add_image(s, f"{DH_MEDIA}/dh_06.png", Inches(5.70), Inches(4.00), w=Inches(2.26))
    pill(5.10, 6.42, 3.20, 0.80,
         "Smoothing animation in frames, using signal filters (low-pass filter)",
         size=9, r=0.14)
    add_chevron(s, Inches(4.85), Inches(4.95), Inches(0.55), Inches(0.50), DH_ARROW, left=True)
    tile(2.26, 3.86, 2.55, 2.55)
    add_image(s, f"{DH_MEDIA}/dh_07.png", Inches(2.55), Inches(4.05), w=Inches(1.96))
    pill(2.00, 6.42, 3.00, 0.60, "Rendering the final animation", size=11, r=0.2)

    _dh_fade_transition(s)
    _dh_build_anim(s)


def slide_proj_digihuman_demos():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s, feature=True)
    add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
             "04 / PROJECTS  ·  DIGIHUMAN", font=H_FONT, size=11, bold=True, color=EYE)
    add_rect(s, Inches(0.7), Inches(0.85), Inches(0.4), Inches(0.04), ACC1)
    add_text(s, Inches(0.7), Inches(1.0), Inches(11), Inches(0.7),
             "Live in action", font=H_FONT, size=32, bold=True, color=INK_F)
    add_text(s, Inches(0.7), Inches(1.65), Inches(11), Inches(0.4),
             "Real-time face, hand and full-body capture transferred to a 3D avatar.",
             font=B_FONT, size=13, italic=True, color=SUB_F)
    # 2 demo images side by side
    add_image(s, f"{MEDIA}/image25.gif", Inches(0.7), Inches(2.4), w=Inches(6))
    add_image(s, f"{MEDIA}/image30.gif", Inches(7.0), Inches(2.4), w=Inches(5.6))
    add_text(s, Inches(0.7), Inches(6.0), Inches(6), Inches(0.4),
             "Facial & body landmark tracking", font=B_FONT, size=12, italic=True, color=SUB_F)
    add_text(s, Inches(7.0), Inches(6.0), Inches(5.6), Inches(0.4),
             "3D avatar animation pipeline", font=B_FONT, size=12, italic=True, color=SUB_F)
    page_footer(s, "PROJECTS", 25)


def slide_proj_safezone_cover():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s)
    project_cover(s, name="SAFEZone AI",
                  tagline="Boson AI × MScAC hackathon — a 3D therapist avatar that listens, reasons and responds in a cloned voice.",
                  year="2025  ·  SEMI-FINALIST",
                  tags=[], slide_num=26)
    add_text(s, Inches(0.7), Inches(3.4), Inches(12), Inches(0.4),
             "WHAT IT IS", font=H_FONT, size=11, bold=True, color=EYE)
    bullets = [
        "[b]Unity Frontend[/b]: lifelike 3D therapist avatar with natural expressions and lip-sync.",
        "[b]Higgs Audio Understanding (ASR)[/b]: transcribes user speech and detects emotional tone.",
        "[b]GPT-4o (LLM)[/b]: combines transcript with onboarding questionnaire to generate empathetic replies.",
        "[b]Higgs Audio Generation (TTS + Voice Cloning)[/b]: converts response to natural speech and lip-sync.",
        "Result: an emotionally-aligned conversation that feels real and supportive.",
    ]
    add_bullets(s, Inches(0.7), Inches(3.8), Inches(12), Inches(3.5),
                bullets, size=14, line_spacing=1.3, bullet_char="▸",
                color=BODY)
    light_footer(s, "PROJECTS", 26)


def slide_proj_safezone_arch():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s, feature=True)
    add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
             "04 / PROJECTS  ·  SAFEZONE AI", font=H_FONT, size=11, bold=True, color=EYE)
    add_rect(s, Inches(0.7), Inches(0.85), Inches(0.4), Inches(0.04), ACC1)
    add_text(s, Inches(0.7), Inches(1.0), Inches(11), Inches(0.7),
             "How it works", font=H_FONT, size=32, bold=True, color=INK_F)
    add_text(s, Inches(0.7), Inches(1.65), Inches(11), Inches(0.4),
             "Voice  →  ASR + emotion  →  LLM grounded in APA PsycInfo  →  cloned-voice TTS  →  blend-shape lip-sync.",
             font=B_FONT, size=13, italic=True, color=SUB_F)
    if DECO == "memphis":
        # image26.png is a dark navy-background diagram; seat it in a dark
        # "exhibit" frame so it reads as deliberate on a white slide.
        path = f"{MEDIA}/image26.png"
        img_w = 10.6
        ix = (13.333 - img_w) / 2.0
        iy = 2.45
        img_h = img_w * _png_ratio(path)
        pad = 0.16
        add_round_rect(s, Inches(ix - pad), Inches(iy - pad),
                       Inches(img_w + 2 * pad), Inches(img_h + 2 * pad),
                       NAVY, radius=0.03)
        add_image(s, path, Inches(ix), Inches(iy), w=Inches(img_w))
    else:
        add_image(s, f"{MEDIA}/image26.png", Inches(0.7), Inches(2.4), w=Inches(12))
    page_footer(s, "PROJECTS", 27)


def slide_proj_caselogic():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s)
    project_cover(s, name="CaseLogic",
                  tagline="3rd-place finisher — source-grounded legal research assistant built in 24 hours at the Toronto Legal-Tech Hackathon.",
                  year="2026  ·  3RD PLACE",
                  tags=[], slide_num=28)
    # left bullets, right architecture sketch
    add_text(s, Inches(0.7), Inches(3.3), Inches(7), Inches(0.4),
             "HIGHLIGHTS", font=H_FONT, size=11, bold=True, color=EYE)
    bullets = [
        "[b]Hybrid retrieval[/b]: Chroma vector DB + SQLite FTS5 keyword search for robust legal lookups.",
        "[b]Per-claim verification[/b]: every assertion is grounded back to a source paragraph.",
        "[b]Multi-agent planning workspace[/b]: orchestrated agents to research, draft and verify.",
        "[b]Demo[/b] in 24h: shipped end-to-end UI + multi-agent backend.",
        "[b]Result[/b]: placed 3rd at the Legal-Tech Hackathon (Toronto, 2026).",
    ]
    add_bullets(s, Inches(0.7), Inches(3.7), Inches(7.7), Inches(3.3),
                bullets, size=12, line_spacing=1.25, bullet_char="▸",
                color=BODY)
    # right column: tech stack pills
    add_rect(s, Inches(8.8), Inches(3.3), Inches(4), Inches(3.5), CARD)
    add_text(s, Inches(9.0), Inches(3.4), Inches(3.6), Inches(0.4),
             "STACK", font=H_FONT, size=11, bold=True, color=EYE)
    techs = ["Chroma", "SQLite FTS5", "LangChain", "OpenAI", "Multi-Agent", "Python", "FastAPI", "RAG", "Hybrid Search"]
    cx = Inches(9.0)
    cy = Inches(3.85)
    for t in techs:
        cw = Inches(0.3 + 0.09 * len(t))
        if cx + cw > Inches(12.7):
            cx = Inches(9.0)
            cy += Inches(0.45)
        add_round_rect(s, cx, cy, cw, Inches(0.35), CHIP_BG, radius=0.4)
        add_text(s, cx, cy + Inches(0.04), cw, Inches(0.3),
                 t, font=B_FONT, size=9, bold=True, color=CHIP_FG,
                 align=PP_ALIGN.CENTER)
        cx += cw + Inches(0.1)
    light_footer(s, "PROJECTS", 28)


def slide_proj_guardian_arch():
    """Guardian — full-bleed System Architecture exhibit (dark, like the diagram)."""
    s = prs.slides.add_slide(blank_layout)
    fill_slide(s, GUARD_BG)
    # architecture PNG is 2400x1500 (ratio 1.6) -> fit to full height, centred
    s.shapes.add_picture(f"{GUARDIAN_MEDIA}/guardian-architecture.png",
                         Inches(0.667), Inches(0.0), Inches(12.0), Inches(7.5))


def slide_proj_guardian_demo():
    """Guardian — embedded 60s live-demo video on a dark exhibit slide."""
    s = prs.slides.add_slide(blank_layout)
    fill_slide(s, GUARD_BG)
    s.shapes.add_movie(
        f"{GUARDIAN_MEDIA}/guardian-combined-demo.mp4",
        Inches(0.917), Inches(0.34), Inches(11.5), Inches(6.469),
        poster_frame_image=f"{GUARDIAN_MEDIA}/guardian-demo-poster.png",
        mime_type="video/mp4")
    add_text(s, Inches(0.5), Inches(6.95), Inches(12.33), Inches(0.4),
             "Live demo  ·  watch detects a fall  →  risk escalates to CRITICAL"
             "  →  autonomous 911 call with full medical context",
             font=B_FONT, size=11, color=RGBColor(0x9A, 0xA6, 0xB8),
             align=PP_ALIGN.CENTER)


def slide_proj_latex_cv():
    """LaTeX CV Builder — animated pipeline exhibit (DigiHuman-style).
    One job description in -> a tailored CV, cover letter & fit score out.
    Mirrors the real repo: Obsidian vault + JD -> 5-stage cv_pipeline -> PDFs."""
    s = prs.slides.add_slide(blank_layout)
    fill_slide(s, M_WHITE)

    # --- decorative blobs (static page frame) ---
    add_oval(s, Inches(11.98), Inches(-0.80), Inches(2.05), Inches(2.05), M_PURPLE)
    add_oval(s, Inches(12.63), Inches(2.92), Inches(0.74), Inches(0.74), M_CYAN)
    add_oval(s, Inches(-0.55), Inches(6.40), Inches(1.45), Inches(1.45), M_YELLOW)
    add_oval(s, Inches(0.45), Inches(0.34), Inches(0.24), Inches(0.24), DH_CORAL)

    # --- title + subtitle (static) ---
    add_text(s, Inches(0.5), Inches(0.10), Inches(12.33), Inches(0.54),
             "LaTeX CV Builder", font=H_FONT, size=30, bold=True,
             color=DH_INK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(0.64), Inches(12.33), Inches(0.34),
             "One job description in   →   a tailored CV, cover letter & fit score out",
             font=B_FONT, size=12.5, italic=True, color=M_TEAL,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ---- nested helpers ----
    def card(x, y, w, h, fill, *, radius=0.10, line=None):
        return add_round_rect(s, Inches(x), Inches(y), Inches(w), Inches(h),
                              fill, radius=radius, line=line)

    def t(x, y, w, h, text, *, size=10, bold=False, color=DH_INK,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=None, italic=False,
          ls=None):
        add_text(s, Inches(x), Inches(y), Inches(w), Inches(h), text,
                 font=font or H_FONT, size=size, bold=bold, color=color,
                 align=align, anchor=anchor, italic=italic, line_spacing=ls)

    def badge(x, y, w, h, text, fill, size=8.5):
        add_round_rect(s, Inches(x), Inches(y), Inches(w), Inches(h), fill, radius=0.5)
        add_text(s, Inches(x), Inches(y), Inches(w), Inches(h), text, font=H_FONT,
                 size=size, bold=True, color=M_WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

    def arrow(x, y, w=0.52, h=0.56, left=False, down=False):
        c = add_chevron(s, Inches(x), Inches(y), Inches(w), Inches(h),
                        DH_ARROW, left=left)
        if down:
            c.rotation = 90
        return c

    TOP_Y, BOT_Y, TH = 1.62, 3.92, 1.88

    # ---- SOURCES lane (read every run) ----
    card(0.40, 1.06, 12.13, 0.46, M_CARD, radius=0.22, line=M_CYAN)
    t(0.62, 1.06, 2.0, 0.46, "READS EVERY RUN", size=8.5, bold=True,
      color=M_PURPLE, anchor=MSO_ANCHOR.MIDDLE)
    t(2.78, 1.06, 4.6, 0.46, "Obsidian vault  ·  experiences / projects / skills",
      size=9.5, bold=True, color=M_TEAL, anchor=MSO_ANCHOR.MIDDLE)
    t(7.55, 1.06, 2.3, 0.46, "canonical CV .tex", size=9.5, bold=True,
      color=DH_INK, anchor=MSO_ANCHOR.MIDDLE)
    t(10.0, 1.06, 2.4, 0.46, "referrals.db (SQLite)", size=9.5, bold=True,
      color=DH_INK, anchor=MSO_ANCHOR.MIDDLE)

    # ---- INPUT: Job Description (the trigger) ----
    card(0.40, TOP_Y, 1.85, TH, M_PURPLE)
    t(0.40, TOP_Y + 0.18, 1.85, 0.28, "INPUT", size=8.5, bold=True,
      color=M_WHITE, align=PP_ALIGN.CENTER)
    t(0.45, TOP_Y + 0.55, 1.75, 0.72, "Job Description", size=15, bold=True,
      color=M_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, ls=1.0)
    t(0.45, TOP_Y + 1.30, 1.75, 0.48, "paste · URL · Gmail queue", size=8.5,
      color=M_WHITE, align=PP_ALIGN.CENTER, ls=1.0)
    arrow(2.27, TOP_Y + 0.66)

    # ---- the 5-stage pipeline (LLM only in stages 1 & 3) ----
    def stage(x, y, num, title, module, desc, out, badge_txt, badge_fill,
              out_color=DH_CORAL):
        w = 2.85
        card(x, y, w, TH, DH_MINT)
        add_oval(s, Inches(x + 0.15), Inches(y + 0.16), Inches(0.44),
                 Inches(0.44), DH_CORAL)
        t(x + 0.15, y + 0.16, 0.44, 0.44, num, size=15, bold=True, color=M_WHITE,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        t(x + 0.68, y + 0.15, w - 0.84, 0.46, title, size=14, bold=True,
          color=DH_INK, anchor=MSO_ANCHOR.MIDDLE)
        t(x + 0.18, y + 0.72, w - 0.36, 0.28, module, size=9.5, bold=True,
          color=M_TEAL, font=B_FONT)
        t(x + 0.18, y + 1.01, w - 0.36, 0.52, desc, size=9, color=DH_INK, ls=1.04)
        t(x + 0.18, y + 1.54, w - 1.15, 0.30, out, size=10, bold=True,
          color=out_color, anchor=MSO_ANCHOR.MIDDLE)
        bw = 0.34 + 0.082 * len(badge_txt)
        badge(x + w - bw - 0.14, y + 1.55, bw, 0.30, badge_txt, badge_fill)

    # top row: 1 Parse JD -> 2 Select -> 3 Generate
    stage(2.88, TOP_Y, "1", "Parse JD", "job_parser.py",
          "LLM reads the JD into domain, role-fit, tags & constraints",
          "→ JobTargets", "LLM", DH_CORAL)
    arrow(5.78, TOP_Y + 0.66)
    stage(6.37, TOP_Y, "2", "Select", "candidate_selector.py",
          "Rule-based scoring of vault notes — anchors & de-dupe",
          "→ shortlist", "no LLM", M_TEAL)
    arrow(9.27, TOP_Y + 0.66)
    stage(9.86, TOP_Y, "3", "Generate", "bullet · cover · fit",
          "One LLM call writes bullets, a cover letter & a fit score",
          "→ bullets + letter", "LLM", DH_CORAL)
    # serpentine turn: stage 3 (top) -> stage 4 (bottom)
    arrow(11.05, 3.46, w=0.46, h=0.52, down=True)

    # bottom row (right -> left): 4 Render LaTeX -> 5 Compile
    stage(9.86, BOT_Y, "4", "Render LaTeX", "latex_formatter.py",
          "Fills résumé + per-entry .tex from the generation",
          "→ resume.tex", ".tex", M_PURPLE)
    arrow(9.27, BOT_Y + 0.66, left=True)
    stage(6.37, BOT_Y, "5", "Compile", "build_runner.py",
          "pdflatex builds the final PDF (skippable)",
          "→ CV.pdf", "pdflatex", M_TEAL)
    arrow(5.78, BOT_Y + 0.66, left=True)

    # ---- OUTPUTS panel ----
    card(0.40, BOT_Y, 5.30, TH, M_CARD, radius=0.10, line=M_GREEN)
    t(0.62, BOT_Y + 0.12, 3.0, 0.30, "OUTPUTS", size=10.5, bold=True,
      color=M_PURPLE)
    outs = [("CV.pdf", "tailored résumé", DH_CORAL),
            ("Cover letter .pdf", "matched tone & structure", M_TEAL),
            ("Fit score", "0–100 with strengths & gaps", M_GREEN),
            ("Archived", "JSON · SQLite · Obsidian note", M_PURPLE)]
    for i, (name, sub, acc) in enumerate(outs):
        ry = BOT_Y + 0.50 + i * 0.335
        add_oval(s, Inches(0.66), Inches(ry + 0.05), Inches(0.15),
                 Inches(0.15), acc)
        t(0.92, ry, 1.95, 0.30, name, size=10.5, bold=True, color=DH_INK,
          anchor=MSO_ANCHOR.MIDDLE)
        t(2.90, ry, 2.70, 0.30, sub, size=8.7, color=DH_INK,
          anchor=MSO_ANCHOR.MIDDLE)

    # ---- backend swap + UI (captions) ----
    card(1.95, 6.04, 9.45, 0.48, M_PURPLE, radius=0.24)
    t(1.95, 6.04, 9.45, 0.48,
      "LLM backend (swappable):   api_llm — LangChain → GPT-4o / Ollama"
      "      ·      claude_code — claude CLI",
      size=9.5, bold=True, color=M_WHITE, align=PP_ALIGN.CENTER,
      anchor=MSO_ANCHOR.MIDDLE)
    t(1.95, 6.60, 9.45, 0.32,
      "CV-UI · Next.js wraps the CLI — paste a JD and watch Step 1/5 … 5/5 stream live",
      size=10, bold=True, color=M_TEAL, align=PP_ALIGN.CENTER,
      anchor=MSO_ANCHOR.MIDDLE)

    _dh_fade_transition(s)
    _dh_build_anim(s)


def slide_proj_techu():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s, feature=True)
    add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
             "04 / PROJECTS  ·  TECHU", font=H_FONT, size=11, bold=True, color=EYE)
    add_rect(s, Inches(0.7), Inches(0.85), Inches(0.4), Inches(0.04), ACC1)
    add_text(s, Inches(0.7), Inches(1.0), Inches(11), Inches(0.9),
             "Techu  —  the AI-driven card game", font=H_FONT, size=32, bold=True, color=INK_F)
    add_text(s, Inches(0.7), Inches(1.95), Inches(11), Inches(0.4),
             "Cross-platform strategic card game with DQN + MCTS opponents and full online multiplayer.",
             font=B_FONT, size=13, italic=True, color=SUB_F)
    add_image(s, f"{MEDIA}/image18.png", Inches(0.7), Inches(2.6), w=Inches(7.0))
    # right metrics
    cards = [
        ("DQN + MCTS", "AI decision-making"),
        ("Photon + PlayFab", "Online multiplayer + Azure matchmaking"),
        ("Firebase", "Crash analytics & diagnostics"),
        ("Cross-platform", "iOS  ·  Android  ·  Desktop"),
    ]
    cy = Inches(2.6)
    for n, lbl in cards:
        add_round_rect(s, Inches(8.4), cy, Inches(4.4), Inches(0.95), CARD2, radius=0.1,
                       line=cline(ACC1))
        add_text(s, Inches(8.6), cy + Inches(0.1), Inches(4), Inches(0.4),
                 n, font=H_FONT, size=15, bold=True, color=ACC1)
        add_text(s, Inches(8.6), cy + Inches(0.5), Inches(4), Inches(0.4),
                 lbl, font=B_FONT, size=11, color=CBODY)
        cy += Inches(1.05)
    page_footer(s, "PROJECTS", 30)


def slide_proj_hypervigilance():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s, feature=True)
    add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
             "04 / PROJECTS  ·  HYPERVIGILANCE", font=H_FONT, size=11, bold=True, color=EYE)
    add_rect(s, Inches(0.7), Inches(0.85), Inches(0.4), Inches(0.04), ACC1)
    add_text(s, Inches(0.7), Inches(1.0), Inches(11), Inches(0.9),
             "HYPERVIGILANCE", font=H_FONT, size=36, bold=True, color=INK_F)
    add_text(s, Inches(0.7), Inches(1.85), Inches(11), Inches(0.5),
             "A psychological horror short — atmospheric Unity scenes + generative video pipeline.",
             font=B_FONT, size=13, italic=True, color=SUB_F)
    # 2x2 grid
    add_image(s, f"{MEDIA}/image35.png", Inches(0.7), Inches(2.6), w=Inches(6))
    add_image(s, f"{MEDIA}/image36.png", Inches(6.9), Inches(2.6), w=Inches(5.9))
    add_image(s, f"{MEDIA}/image29.png", Inches(0.7), Inches(4.85), w=Inches(6))
    add_image(s, f"{MEDIA}/image31.png", Inches(6.9), Inches(4.85), w=Inches(5.9))
    page_footer(s, "PROJECTS", 31)


def slide_proj_backgammon():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s)
    add_section_marker(s, "04 / PROJECTS")
    add_text(s, Inches(0.7), Inches(1.05), Inches(8), Inches(0.4),
             "PROJECT  ·  2021", font=H_FONT, size=12, bold=True, color=EYE)
    add_text(s, Inches(0.7), Inches(1.45), Inches(11), Inches(1.0),
             "Backgammon 3D",
             font=H_FONT, size=40, bold=True, color=INK)
    add_text(s, Inches(0.7), Inches(2.5), Inches(11), Inches(0.4),
             "Turn-based multiplayer with matchmaking and Monte-Carlo AI opponents.",
             font=B_FONT, size=14, italic=True, color=SUB)
    add_image(s, f"{MEDIA}/image49.png", Inches(0.7), Inches(3.2), w=Inches(7.8))
    # right bullets
    add_text(s, Inches(8.7), Inches(3.2), Inches(4), Inches(0.4),
             "HIGHLIGHTS", font=H_FONT, size=11, bold=True, color=EYE)
    add_bullets(s, Inches(8.7), Inches(3.6), Inches(4.2), Inches(3.2),
                [
                    "Online matchmaking system",
                    "[b]Monte Carlo[/b] AI decision-making",
                    "Logging + replay for analysis",
                    "Multi-platform deployment",
                ],
                size=12, line_spacing=1.3, bullet_char="▸",
                color=BODY)
    light_footer(s, "PROJECTS", 32)


def slide_proj_cv():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s)
    add_section_marker(s, "04 / PROJECTS")
    add_text(s, Inches(0.7), Inches(1.05), Inches(11), Inches(0.9),
             "Computer Vision projects", font=H_FONT, size=32, bold=True, color=INK)
    add_text(s, Inches(0.7), Inches(1.85), Inches(11), Inches(0.4),
             "SIFT, RANSAC, homography, seam carving, 3D reconstruction — from-scratch in Python.",
             font=B_FONT, size=13, italic=True, color=SUB)
    # 2 image panels with captions
    add_rect(s, Inches(0.7), Inches(2.5), Inches(6.0), Inches(4.3), FRAME)
    add_image(s, f"{MEDIA}/image43.png", Inches(0.95), Inches(2.85), h=Inches(2.5))
    add_text(s, Inches(0.95), Inches(5.5), Inches(5.6), Inches(0.45),
             "Panorama Reconstruction", font=H_FONT, size=15, bold=True, color=ACC1)
    add_text(s, Inches(0.95), Inches(5.9), Inches(5.6), Inches(0.4),
             "SIFT + RANSAC homography for seamless stitching",
             font=B_FONT, size=11, color=CBODY)

    add_rect(s, Inches(6.95), Inches(2.5), Inches(5.85), Inches(4.3), FRAME)
    add_image(s, f"{MEDIA}/image53.png", Inches(7.6), Inches(2.85), h=Inches(2.5))
    add_text(s, Inches(7.15), Inches(5.5), Inches(5.5), Inches(0.45),
             "3D Object Reconstruction", font=H_FONT, size=15, bold=True, color=ACC1)
    add_text(s, Inches(7.15), Inches(5.9), Inches(5.5), Inches(0.4),
             "SIFT → RANSAC → PoinTr point-cloud up-sampling",
             font=B_FONT, size=11, color=CBODY)
    light_footer(s, "PROJECTS", 33)


def slide_proj_ai_viz():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s, feature=True)
    add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
             "04 / PROJECTS", font=H_FONT, size=11, bold=True, color=EYE)
    add_rect(s, Inches(0.7), Inches(0.85), Inches(0.4), Inches(0.04), ACC1)
    add_text(s, Inches(0.7), Inches(1.0), Inches(11), Inches(0.9),
             "AI Algorithms Visualization", font=H_FONT, size=32, bold=True, color=INK_F)
    add_text(s, Inches(0.7), Inches(1.85), Inches(11), Inches(0.4),
             "Educational animations of search & graph algorithms with Manim.",
             font=B_FONT, size=13, italic=True, color=SUB_F)
    add_image(s, f"{MEDIA}/image48.jpg", Inches(0.7), Inches(2.4), w=Inches(8.5))
    add_text(s, Inches(9.5), Inches(2.6), Inches(3.5), Inches(0.4),
             "HIGHLIGHTS", font=H_FONT, size=11, bold=True, color=EYE)
    add_bullets(s, Inches(9.5), Inches(3.0), Inches(3.5), Inches(3.5),
                [
                    "Built with [b]Manim[/b]",
                    "Iterative deepening, A*, BFS/DFS",
                    "Frontier vs explored state",
                    "Designed for teaching",
                ],
                size=11, line_spacing=1.3, bullet_char="▸",
                color=CHEAD)
    page_footer(s, "PROJECTS", 34)


def slide_proj_search():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s)
    add_section_marker(s, "04 / PROJECTS")
    add_text(s, Inches(0.7), Inches(1.05), Inches(11), Inches(0.9),
             "Search Engine from Scratch", font=H_FONT, size=32, bold=True, color=INK)
    add_text(s, Inches(0.7), Inches(1.85), Inches(11), Inches(0.4),
             "tf-idf  ·  champion lists  ·  KNN  ·  K-means — Python implementation tuned for latency.",
             font=B_FONT, size=13, italic=True, color=SUB)
    add_image(s, f"{MEDIA}/image34.png", Inches(0.7), Inches(2.6), w=Inches(6.5))
    add_text(s, Inches(7.6), Inches(2.6), Inches(5.2), Inches(0.4),
             "TECHNIQUES", font=H_FONT, size=11, bold=True, color=EYE)
    add_bullets(s, Inches(7.6), Inches(3.0), Inches(5.2), Inches(3.8),
                [
                    "Inverted index + champion lists",
                    "tf-idf scoring",
                    "[b]KNN[/b] and [b]K-means[/b] clustering",
                    "Aggressive pre-processing & normalization",
                    "Open source on GitHub",
                ],
                size=12, line_spacing=1.3, bullet_char="▸",
                color=BODY)
    light_footer(s, "PROJECTS", 35)


def slide_proj_solar():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s, feature=True)
    add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
             "04 / PROJECTS", font=H_FONT, size=11, bold=True, color=EYE)
    add_rect(s, Inches(0.7), Inches(0.85), Inches(0.4), Inches(0.04), ACC1)
    add_text(s, Inches(0.7), Inches(1.0), Inches(11), Inches(0.9),
             "Solar System Simulation", font=H_FONT, size=32, bold=True, color=INK_F)
    add_text(s, Inches(0.7), Inches(1.85), Inches(11), Inches(0.4),
             "C++ · OpenGL — low-level graphics with rendering-pipeline optimizations.",
             font=B_FONT, size=13, italic=True, color=SUB_F)
    add_image(s, f"{MEDIA}/image56.gif", Inches(3.7), Inches(2.5), h=Inches(3.6))
    # tags below image
    tags = ["C++", "OpenGL", "Multi-Core: OpenMP + CUDA"]
    tx = Inches(0.7)
    ty = Inches(6.4)
    chip_total_w = 0
    chips = []
    for t in tags:
        cw = Inches(0.3 + 0.09 * len(t))
        chips.append((t, cw))
        chip_total_w += cw + Inches(0.15)
    start_x = (SLIDE_W - chip_total_w) / 2
    cx = start_x
    for t, cw in chips:
        add_round_rect(s, cx, ty, cw, Inches(0.35), CHIP_BG, radius=0.4)
        add_text(s, cx, ty + Inches(0.05), cw, Inches(0.28),
                 t, font=B_FONT, size=10, bold=True, color=CHIP_FG,
                 align=PP_ALIGN.CENTER)
        cx += cw + Inches(0.15)
    page_footer(s, "PROJECTS", 36)


def slide_proj_minigames():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s)
    add_section_marker(s, "04 / PROJECTS  ·  MOBILE GAMES")
    add_text(s, Inches(0.7), Inches(1.05), Inches(11), Inches(0.9),
             "Shipped mobile games", font=H_FONT, size=32, bold=True, color=INK)
    add_text(s, Inches(0.7), Inches(1.85), Inches(11), Inches(0.4),
             "Published Unity titles on the App Store and Google Play during the Pherma / freelance years.",
             font=B_FONT, size=13, italic=True, color=SUB)
    # 4-card grid
    games = [
        ("Ring and Ball",  "App Store — hyper-casual",  f"{MEDIA}/image7.png"),
        ("E' Gadung",      "Google Play — 2D platformer", f"{MEDIA}/image58.png"),
        ("LittleBounty",   "Hyper-casual platformer",        f"{MEDIA}/image47.png"),
        ("Stone Thrower",  "WebGL IoT multiplayer",          None),
    ]
    cw = Inches(2.95)
    sx = Inches(0.7)
    sy = Inches(2.6)
    ch = Inches(4.3)
    gap = Inches(0.1)
    # image area inside each card
    img_h = Inches(2.8)
    img_top = sy + Inches(0.2)
    for i, (name, desc, img) in enumerate(games):
        x = sx + i * (cw + gap)
        # card
        add_round_rect(s, x, sy, cw, ch, CARD, radius=0.05, line=cline(ACC1))
        # image container box (lighter slot so misaligned aspect doesn't look broken)
        slot_x = x + Inches(0.25)
        slot_w = cw - Inches(0.5)
        if img and os.path.exists(img):
            # constrain by height so tall portraits fit
            pic = add_image(s, img, slot_x, img_top, h=img_h)
            # if pic is wider than slot, re-add with width=slot_w instead
            if pic is not None and pic.width > slot_w:
                pic._element.getparent().remove(pic._element)
                add_image(s, img, slot_x, img_top, w=slot_w)
        else:
            # placeholder
            add_rect(s, slot_x, img_top, slot_w, img_h, CARD2)
            add_text(s, slot_x, img_top + Inches(1.15), slot_w, Inches(0.5),
                     "IoT  ·  WebGL", font=H_FONT, size=18, bold=True, color=ACC2,
                     align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.25), sy + Inches(3.25), cw - Inches(0.5), Inches(0.4),
                 name, font=H_FONT, size=15, bold=True, color=ACC1)
        add_text(s, x + Inches(0.25), sy + Inches(3.65), cw - Inches(0.5), Inches(0.6),
                 desc, font=B_FONT, size=10, color=CBODY,
                 line_spacing=1.2)
    light_footer(s, "PROJECTS", 37)


def slide_proj_gaugan():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s, feature=True)
    add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
             "04 / PROJECTS", font=H_FONT, size=11, bold=True, color=EYE)
    add_rect(s, Inches(0.7), Inches(0.85), Inches(0.4), Inches(0.04), ACC1)
    add_text(s, Inches(0.7), Inches(1.0), Inches(11), Inches(0.9),
             "GauGan Painter", font=H_FONT, size=32, bold=True, color=INK_F)
    add_text(s, Inches(0.7), Inches(1.85), Inches(11), Inches(0.4),
             "Companion DigiHuman backend: paint a semantic segmentation → GauGAN renders a photoreal background.",
             font=B_FONT, size=13, italic=True, color=SUB_F)
    add_image(s, f"{MEDIA}/image57.png", Inches(1.5), Inches(2.7), w=Inches(10))
    page_footer(s, "PROJECTS", 38)


def slide_proj_xv6():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s)
    add_section_marker(s, "04 / PROJECTS  ·  SYSTEMS")
    add_text(s, Inches(0.7), Inches(1.05), Inches(11), Inches(0.9),
             "Low-level systems work", font=H_FONT, size=32, bold=True, color=INK)
    add_text(s, Inches(0.7), Inches(1.85), Inches(11), Inches(0.4),
             "Kernel hacking, multi-core acceleration and cloud-scale microservices.",
             font=B_FONT, size=13, italic=True, color=SUB)
    items = [
        ("XV6 Customization", "Added custom system calls and modified CPU scheduling in the XV6 teaching kernel.", ACC1),
        ("Multi-Core Computing", "Optimized applications in C with [b]OpenMP[/b] and [b]CUDA[/b]; multi-threaded debugging on CPU+GPU.", ACC2),
        ("URL Shortener SaaS", "[b]Java + MySQL[/b] microservices on [b]Docker / Kubernetes / AWS[/b]; Hadoop for redundant distributed storage.", ACC1),
        ("Sitadu", "Restaurant database management app in [b]JavaFX + MySQL via JDBC[/b].", ACC2),
    ]
    y = Inches(2.7)
    for title, desc, color in items:
        add_round_rect(s, Inches(0.7), y, Inches(12), Inches(0.9), CARD, radius=0.05,
                       line=cline(color))
        add_rect(s, Inches(0.7), y, Inches(0.1), Inches(0.9), color)
        add_text(s, Inches(0.95), y + Inches(0.15), Inches(11), Inches(0.4),
                 title, font=H_FONT, size=15, bold=True, color=color)
        # render description with bold support
        tb = s.shapes.add_textbox(Inches(0.95), y + Inches(0.5), Inches(11.5), Inches(0.4))
        tf = tb.text_frame
        tf.margin_left = Pt(0); tf.margin_right = Pt(0)
        tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        for txt, is_bold in _split_bold(desc):
            run = p.add_run()
            run.text = txt
            run.font.name = B_FONT
            run.font.size = Pt(11)
            run.font.bold = is_bold
            run.font.color.rgb = CBODY
        y += Inches(1.0)
    light_footer(s, "PROJECTS", 39)


def slide_future_digital_earth():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s, feature=True)
    add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
             "04 / PROJECTS  ·  FUTURE WORK", font=H_FONT, size=11, bold=True, color=EYE)
    add_rect(s, Inches(0.7), Inches(0.85), Inches(0.4), Inches(0.04), ACC1)
    add_text(s, Inches(0.7), Inches(1.0), Inches(11), Inches(0.9),
             "Digital Earth", font=H_FONT, size=40, bold=True, color=INK_F)
    add_text(s, Inches(0.7), Inches(1.85), Inches(11), Inches(0.4),
             "Research direction: a multi-resolution interactive 3D Earth combining DGGS, subdivision surfaces and sketch-based modeling.",
             font=B_FONT, size=13, italic=True, color=SUB_F)
    add_image(s, f"{MEDIA}/image50.png", Inches(0.7), Inches(2.6), w=Inches(5.5))
    # right column
    pillars = [
        ("Discrete Global Grid System", "Hexagonal, hierarchical tessellation of the globe for uniform sampling and zoom."),
        ("Catmull–Clark Subdivision", "Smooth reverse-loop and Catmull–Clark surfaces for terrain detail."),
        ("Real-time Panorama Maps", "Stream high-resolution panoramas into the globe surface in real time."),
        ("Sketch-based 3D Content", "Inspired by NaturaSketch — draw to author 3D assets directly on the globe."),
    ]
    y = Inches(2.6)
    for title, desc in pillars:
        add_round_rect(s, Inches(6.7), y, Inches(6.1), Inches(1.0), CARD2, radius=0.08,
                       line=cline(ACC1))
        add_text(s, Inches(6.9), y + Inches(0.1), Inches(5.7), Inches(0.4),
                 title, font=H_FONT, size=13, bold=True, color=ACC1)
        add_text(s, Inches(6.9), y + Inches(0.45), Inches(5.7), Inches(0.5),
                 desc, font=B_FONT, size=10, color=CBODY,
                 line_spacing=1.2)
        y += Inches(1.1)
    page_footer(s, "PROJECTS", 40)


# ====================================================================
# HOBBIES
# ====================================================================
def slide_hobbies():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s)
    add_section_marker(s, "05 / HOBBIES")
    add_text(s, Inches(0.7), Inches(1.05), Inches(11), Inches(0.9),
             "Off the keyboard", font=H_FONT, size=40, bold=True, color=INK)
    add_text(s, Inches(0.7), Inches(1.95), Inches(11), Inches(0.4),
             "How I refuel — stories, strategy and outdoor games.",
             font=B_FONT, size=14, italic=True, color=SUB)
    hobbies = [
        ("Movies",     "Inception  ·  About Time  ·  Knives Out  ·  Memento", ACC1),
        ("Video Games","Rainbow Six Siege  ·  Apex Legends", ACC2),
        ("Animes",     "Your Name  ·  Howl's Moving Castle  ·  Steins;Gate", ACC1),
        ("Animations", "Wall-E  ·  Finding Nemo  ·  Shrek", ACC2),
        ("Board Games","Cluedo  ·  Mafia  ·  Code Names", ACC1),
        ("Sports",     "Chess  ·  Football  ·  Badminton", ACC2),
    ]
    cw = Inches(4.05)
    ch = Inches(1.4)
    sx = Inches(0.7)
    sy = Inches(2.85)
    gap = Inches(0.1)
    for i, (title, desc, color) in enumerate(hobbies):
        r = i // 3
        c = i % 3
        x = sx + c * (cw + gap)
        y = sy + r * (ch + gap)
        add_round_rect(s, x, y, cw, ch, CARD, radius=0.08, line=cline(color))
        # left accent
        add_rect(s, x, y, Inches(0.1), ch, color)
        add_text(s, x + Inches(0.3), y + Inches(0.2), cw - Inches(0.4), Inches(0.4),
                 title, font=H_FONT, size=16, bold=True, color=color)
        add_text(s, x + Inches(0.3), y + Inches(0.7), cw - Inches(0.4), Inches(0.6),
                 desc, font=B_FONT, size=11, color=CBODY,
                 line_spacing=1.25)
    light_footer(s, "HOBBIES", 41)


# ====================================================================
# CONTACT
# ====================================================================
def slide_contact():
    s = prs.slides.add_slide(blank_layout)
    page_bg(s, feature=True)
    add_text(s, Inches(0.7), Inches(0.7), Inches(8), Inches(0.4),
             "GET IN TOUCH", font=H_FONT, size=12, bold=True, color=EYE)
    add_rect(s, Inches(0.7), Inches(1.05), Inches(0.6), Inches(0.06), ACC1)
    add_text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(1.0),
             "Let's build something.", font=H_FONT, size=52, bold=True, color=INK_F)
    add_text(s, Inches(0.7), Inches(2.45), Inches(12), Inches(0.5),
             "Open to research collaborations, full-time roles and interesting hackathons.",
             font=B_FONT, size=16, italic=True, color=SUB_F)

    # contact cards
    contacts = [
        ("Email",    "danielkordm@gmail.com",                       ACC1),
        ("Phone",    "+1 (437) 559 3462",                           ACC1),
        ("LinkedIn", "linkedin.com/in/danial-kord",                 ACC2),
        ("GitHub",   "github.com/Danial-Kord",                      ACC2),
        ("Website",  "danial-kord.github.io",                       ACC1),
        ("Location", "Toronto, Ontario, Canada",                    ACC1),
    ]
    cw = Inches(6.1)
    ch = Inches(1.0)
    sx = Inches(0.7)
    sy = Inches(3.4)
    gap_x = Inches(0.15)
    gap_y = Inches(0.15)
    for i, (label, value, color) in enumerate(contacts):
        r = i // 2
        c = i % 2
        x = sx + c * (cw + gap_x)
        y = sy + r * (ch + gap_y)
        add_round_rect(s, x, y, cw, ch, CARD2, radius=0.08, line=cline(color))
        add_rect(s, x, y, Inches(0.1), ch, color)
        add_text(s, x + Inches(0.3), y + Inches(0.15), cw - Inches(0.4), Inches(0.35),
                 label.upper(), font=H_FONT, size=10, bold=True, color=color)
        add_text(s, x + Inches(0.3), y + Inches(0.5), cw - Inches(0.4), Inches(0.45),
                 value, font=H_FONT, size=15, bold=True, color=CHEAD)


# ====================================================================
# THANKS
# ====================================================================
def slide_thanks():
    s = prs.slides.add_slide(blank_layout)
    if DECO == "memphis":
        fill_slide(s, M_WHITE)
        add_oval(s, Inches(-1.3), Inches(-1.3), Inches(3.2), Inches(3.2), M_CYAN)
        add_oval(s, Inches(11.4), Inches(-1.0), Inches(2.8), Inches(2.8), M_CORAL)
        add_oval(s, Inches(12.3), Inches(6.6), Inches(2.0), Inches(2.0), M_YELLOW)
        add_oval(s, Inches(0.7), Inches(6.4), Inches(0.4), Inches(0.4), M_PURPLE)
        add_oval(s, Inches(2.0), Inches(1.0), Inches(0.3), Inches(0.3), M_PINK)
        add_text(s, Inches(1.0), Inches(1.6), Inches(11), Inches(0.5),
                 "THE END", font=H_FONT, size=14, bold=True, color=M_CORAL)
        add_text(s, Inches(1.0), Inches(2.2), Inches(11), Inches(2.0),
                 "Thank you.", font=H_FONT, size=96, bold=True, color=M_INK)
        add_rect(s, Inches(1.05), Inches(4.3), Inches(2.6), Inches(0.09), M_CYAN)
        add_text(s, Inches(1.0), Inches(4.6), Inches(11), Inches(0.5),
                 "Questions, ideas, or coffee? — I'd love to hear from you.",
                 font=B_FONT, size=18, italic=True, color=M_GREY_D)
        add_text(s, Inches(1.0), Inches(5.8), Inches(11.5), Inches(0.5),
                 "danielkordm@gmail.com   ·   linkedin.com/in/danial-kord   ·   github.com/Danial-Kord   ·   danial-kord.github.io",
                 font=B_FONT, size=14, color=M_INK)
        return

    fill_slide(s, NAVY)
    add_rect(s, Inches(0.6), Inches(1.4), Inches(0.12), Inches(4.5), AMBER)
    add_text(s, Inches(1.0), Inches(1.6), Inches(11), Inches(0.5),
             "THE END", font=H_FONT, size=14, bold=True, color=AMBER)
    add_text(s, Inches(1.0), Inches(2.2), Inches(11), Inches(2.0),
             "Thank you.", font=H_FONT, size=110, bold=True, color=IVORY)
    add_text(s, Inches(1.0), Inches(4.6), Inches(11), Inches(0.5),
             "Questions, ideas, or coffee? — I'd love to hear from you.",
             font=B_FONT, size=18, italic=True, color=GOLD)
    add_text(s, Inches(1.0), Inches(5.8), Inches(11.5), Inches(0.5),
             "danielkordm@gmail.com   ·   linkedin.com/in/danial-kord   ·   github.com/Danial-Kord   ·   danial-kord.github.io",
             font=B_FONT, size=14, color=IVORY)


# ====================================================================
# BUILD
# ====================================================================
def build(theme="midnight", role="general", out=None):
    apply_theme(theme)
    apply_role(role)
    _reset_deck()
    hide = rp("hide_proj", [])               # role may drop off-target projects

    slide_title()                            # 1
    slide_about()                            # 2
    slide_quote()                            # 3
    slide_toc()                              # 4
    section_divider("01", "Education", "Foundations — NODET  ·  AUT  ·  York", 5)
    slide_education()                        # 6
    slide_edu_york()                         # 7
    slide_edu_aut()                          # 8
    section_divider("02", "Experience", "Industry & research, 2017 — Today", 9)
    slide_exp_overview()                     # 10
    slide_exp_mla()                          # 11
    slide_exp_dreamforge()                   # 12
    slide_dreamforge_cicd()                  # DreamForge — self-healing CI/CD flowchart
    slide_exp_biomotion()                    # 13
    slide_exp_tectotrack()                   # 14
    slide_exp_techu()                        # 15
    slide_exp_iaeste()                       # 16
    slide_exp_sepantab()                     # 17
    slide_skills()                           # 18
    slide_languages()                        # 19
    slide_certifications()                   # 20
    slide_honors()                           # 21
    section_divider("04", "Projects",
                    rp("proj_div_sub", "Open-source, hackathons & academic builds"), 22)
    slide_proj_digihuman_cover()             # 23
    slide_proj_digihuman_arch()              # 24
    slide_proj_digihuman_demos()             # 25
    slide_proj_safezone_cover()              # 26
    slide_proj_safezone_arch()               # 27
    slide_proj_caselogic()                   # 28
    slide_proj_guardian_arch()               # Guardian — system architecture
    slide_proj_guardian_demo()               # Guardian — live demo (embedded video)
    slide_proj_latex_cv()                    # 29
    slide_proj_techu()                       # 30
    if "hypervigilance" not in hide: slide_proj_hypervigilance()  # 31
    slide_proj_backgammon()                  # 32
    slide_proj_cv()                          # 33
    slide_proj_ai_viz()                      # 34
    slide_proj_search()                      # 35
    slide_proj_solar()                       # 36
    if "minigames" not in hide: slide_proj_minigames()           # 37
    if "gaugan" not in hide: slide_proj_gaugan()                  # 38
    slide_proj_xv6()                         # 39
    slide_future_digital_earth()             # 40
    slide_hobbies()                          # 41
    slide_contact()                          # 42
    slide_thanks()                           # 43

    # --- Every page: a fade transition + a gentle auto fade-in entrance, so
    #     the whole CV feels animated. The DigiHuman pipeline already carries
    #     its grouped click reveals (it has a <p:timing>), so we leave it be.
    for slide in prs.slides:
        _dh_fade_transition(slide)
        if slide.element.find(qn('p:timing')) is None:
            _apply_auto_entrance(slide)

    # Renumber footers to true slide position (robust to inserted exhibits).
    for idx, slide in enumerate(prs.slides, start=1):
        for shp in slide.shapes:
            if getattr(shp, "name", "") != "FOOTNUM" or not shp.has_text_frame:
                continue
            label = shp.text_frame.text.split("/")[-1].strip()
            runs = shp.text_frame.paragraphs[0].runs
            if runs:
                runs[0].text = f"{idx:02d}  /  {label}"

    out = out or OUT_FILE.get(theme, "Daniel's CV (1).pptx")
    sfx = rp("suffix")
    if sfx and out.endswith(".pptx"):
        out = f"{out[:-5]} - {sfx}.pptx"
    prs.save(out)
    print(f"[{theme} | {ROLE}] Saved {out} with {len(prs.slides)} slides")


if __name__ == "__main__":
    theme_arg = (sys.argv[1] if len(sys.argv) > 1 else "midnight").lower()
    role_arg = (sys.argv[2] if len(sys.argv) > 2 else "general").lower()
    themes = (["midnight", "memphis"] if theme_arg == "both"
              else [theme_arg] if theme_arg in ("midnight", "memphis") else None)
    if themes is None or role_arg not in ROLE_PROFILES:
        print("usage: python build_cv.py [midnight|memphis|both] "
              f"[{'|'.join(ROLE_PROFILES)}]")
        sys.exit(1)
    for th in themes:
        build(th, role_arg)
