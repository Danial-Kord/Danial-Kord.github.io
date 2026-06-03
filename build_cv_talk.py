"""
Daniel Kordmodanlou — Talk / CV Deck
Focus: Agentic AI · LLMs · System Design

This deck is COMPREHENSIVE by default: it contains *all* of my education,
experience, projects and honors. Trim it for any given talk by HIDING what
you don't want — nothing is ever deleted, just toggled off.

────────────────────────────────────────────────────────────────────────
TWO VISUAL THEMES  (same content, different look)
────────────────────────────────────────────────────────────────────────
  python build_cv_talk.py            # midnight  (dark "Midnight Executive")
  python build_cv_talk.py midnight   #  → Daniel-CV-5min-Talk.pptx
  python build_cv_talk.py memphis    # bright playful "Memphis" (Slidesgo-style)
  python build_cv_talk.py both       #  → also Daniel-CV-5min-Talk-Memphis.pptx

Every slide is drawn through semantic theme tokens (apply_theme), so the exact
same builders render either style. "midnight" = deep navy + amber/gold, ivory
cards. "memphis" = white background, coral/cyan/yellow/purple, Lexend Deca +
Questrial fonts, big corner circles + dot clusters, outlined cards.

────────────────────────────────────────────────────────────────────────
HOW TO HIDE THINGS  (two levels of control)
────────────────────────────────────────────────────────────────────────
1. Hide an individual entry — set  "show": False  on any item inside the
   EDUCATION / EXPERIENCES / PROJECTS / HONORS lists in the CONTENT section.
2. Hide a whole slide / section — flip the matching key to False in SHOW_SLIDES.

The Title and Thank-you slides always render. Page numbers, the footer and
the entrance animations all adjust to whatever is left visible.
"""
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---- Raw palette · MIDNIGHT EXECUTIVE -------------------------------------
NAVY     = RGBColor(0x0F, 0x17, 0x2A)
INDIGO   = RGBColor(0x1E, 0x2A, 0x55)
ROYAL    = RGBColor(0x2E, 0x42, 0x82)
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)
GOLD     = RGBColor(0xFB, 0xBF, 0x24)
TEAL     = RGBColor(0x14, 0xB8, 0xA6)
IVORY    = RGBColor(0xF8, 0xFA, 0xFC)
SLATE    = RGBColor(0x94, 0xA3, 0xB8)
DARK_TEXT = RGBColor(0x1F, 0x29, 0x37)
LIGHT_GREY = RGBColor(0xE2, 0xE8, 0xF0)

# ---- Raw palette · MEMPHIS (bright / playful, from "Daniel's CV.pptx") -----
M_CORAL  = RGBColor(0xEC, 0x66, 0x3D)   # dominant accent
M_CYAN   = RGBColor(0x78, 0xDD, 0xF4)   # light cyan — underline / pills / dots
M_TEAL   = RGBColor(0x1C, 0x92, 0xBD)   # deeper cyan — readable accent text
M_YELLOW = RGBColor(0xF7, 0xB8, 0x45)   # decorative only
M_PURPLE = RGBColor(0x67, 0x58, 0x9D)
M_GREEN  = RGBColor(0x69, 0xD9, 0x71)   # decorative only
M_PINK   = RGBColor(0xEF, 0xA1, 0xBD)   # decorative only
M_INK    = RGBColor(0x1C, 0x1C, 0x1C)
M_GREY_D = RGBColor(0x4A, 0x4A, 0x4A)
M_GREY_M = RGBColor(0x80, 0x80, 0x80)
M_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
M_CARD   = RGBColor(0xF1, 0xF3, 0xF5)   # near-white card surface
M_LAV    = RGBColor(0xCF, 0xC7, 0xE6)   # muted text on purple panel

MEDIA = "extracted/ppt/media"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Decorative shape colours used by the Memphis corner/edge ornaments.
_DECO_PAL = [M_CORAL, M_CYAN, M_YELLOW, M_PURPLE, M_GREEN, M_PINK]

# Output filename per theme.
OUT_FILE = {
    "midnight": "Daniel-CV-5min-Talk.pptx",
    "memphis":  "Daniel-CV-5min-Talk-Memphis.pptx",
}


# ====================================================================
# THEME  ·  semantic tokens, rebound per theme by apply_theme()
# ====================================================================
# Builders never name a raw colour; they use these tokens. apply_theme()
# swaps them so one set of builders renders either visual style.
def apply_theme(name):
    midnight = dict(
        H_FONT="Calibri", B_FONT="Calibri Light",
        PAGE=IVORY, FEAT=NAVY, PANEL=NAVY, FRAME=NAVY,
        INK=NAVY, INK_F=IVORY, ON_DARK=IVORY,
        SUB=ROYAL, SUB_F=GOLD, BODY=DARK_TEXT, MUTE=SLATE, MUTE_D=SLATE,
        CARD=NAVY, CARD2=INDIGO, CHEAD=IVORY, CBODY=LIGHT_GREY,
        ACC1=AMBER, ACC2=GOLD, ACC3=TEAL, EYE=AMBER, PANEL_ACC=GOLD,
        CHIP_BG=INDIGO, CHIP_FG=GOLD,
        STACKHDR=ROYAL, STACKCHIP_BG=NAVY, STACKCHIP_FG=IVORY,
        PILL_BG=INDIGO, PILL_FG=GOLD, PILL2_BG=ROYAL, PILL2_FG=IVORY,
        FOOT=SLATE, UNDER=None,
        CARD_LINE="none", MARKER="bar", DECO="none",
        CATC={"ai": AMBER, "game": GOLD, "sys": TEAL},
    )
    memphis = dict(
        H_FONT="Lexend Deca", B_FONT="Questrial",
        PAGE=M_WHITE, FEAT=M_WHITE, PANEL=M_PURPLE, FRAME=M_CARD,
        INK=M_INK, INK_F=M_INK, ON_DARK=M_WHITE,
        SUB=M_GREY_D, SUB_F=M_GREY_D, BODY=M_INK, MUTE=M_GREY_M, MUTE_D=M_LAV,
        CARD=M_CARD, CARD2=M_CARD, CHEAD=M_INK, CBODY=M_GREY_D,
        ACC1=M_CORAL, ACC2=M_PURPLE, ACC3=M_TEAL, EYE=M_CORAL, PANEL_ACC=M_CYAN,
        CHIP_BG=M_PURPLE, CHIP_FG=M_WHITE,
        STACKHDR=M_CORAL, STACKCHIP_BG=M_PURPLE, STACKCHIP_FG=M_WHITE,
        PILL_BG=M_CYAN, PILL_FG=M_INK, PILL2_BG=M_PURPLE, PILL2_FG=M_WHITE,
        FOOT=M_GREY_M, UNDER=M_CYAN,
        CARD_LINE="accent", MARKER="dot", DECO="memphis",
        CATC={"ai": M_CORAL, "game": M_PURPLE, "sys": M_TEAL},
    )
    chosen = memphis if name == "memphis" else midnight
    globals().update(chosen)
    globals()["THEME_NAME"] = name


# Initialise to midnight so module-level defaults resolve before build().
apply_theme("midnight")


# ====================================================================
# CONTENT  ·  edit these lists / toggles to curate the deck
# ====================================================================

# --- Whole-slide toggles. Flip to False to drop an entire slide/section. --
SHOW_SLIDES = {
    "about":         True,
    "pillars":       True,
    "education":     True,   # comprehensive Education section
    "experience":    True,   # comprehensive Experience section
    "vector":        True,   # deep-dive: Vector Institute (current role)
    "latex_cv":      True,   # deep-dive: LaTeX CV Builder
    "dreamforge":    True,   # deep-dive: DreamForge
    "guardian":      True,   # deep-dive: Guardian (NVIDIA Spark Hack)
    "caselogic":     True,   # deep-dive: CaseLogic hackathon
    "safezone":      True,   # deep-dive: SAFEZone AI hackathon
    "safezone_arch": True,   # deep-dive: SAFEZone architecture diagram
    "digihuman":     True,   # deep-dive: animated DigiHuman pipeline
    "projects":      True,   # comprehensive Projects section
    "stack":         True,   # the toolkit
    "honors":        True,   # comprehensive Honors & Awards section
}

# ====================================================================
# ROLE MODES  ·  tailor the deck's wording to a target role
# ====================================================================
# "general" (default) uses the content written throughout this file. Pass a
# role on the CLI — e.g.  `python build_cv_talk.py memphis backend` — to
# re-frame the positioning, roles and projects toward that role. Add another
# dict to ROLE_PROFILES to support more roles (frontend, ml, ...).
ROLE = "general"

ROLE_PROFILES = {
    "general": {},
    "backend": {
        "suffix": "Backend",
        "footer_tag": "Backend Software Engineering  ·  APIs · Data · Cloud",
        "eyebrow": "BACKEND SOFTWARE ENGINEER  ·  APIs · DISTRIBUTED SYSTEMS  ·  2026",
        "title_role": "Backend software engineering  ·  APIs · data · cloud infrastructure",
        "about_lead": "I build [b]reliable, scalable backend systems[/b] — APIs, data pipelines and cloud infrastructure.",
        "about_bullets": [
            "Backend engineer across research labs and startups — [b]FastAPI, Spring Boot, Django, gRPC[/b] services over [b]PostgreSQL / MySQL / Redis[/b], containerized on [b]Docker / Kubernetes / AWS[/b].",
            "Shipped a [b]URL-shortener SaaS[/b] (Java microservices on K8s + AWS, Hadoop redundancy) and a [b]from-scratch search engine[/b] (inverted index, tf-idf, KNN).",
            "Built the [b]FastAPI + LangGraph[/b] backends for Guardian (multi-agent, idempotent 911 tooling) and CaseLogic (hybrid Chroma + SQLite FTS5 retrieval).",
            "Solid CS fundamentals: [b]XV6 kernel[/b] syscalls + scheduling, [b]multi-core C[/b] (OpenMP / CUDA), and [b]CI/CD[/b] on .NET / Azure / GitHub Actions.",
        ],
        "pillars_sub": "Three pillars I bring to every backend I build.",
        "pillars": [
            ("01", "APIs & Services",
             "REST / gRPC services, multi-agent orchestration and idempotent, fault-tolerant tooling.",
             ["FastAPI", "Spring Boot", "Django", "gRPC", "LangGraph"]),
            ("02", "Data & Storage",
             "Relational + vector stores, caching, search and retrieval — built to scale.",
             ["PostgreSQL", "MySQL", "Redis", "pgvector", "Chroma", "FTS5"]),
            ("03", "Infra & Delivery",
             "Containers, orchestration and CI/CD — built to ship and stay up.",
             ["Docker", "Kubernetes", "AWS", "GitHub Actions", ".NET / Azure"]),
        ],
        "stack": [
            ("Languages & APIs",
             ["Java", "Python", "C#", "C++", "Go", "FastAPI", "Spring Boot", "Django", "gRPC", "REST"]),
            ("Data & Messaging",
             ["PostgreSQL", "MySQL", "Redis", "pgvector", "Chroma", "SQLite FTS5", "MQTT", "Protobuf", "Hadoop"]),
            ("Infra & Delivery",
             ["Docker", "Kubernetes", "AWS (ECR)", "Azure", ".NET", "GitHub Actions", "Nginx", "Firebase", "CI/CD"]),
        ],
        "exp_hide": ["Pherma"],
        "exp_bullets": {
            "Vector Institute": [
                "Engineered the [b]FastAPI + Ollama / LangChain[/b] services and the [b]RAG[/b] retrieval pipeline (FAISS) behind a real-time VR coaching system.",
                "Built the data-collection & normalization pipeline mapping headset telemetry to NTU-25 skeletons for cross-platform robustness.",
            ],
            "BioMotion Lab, York University": [
                "Root-caused a Meta XR SDK stereo-geometry bug and shipped a reproducible [b]test harness[/b] documenting it.",
                "TA for Object-Oriented Programming (Java), 5 semesters — design patterns and systems thinking.",
            ],
            "DreamForge": [
                "Designed [b]CI/CD delivery[/b] components and [b]Docker[/b] workflows on [b].NET / Azure / GitHub Actions[/b] for automated deployment.",
                "Built Claude-assisted tooling that intercepts exceptions, isolates root causes across C# / Python / Docker layers, and commits fixes.",
            ],
            "TectoTrack": [
                "Built [b]simulation backends[/b] for high-traffic environments with fault-tolerant pathfinding.",
                "Shipped real-time [b]logging & monitoring dashboards[/b] for system health and agent-behaviour metrics.",
            ],
            "Techu": [
                "Built the online-multiplayer backend: [b]PlayFab (Azure)[/b] matchmaking and [b]Firebase[/b] real-time diagnostics & crash alerts.",
                "Engineered robust client/server game logic and AI agents (DQN + MCTS).",
            ],
            "IAESTE": [
                "Maintained and optimized a [b]Django[/b] backend (with a Vue.js frontend) for a student-exchange platform spanning 100+ countries.",
            ],
            "Amirkabir Univ. of Technology": [
                "TA for [b]Operating Systems, C Programming, Advanced Java[/b] and AI.",
                "Technical staff for AUT game-development events.",
            ],
            "Sepantab": [
                "Designed resilient distributed networking with [b]gRPC, Protocol Buffers, MQTT[/b] and WebSockets in C#.",
                "Built an IoT-based real-time multiplayer backend with low-latency state sync.",
            ],
        },
        "cat_order": ["sys", "ai", "game"],
        "projects_sub": "Backend & systems first — services, data, infra and low-level.",
        "project_desc": {
            "BakeryPilot": "FastAPI + LangGraph backend; PostgreSQL + pgvector, Redis streams, OR-Tools.",
            "Guardian": "FastAPI + LangGraph backend; deterministic risk monitor + idempotent Twilio tools.",
            "LaTeX CV Builder": "Spring Boot orchestration; Node / Python microservices; Docker + AWS; Nginx.",
            "URL Shortener SaaS": "Java microservices on Docker / K8s / AWS; Hadoop redundancy layer.",
            "Search Engine": "Inverted index · tf-idf · champion lists · KNN / K-means — from scratch.",
            "Stone Thrower": "IoT real-time backend — WebSockets / gRPC / MQTT networking (Unity client).",
        },
        "show_slides": {"safezone": False, "safezone_arch": False, "digihuman": False},
    },
}


def apply_role(name):
    global ROLE
    ROLE = name if name in ROLE_PROFILES else "general"


def rp(key, default=None):
    """Return the active role's override for `key`, or `default`."""
    return ROLE_PROFILES.get(ROLE, {}).get(key, default)


# --- EDUCATION.  "show": False to hide a degree. -------------------------
EDUCATION = [
    {"show": True,
     "degree": "M.Sc. Computer Science",
     "school": "York University",
     "loc":    "Toronto, Canada",
     "dates":  "Jan 2024 — Jun 2026",
     "gpa":    "GPA 4.0 / 4.0",
     "note":   "Fully funded graduate position · BioMotion Lab."},
    {"show": True,
     "degree": "B.Sc. Computer Science",
     "school": "Amirkabir University of Technology (AUT)",
     "loc":    "Tehran, Iran",
     "dates":  "2017 — 2022",
     "gpa":    "GPA 3.86 / 4.0",
     "note":   "Top-tier national polytechnic ('MIT of Iran')."},
    {"show": True,  # was commented-out in the LaTeX CV — hide if not needed
     "degree": "Mathematics & Physics Diploma",
     "school": "Allameh Helli High School (NODET)",
     "loc":    "Iran",
     "dates":  "2013 — 2017",
     "gpa":    "GPA 4.0 / 4.0",
     "note":   "National Org. for Development of Exceptional Talents."},
]

# --- EXPERIENCE.  "show": False to hide a role. --------------------------
EXPERIENCES = [
    {"show": True,
     "role": "Machine Learning Associate", "org": "Vector Institute",
     "loc": "Toronto, Canada", "dates": "Jan 2026 — May 2026",
     "bullets": [
        "VR firefighter training with DXTR — a deviation engine over 3D skeletal telemetry (MPJPE + quaternion distance).",
        "LLM coaching layer: RAG over training manuals (FAISS) via Ollama / LangChain for grounded feedback.",
     ]},
    {"show": True,
     "role": "Research & Teaching Assistant", "org": "BioMotion Lab, York University",
     "loc": "Toronto, Canada", "dates": "Jan 2024 — Apr 2026",
     "bullets": [
        "VR & motion-capture depth-perception studies; presented illusory-parallax research at ECVP 2025, Mainz.",
        "Diagnosed a Meta XR SDK stereo-geometry bug; TA for Object-Oriented Programming (Java), 5 semesters.",
     ]},
    {"show": True,
     "role": "Senior Unity Developer", "org": "DreamForge",
     "loc": "Miami, FL (Remote)", "dates": "Dec 2025 — Feb 2026",
     "bullets": [
        "Core contributor to a C# / Python AI-driven engine — tooling, procedural generation, CI/CD pipelines.",
        "Built Claude-assisted tools that intercept exceptions, investigate build issues, and commit fixes.",
     ]},
    {"show": True,
     "role": "Software Developer", "org": "TectoTrack",
     "loc": "Toronto, Canada", "dates": "Aug 2023 — Feb 2025",
     "bullets": [
        "Simulation frameworks for high-traffic environments with fault-tolerant, adaptive pathfinding.",
        "Real-time logging & monitoring dashboards for system health and agent-behaviour metrics.",
     ]},
    {"show": True,
     "role": "Lead Unity3D Game Developer", "org": "Techu",
     "loc": "Japan (Remote)", "dates": "Jan 2023 — Apr 2024",
     "bullets": [
        "Led a cross-platform multiplayer game; AI agents via DQN + Monte-Carlo Tree Search.",
        "PlayFab Azure matchmaking and Firebase monitoring for real-time diagnostics & crash alerts.",
     ]},
    {"show": True,
     "role": "Volunteer Full-Stack Developer", "org": "IAESTE",
     "loc": "Tehran, Iran", "dates": "Jan 2022 — Dec 2022",
     "bullets": [
        "Maintained and enhanced IAESTE's web platform (Vue.js + Django), improving performance and UX.",
     ]},
    {"show": True,
     "role": "Technical Staff & Teaching Assistant", "org": "Amirkabir Univ. of Technology",
     "loc": "Tehran, Iran", "dates": "Sep 2017 — Aug 2022",
     "bullets": [
        "Technical staff for AUT game-development events.",
        "TA for AI, Advanced Java, C Programming, and Operating Systems.",
     ]},
    {"show": True,
     "role": "Software Engineer Intern", "org": "Sepantab",
     "loc": "Tehran, Iran", "dates": "Jun 2021 — Aug 2021",
     "bullets": [
        "Built 'Stone Thrower', an IoT-based WebGL multiplayer game (Unity3D, WebSockets).",
        "Resilient networking via C#, gRPC, Protocol Buffers, and MQTT.",
     ]},
    {"show": True,  # was commented-out in the LaTeX CV — hide if not needed
     "role": "Game Developer", "org": "Pherma",
     "loc": "Iran", "dates": "Jan 2019 — Apr 2020",
     "bullets": [
        "Hyper-casual and adventure games for iOS, Android, and Windows using Unity and C#.",
     ]},
]

# --- PROJECTS.  "cat" sets the colour band (ai / game / sys). ------------
#     "show": False to hide a project.  Colours come from the active theme.
CAT_ORDER = ["ai", "game", "sys"]
CAT_LEGEND = [
    ("ai",   "AI · ML · VISION"),
    ("game", "GAMES · GRAPHICS · GAME-AI"),
    ("sys",  "SYSTEMS · IR · LOW-LEVEL"),
]
PROJECTS = [
    # AI / ML / Vision
    {"show": True, "cat": "ai", "name": "BakeryPilot",
     "desc": "Agentic supply-chain copilot; LangGraph + Next.js. 4th @ TMLS."},
    {"show": True, "cat": "ai", "name": "Guardian",
     "desc": "Local-first emergency AI; LangGraph agents on NVIDIA DGX Spark."},
    {"show": True, "cat": "ai", "name": "DigiHuman  ⭐ 500",
     "desc": "Open-source real-time 3D avatar animation from monocular video."},
    {"show": True, "cat": "ai", "name": "LaTeX CV Builder",
     "desc": "LangChain agents + RAG; Next.js streaming UI on Docker / AWS."},
    {"show": True, "cat": "ai", "name": "PISE (GAN)",
     "desc": "Reimplemented CVPR'21 decoupled-GAN person-image synthesis."},
    {"show": True, "cat": "ai", "name": "3D Reconstruction",
     "desc": "SIFT + RANSAC + PoinTr point-cloud upsampling from 2D images."},
    {"show": True, "cat": "ai", "name": "Panorama Stitching",
     "desc": "SIFT matching + RANSAC homography stitching (OpenCV)."},
    {"show": True, "cat": "ai", "name": "Harris Corner Detector",
     "desc": "Corner detection from scratch via image-derivative analysis."},
    {"show": True, "cat": "ai", "name": "AI Algorithms Viz",
     "desc": "Manim animations of graph-search algorithms for teaching."},
    {"show": True, "cat": "ai", "name": "GauGAN Painter",
     "desc": "Semantic-segmentation → photorealistic scene synthesis."},
    # Games / graphics / game-AI
    {"show": True, "cat": "game", "name": "Techu",
     "desc": "Cross-platform multiplayer card game — DQN + MCTS opponents."},
    {"show": True, "cat": "game", "name": "Backgammon 3D",
     "desc": "Turn-based multiplayer with Monte-Carlo AI + matchmaking."},
    {"show": True, "cat": "game", "name": "Solar System Sim",
     "desc": "C++ / OpenGL renderer with frame-rate optimization."},
    {"show": True, "cat": "game", "name": "HYPERVIGILANCE",
     "desc": "Generative-video + Unity psychological-horror short."},
    {"show": True, "cat": "game", "name": "RingBall (iOS)",
     "desc": "Commercial Unity hyper-casual game; custom mesh pathfinding."},
    {"show": True, "cat": "game", "name": "Mobile Games",
     "desc": "Unity titles shipped to the App Store / Google Play."},
    {"show": True, "cat": "game", "name": "E' Gadung",
     "desc": "2D platformer (Unity), published on Google Play."},
    {"show": True, "cat": "game", "name": "LittleBounty",
     "desc": "Multiplayer capture-the-flag (Java, UDP / TCP)."},
    {"show": True, "cat": "game", "name": "Stone Thrower",
     "desc": "IoT WebGL multiplayer cafe game (Unity, WebSockets)."},
    # Systems / IR / low-level
    {"show": True, "cat": "sys", "name": "Search Engine",
     "desc": "tf-idf · champion lists · KNN / K-means — from scratch."},
    {"show": True, "cat": "sys", "name": "URL Shortener SaaS",
     "desc": "Java · MySQL · Docker · K8s · AWS · Hadoop microservices."},
    {"show": True, "cat": "sys", "name": "Multi-Core Computing",
     "desc": "Optimized C with OpenMP + CUDA across CPU & GPU."},
    {"show": True, "cat": "sys", "name": "XV6 Kernel",
     "desc": "Custom syscalls + CPU-scheduling mods on the XV6 OS."},
    {"show": True, "cat": "sys", "name": "Sitadu",
     "desc": "Restaurant DB management — JavaFX + MySQL via JDBC."},
]

# --- HONORS & AWARDS.  "show": False to hide one. ------------------------
HONORS = [
    {"show": True, "big": "4th",        "title": "TMLS Toronto Hackathon",
     "desc": "BakeryPilot — agentic supply-chain copilot.", "year": "2026"},
    {"show": True, "big": "3rd",        "title": "Legal-Tech Hackathon",
     "desc": "CaseLogic — source-grounded legal research (24h).", "year": "2026"},
    {"show": True, "big": "Semi-Finals", "title": "Boson AI × MScAC",
     "desc": "SAFEZone AI — voice-first therapist agent (48h).", "year": "2025"},
    {"show": True, "big": "$10K",       "title": "L2M Lab to Market",
     "desc": "Accepted into the program with funding.", "year": "2025"},
    {"show": True, "big": "Funded",     "title": "M.Sc. @ York University",
     "desc": "Fully funded graduate position in CS.", "year": "2024"},
    {"show": True, "big": "500⭐",      "title": "DigiHuman (Open Source)",
     "desc": "Open-source 3D character-animation pipeline.", "year": "2023"},
    {"show": True, "big": "Top 0.5%",   "title": "AI Graduate Exam",
     "desc": "Iranian national graduate exam.", "year": "2023"},
    {"show": True, "big": "Top 1%",     "title": "Konkur Entrance Exam",
     "desc": "Iranian national university entrance exam.", "year": "2017"},
    {"show": True, "big": "4th",        "title": "National RoboCup",
     "desc": "Robotics competition — out of 32 teams.", "year": "2015"},
]


# ---- primitive shape / text helpers --------------------------------------
def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    spPr = shp.fill._xPr
    for el in spPr.findall(qn('a:effectLst')): spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))
    return shp


def add_round_rect(slide, x, y, w, h, fill, radius=0.08, line=None, line_w=1.25):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    try: shp.adjustments[0] = radius
    except Exception: pass
    spPr = shp.fill._xPr
    for el in spPr.findall(qn('a:effectLst')): spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))
    return shp


def add_oval(slide, x, y, w, h, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    spPr = shp.fill._xPr
    for el in spPr.findall(qn('a:effectLst')): spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))
    return shp


def add_chevron(slide, x, y, w, h, fill, left=False):
    """A chunky chevron arrow (points right; left=True flips it)."""
    shp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    if left:
        shp.rotation = 180
    spPr = shp.fill._xPr
    for el in spPr.findall(qn('a:effectLst')): spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))
    return shp


# --- DigiHuman pipeline assets (illustrations extracted from the template) --
DH_MEDIA = "digihuman_assets"
DH_MINT  = RGBColor(0xD4, 0xF0, 0xD5)   # node tiles / landmark bars
DH_CORAL = RGBColor(0xEE, 0x6B, 0x6B)   # process pills
DH_ARROW = RGBColor(0xFC, 0xA0, 0x8E)   # connector arrows
DH_INK   = RGBColor(0x2B, 0x2B, 0x2B)   # labels


def _split_bold(s):
    out, i = [], 0
    while i < len(s):
        start = s.find("[b]", i)
        if start == -1: out.append((s[i:], False)); break
        if start > i: out.append((s[i:start], False))
        end = s.find("[/b]", start)
        if end == -1: out.append((s[start:], False)); break
        out.append((s[start+3:end], True))
        i = end + 4
    return out


def add_text(slide, x, y, w, h, text, *, font=None, size=14, bold=False,
             color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             italic=False, line_spacing=None):
    if font is None: font = H_FONT
    if color is None: color = BODY
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing: p.line_spacing = line_spacing
    for txt, is_bold in _split_bold(text):
        run = p.add_run()
        run.text = txt
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold or is_bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, bullets, *, font=None, size=14,
                color=None, line_spacing=1.25, bullet_char="▸",
                spacing_after=6):
    if font is None: font = B_FONT
    if color is None: color = BODY
    tb = slide.shapes.add_textbox(x, y, w, h)
    try:
        tb.name = "BULLETS"          # tag so the animation reveals it per bullet
    except Exception:
        pass
    tf = tb.text_frame
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(spacing_after)
        if bullet_char:
            run = p.add_run()
            run.text = bullet_char + "  "
            run.font.name = font
            run.font.size = Pt(size)
            run.font.color.rgb = color
        for txt, is_bold in _split_bold(b):
            run = p.add_run()
            run.text = txt
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = is_bold
            run.font.color.rgb = color
    return tb


def fill_slide(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    spPr = bg.fill._xPr
    for el in spPr.findall(qn('a:effectLst')): spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))
    return bg


def add_image(slide, path, x, y, w=None, h=None):
    if not os.path.exists(path): return None
    if w and h: return slide.shapes.add_picture(path, x, y, w, h)
    if w: return slide.shapes.add_picture(path, x, y, width=w)
    if h: return slide.shapes.add_picture(path, x, y, height=h)
    return slide.shapes.add_picture(path, x, y)


# ---- theme-aware composite helpers ---------------------------------------
_deco_i = [0]


def page_bg(s, feature=False):
    """Fill the slide background for the active theme and (memphis) scatter
    playful corner ornaments that stay clear of the content area."""
    fill_slide(s, FEAT if feature else PAGE)
    if DECO == "memphis":
        idx = _deco_i[0]; _deco_i[0] += 1
        a = _DECO_PAL[idx % len(_DECO_PAL)]
        b = _DECO_PAL[(idx + 3) % len(_DECO_PAL)]
        # big quarter-circle hugging the (empty) top-right corner
        add_oval(s, Inches(11.85), Inches(-1.05), Inches(2.6), Inches(2.6), a)
        add_oval(s, Inches(11.35), Inches(1.78), Inches(0.26), Inches(0.26), b)
        add_oval(s, Inches(12.98), Inches(1.46), Inches(0.16), Inches(0.16), M_CYAN)


def cline(accent):
    """Outline colour for a card: the accent (memphis) or none (midnight)."""
    return accent if CARD_LINE == "accent" else None


def marker(s, label):
    """Section eyebrow marker: amber bar (midnight) or coral dot (memphis)."""
    if MARKER == "dot":
        add_oval(s, Inches(0.5), Inches(0.54), Inches(0.22), Inches(0.22), EYE)
        lx = Inches(0.82)
    else:
        add_rect(s, Inches(0.5), Inches(0.5), Inches(0.06), Inches(0.5), EYE)
        lx = Inches(0.7)
    add_text(s, lx, Inches(0.5), Inches(8), Inches(0.5),
             label, font=H_FONT, size=11, bold=True, color=EYE,
             anchor=MSO_ANCHOR.MIDDLE)


def section_header(s, mk, title, subtitle=""):
    """Standard section intro: themed bg + marker + big title + subtitle."""
    page_bg(s)
    marker(s, mk)
    add_text(s, Inches(0.7), Inches(1.05), Inches(12), Inches(0.9),
             title, font=H_FONT, size=40, bold=True, color=INK)
    if subtitle:
        add_text(s, Inches(0.7), Inches(1.97), Inches(12.3), Inches(0.5),
                 subtitle, font=B_FONT, size=14, italic=True, color=SUB)


def pill(s, x, y, text, *, fill=None, fg=None, size=10, h=0.32, bold=True):
    """A small rounded chip; returns its width (EMU) for right-alignment."""
    if fill is None: fill = PILL_BG
    if fg is None: fg = PILL_FG
    w = Inches(0.28 + 0.085 * len(text))
    add_round_rect(s, x, y, w, Inches(h), fill, radius=0.5)
    add_text(s, x, y, w, Inches(h), text, font=B_FONT, size=size, bold=bold,
             color=fg, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return w


def chip_row(s, x_start, y, items, color_fill=None, text_color=None,
             size=10, max_x=None):
    if color_fill is None: color_fill = CHIP_BG
    if text_color is None: text_color = CHIP_FG
    cx = x_start
    if max_x is None: max_x = Inches(13.0)
    for t in items:
        cw = Inches(0.3 + 0.09 * len(t))
        if cx + cw > max_x:
            cx = x_start
            y += Inches(0.45)
        add_round_rect(s, cx, y, cw, Inches(0.34), color_fill, radius=0.4)
        add_text(s, cx, y + Inches(0.045), cw, Inches(0.3),
                 t, font=B_FONT, size=size, bold=True, color=text_color,
                 align=PP_ALIGN.CENTER)
        cx += cw + Inches(0.1)
    return y


def add_footer(s, num, total):
    """Unified footer (name strip + page x/y), themed text colour."""
    add_text(s, Inches(0.5), Inches(7.05), Inches(8), Inches(0.4),
             f"Danial Kordmodanlou  ·  {rp('footer_tag', 'Agentic AI / LLMs / System Design')}",
             font=B_FONT, size=9, color=FOOT)
    add_text(s, Inches(11.6), Inches(7.05), Inches(1.3), Inches(0.4),
             f"{num:02d}  /  {total:02d}", font=B_FONT, size=9, color=FOOT,
             align=PP_ALIGN.RIGHT)


def shown(items):
    """Filter a content list down to the entries flagged show=True."""
    return [it for it in items if it.get("show", True)]


def chunk(items, n):
    """Yield successive n-sized pages from a list."""
    for i in range(0, len(items), n):
        yield items[i:i + n]


# Per-slide entrance-effect tuples, recorded as slides are created so the
# animation pass stays aligned even when slides are toggled off.
_slide_effects = []
prs = None
blank_layout = None


_no_footer_ids = set()


def _reset_deck():
    """Fresh Presentation + per-build state (lets us render both themes)."""
    global prs, blank_layout, _slide_effects, _SLIDE_BREAKS, _deco_i, _no_footer_ids
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]
    _slide_effects = []
    _SLIDE_BREAKS = {}
    _deco_i = [0]
    _no_footer_ids = set()


def new_slide(effects=("fade",)):
    s = prs.slides.add_slide(blank_layout)
    _slide_effects.append(effects)
    return s


# ====================================================================
# ANIMATIONS  ·  inject OOXML <p:transition> and <p:timing> trees
# ====================================================================
_PNS = "http://schemas.openxmlformats.org/presentationml/2006/main"

# Effect preset registry (PowerPoint built-in entrance effects).
# (presetID, presetClass, presetSubtype) — see ECMA-376 §19.5.51.
_EFFECTS = {
    "fade":         (10, "entr", 0),
    "float_up":     (42, "entr", 0),
    "zoom":         (23, "entr", 0),
    "wipe_up":      (22, "entr", 4),
    "fly_from_bot": (2,  "entr", 4),
}


def apply_fade_transition(slide, speed="med"):
    """Add a fade transition (inserted in correct schema slot: after cSld/clrMapOvr)."""
    sld = slide.element
    for el in sld.findall(qn('p:transition')):
        sld.remove(el)
    transition_el = etree.fromstring(
        f'<p:transition xmlns:p="{_PNS}" spd="{speed}"><p:fade/></p:transition>'
    )
    csld = sld.find(qn('p:cSld'))
    idx = list(sld).index(csld) + 1
    while idx < len(sld) and sld[idx].tag == qn('p:clrMapOvr'):
        idx += 1
    sld.insert(idx, transition_el)


def _shape_anim_ids(slide):
    """Shape IDs in z-order, skipping full-slide background fills."""
    out = []
    for sp in slide.shapes:
        try:
            l = int(sp.left or 0); t = int(sp.top or 0)
            w = int(sp.width or 0); h = int(sp.height or 0)
        except Exception:
            l = t = w = h = 0
        if l == 0 and t == 0 and w == int(SLIDE_W) and h == int(SLIDE_H):
            continue  # full-slide background — keep visible from slide load
        out.append(sp.shape_id)
    return out


# Manual section breaks: id(slide) -> list of shape_ids that END a section.
_SLIDE_BREAKS = {}


def mark_break(slide):
    """Mark the most recently added shape as the end of an animation section.

    On playback the cascade plays up to and including this shape, then waits
    for a click/space before revealing the next section. Use it in slide
    builders between logical blocks (header, each card, footer, ...).
    """
    if len(slide.shapes) == 0:
        return
    last_sid = slide.shapes[-1].shape_id
    _SLIDE_BREAKS.setdefault(id(slide), []).append(last_sid)


def _detect_group_boundaries(slide):
    """Auto-detect shape_ids that should START a new animation section.

    Heuristic: a rounded rectangle larger than ~2.5"x1.5" is treated as a
    'card' container and starts a fresh section. Smaller rounded rects
    (chips, buttons) are ignored.
    """
    THRESH_W = int(Inches(2.5))
    THRESH_H = int(Inches(1.5))
    boundary_ids = set()
    for sp in slide.shapes:
        try:
            if sp.auto_shape_type != MSO_SHAPE.ROUNDED_RECTANGLE:
                continue
            if int(sp.width or 0) >= THRESH_W and int(sp.height or 0) >= THRESH_H:
                boundary_ids.add(sp.shape_id)
        except Exception:
            continue
    return boundary_ids


def _split_into_groups(slide, animated_ids):
    """Partition animated shape_ids into ordered groups using manual breaks
    and auto-detected card boundaries."""
    manual_after = set(_SLIDE_BREAKS.get(id(slide), []))
    auto_before = _detect_group_boundaries(slide)
    groups, current = [], []
    for sid in animated_ids:
        if sid in auto_before and current:
            groups.append(current)
            current = []
        current.append(sid)
        if sid in manual_after:
            groups.append(current)
            current = []
    if current:
        groups.append(current)
    return groups


def _is_container(sp):
    """A 'box' that holds content: a rounded rectangle wide enough to be a
    card / tile / pill / bar (≥1.5in). Small rounded rects (chips) are NOT
    containers — they ride along with their parent card."""
    try:
        return (sp.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE
                and int(sp.width or 0) >= int(Inches(1.5)))
    except Exception:
        return False


def _center_in(sp, cont):
    """True if sp's center sits inside cont's bounds (sp is 'inside' the box)."""
    try:
        cx = int(sp.left or 0) + int(sp.width or 0) // 2
        cy = int(sp.top or 0) + int(sp.height or 0) // 2
        return (int(cont.left) <= cx <= int(cont.left) + int(cont.width)
                and int(cont.top) <= cy <= int(cont.top) + int(cont.height))
    except Exception:
        return False


def _content_groups(slide):
    """Partition a slide into click-reveal GROUPS of 'a box + everything in it'.

    Everything created before the first container box (background, decorative
    shapes, section marker, title, subtitle — the page 'style') is left OUT of
    every group, so it is STATIC: present the moment the slide opens. From the
    first box onward, each STANDALONE box starts a new group; shapes that
    follow it — including chips/pills whose centre sits INSIDE that box — ride
    along in the same group. So a box and its text reveal together on one
    click, while a separate box (e.g. a pipeline pill) gets its own click.
    """
    shapes = list(slide.shapes)
    first = next((i for i, sp in enumerate(shapes) if _is_container(sp)), None)
    if first is None:
        return []                       # no boxes → whole slide is static
    groups, cur, cont = [], [], None
    for sp in shapes[first:]:
        try:
            if (int(sp.left or 0) == 0 and int(sp.top or 0) == 0
                    and int(sp.width or 0) == int(SLIDE_W)
                    and int(sp.height or 0) == int(SLIDE_H)):
                continue                # skip a full-slide background
        except Exception:
            pass
        if _is_container(sp) and not (cont is not None and _center_in(sp, cont)):
            if cur:
                groups.append(cur)
            cur = [sp.shape_id]; cont = sp     # standalone box → new group
        else:
            cur.append(sp.shape_id)            # rides with the current box
    if cur:
        groups.append(cur)
    return groups


def _content_steps(slide):
    """Like _content_groups, but a bullet list (a text box tagged 'BULLETS')
    is split into ONE step per bullet (paragraph). Returns a list of steps;
    each step is a list of (shape_id, para) targets (para=None → whole shape,
    para=k → the k-th paragraph of that shape)."""
    shapes = list(slide.shapes)

    def is_bullet(sp):
        try:
            return sp.name == "BULLETS"
        except Exception:
            return False

    def npar(sp):
        try:
            return max(1, len(list(sp.text_frame.paragraphs)))
        except Exception:
            return 1

    first = next((i for i, sp in enumerate(shapes)
                  if _is_container(sp) or is_bullet(sp)), None)
    if first is None:
        return []
    steps, cur, cont = [], [], None

    def flush():
        if cur:
            steps.append(list(cur))
            cur.clear()

    for sp in shapes[first:]:
        try:
            if (int(sp.left or 0) == 0 and int(sp.top or 0) == 0
                    and int(sp.width or 0) == int(SLIDE_W)
                    and int(sp.height or 0) == int(SLIDE_H)):
                continue
        except Exception:
            pass
        if is_bullet(sp):
            flush()
            for k in range(npar(sp)):
                steps.append([(sp.shape_id, k)])   # one click per bullet
            cont = None
        elif _is_container(sp) and not (cont is not None and _center_in(sp, cont)):
            flush()
            cur.append((sp.shape_id, None)); cont = sp
        else:
            cur.append((sp.shape_id, None))
    flush()
    return steps


def _steps_timing_xml(steps, dur_ms):
    """One On-Click step per entry in `steps`. Each step fades its targets in
    together; paragraph targets get a build='p' bldP entry."""
    counter = [3]
    def nid():
        v = counter[0]; counter[0] += 1; return v
    click_steps = []
    para_spids, whole_spids = set(), set()
    for step in steps:
        click_id = nid(); inner_id = nid()
        anims = []
        for j, (spid, para) in enumerate(step):
            nt = "clickEffect" if j == 0 else "withEffect"
            anims.append(_anim_par_xml(spid, 10, "entr", 0, nt, 0, dur_ms, nid, para=para))
            (para_spids if para is not None else whole_spids).add(spid)
        click_steps.append(
            f'<p:par><p:cTn id="{click_id}" fill="hold">'
            f'<p:stCondLst><p:cond delay="indefinite"/></p:stCondLst><p:childTnLst>'
            f'<p:par><p:cTn id="{inner_id}" fill="hold">'
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst>{"".join(anims)}</p:childTnLst></p:cTn></p:par>'
            f'</p:childTnLst></p:cTn></p:par>')
    builds = ("".join(f'<p:bldP spid="{s}" grpId="0" build="p"/>' for s in para_spids)
              + "".join(f'<p:bldP spid="{s}" grpId="0"/>'
                        for s in (whole_spids - para_spids)))
    return (
        f'<p:timing xmlns:p="{_PNS}"><p:tnLst><p:par>'
        f'<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>'
        f'<p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>'
        f'{"".join(click_steps)}</p:childTnLst></p:cTn>'
        f'<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
        f'<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
        f'</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst>'
        f'<p:bldLst>{builds}</p:bldLst></p:timing>')


def apply_auto_entrance(slide, *, dur_ms=450):
    """Gentle fade-in of all content the moment the slide opens (no clicks).
    Used for slides with no content boxes so every page still animates in.
    Skips the full-slide background and the footer strip (kept static)."""
    ids = []
    for sp in slide.shapes:
        try:
            l = int(sp.left or 0); t = int(sp.top or 0)
            w = int(sp.width or 0); h = int(sp.height or 0)
        except Exception:
            l = t = w = h = 0
        if l == 0 and t == 0 and w == int(SLIDE_W) and h == int(SLIDE_H):
            continue                      # full-slide background
        if t >= int(Inches(7.0)):
            continue                      # footer strip
        ids.append(sp.shape_id)
    if not ids:
        return
    counter = [3]
    def nid():
        v = counter[0]; counter[0] += 1; return v
    click = nid(); inner = nid()
    anims, builds = [], []
    for i, spid in enumerate(ids):
        nt = "afterEffect" if i == 0 else "withEffect"   # auto-play on enter
        anims.append(_anim_par_xml(spid, 10, "entr", 0, nt, 0, dur_ms, nid))
        builds.append(f'<p:bldP spid="{spid}" grpId="0"/>')
    xml = (f'<p:timing xmlns:p="{_PNS}"><p:tnLst><p:par>'
           f'<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>'
           f'<p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>'
           f'<p:par><p:cTn id="{click}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>'
           f'<p:par><p:cTn id="{inner}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>'
           f'{"".join(anims)}</p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par>'
           f'</p:childTnLst></p:cTn>'
           f'<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
           f'<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
           f'</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst>'
           f'<p:bldLst>{"".join(builds)}</p:bldLst></p:timing>')
    sld = slide.element
    for el in sld.findall(qn('p:timing')):
        sld.remove(el)
    sld.append(etree.fromstring(xml))


def apply_cascade_anim(slide, *, effects=("fade",), step_ms=120, dur_ms=420):
    """Grouped click reveals: the page frame is static; each content box and
    its text fade in TOGETHER on one click, and a bullet list reveals one
    bullet per click. Slides with no boxes/bullets get a gentle auto fade-in
    instead, so every page still animates."""
    steps = _content_steps(slide)
    if not steps:
        apply_auto_entrance(slide, dur_ms=dur_ms)   # nothing to reveal → fade page in
        return
    sld = slide.element
    for el in sld.findall(qn('p:timing')):
        sld.remove(el)
    xml = _steps_timing_xml(steps, dur_ms)
    sld.append(etree.fromstring(xml))


def apply_single_cascade(slide, *, step_ms=0, dur_ms=420):
    """DigiHuman pipeline uses the same grouped click reveals."""
    apply_cascade_anim(slide, dur_ms=dur_ms)


def _anim_par_xml(spid, preset_id, preset_class, preset_sub,
                  node_type, offset_ms, dur_ms, nid, para=None):
    """Build the <p:par> for a single entrance animation on one shape, or on a
    single paragraph (para index) of a shape when `para` is given."""
    ctn_id = nid(); set_id = nid(); anim_id = nid()
    if para is None:
        tgt = f'<p:spTgt spid="{spid}"/>'
    else:
        tgt = (f'<p:spTgt spid="{spid}"><p:txEl>'
               f'<p:pRg st="{para}" end="{para}"/></p:txEl></p:spTgt>')
    return (
        f'<p:par><p:cTn id="{ctn_id}" presetID="{preset_id}" '
        f'presetClass="{preset_class}" presetSubtype="{preset_sub}" '
        f'fill="hold" grpId="0" nodeType="{node_type}">'
        f'<p:stCondLst><p:cond delay="{offset_ms}"/></p:stCondLst>'
        f'<p:childTnLst>'
        f'<p:set><p:cBhvr>'
        f'<p:cTn id="{set_id}" dur="1" fill="hold">'
        f'<p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
        f'<p:tgtEl>{tgt}</p:tgtEl>'
        f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
        f'</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>'
        f'<p:anim calcmode="lin" valueType="num">'
        f'<p:cBhvr additive="base">'
        f'<p:cTn id="{anim_id}" dur="{dur_ms}" fill="hold"/>'
        f'<p:tgtEl>{tgt}</p:tgtEl>'
        f'<p:attrNameLst><p:attrName>style.opacity</p:attrName></p:attrNameLst>'
        f'</p:cBhvr><p:tavLst>'
        f'<p:tav tm="0"><p:val><p:fltVal val="0"/></p:val></p:tav>'
        f'<p:tav tm="100000"><p:val><p:fltVal val="1"/></p:val></p:tav>'
        f'</p:tavLst></p:anim>'
        f'</p:childTnLst></p:cTn></p:par>'
    )


def _cascade_timing_xml(groups, effects, step_ms, dur_ms):
    """Build <p:timing> with one click step per section.

    Section 0: auto-starts on slide enter (cond delay='0').
    Sections 1+: each waits for a click event (cond delay='indefinite').
    """
    counter = [3]  # 1=tmRoot, 2=mainSeq; click step + inner par ids issued below
    def nid():
        v = counter[0]; counter[0] += 1; return v

    click_steps_xml = []
    builds = []
    eff_idx = 0
    for gi, group in enumerate(groups):
        click_id = nid()
        inner_id = nid()
        anims = []
        for i, spid in enumerate(group):
            eff_name = effects[eff_idx % len(effects)]
            preset_id, preset_class, preset_sub = _EFFECTS.get(eff_name, _EFFECTS["fade"])
            eff_idx += 1
            # First anim in each section is tied to the click trigger; the
            # rest run "with" it but offset for an overlapping cascade.
            node_type = "clickEffect" if i == 0 else "withEffect"
            offset = i * step_ms
            anims.append(_anim_par_xml(
                spid, preset_id, preset_class, preset_sub,
                node_type, offset, dur_ms, nid))
            builds.append(f'<p:bldP spid="{spid}" grpId="0"/>')

        step_delay = "indefinite"   # every step waits for its own click
        click_steps_xml.append(
            f'<p:par>'
            f'<p:cTn id="{click_id}" fill="hold">'
            f'<p:stCondLst><p:cond delay="{step_delay}"/></p:stCondLst>'
            f'<p:childTnLst><p:par>'
            f'<p:cTn id="{inner_id}" fill="hold">'
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst>{"".join(anims)}</p:childTnLst>'
            f'</p:cTn></p:par></p:childTnLst></p:cTn></p:par>'
        )

    return (
        f'<p:timing xmlns:p="{_PNS}">'
        f'<p:tnLst><p:par>'
        f'<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
        f'<p:childTnLst><p:seq concurrent="1" nextAc="seek">'
        f'<p:cTn id="2" dur="indefinite" nodeType="mainSeq">'
        f'<p:childTnLst>{"".join(click_steps_xml)}</p:childTnLst></p:cTn>'
        f'<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
        f'<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
        f'</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst>'
        f'<p:bldLst>{"".join(builds)}</p:bldLst>'
        f'</p:timing>'
    )


# ====================================================================
# 1 — TITLE
# ====================================================================
def slide_title():
    s = new_slide(("fade", "float_up"))
    if DECO == "memphis":
        fill_slide(s, M_WHITE)
        # playful geometric composition
        add_oval(s, Inches(-1.2), Inches(-1.2), Inches(3.0), Inches(3.0), M_CORAL)
        add_oval(s, Inches(1.75), Inches(0.35), Inches(0.42), Inches(0.42), M_CYAN)
        add_oval(s, Inches(11.3), Inches(0.5), Inches(1.5), Inches(1.5), M_YELLOW)
        add_oval(s, Inches(12.75), Inches(2.25), Inches(0.3), Inches(0.3), M_CYAN)
        add_oval(s, Inches(11.5), Inches(5.9), Inches(2.0), Inches(2.0), M_PURPLE)
        add_oval(s, Inches(10.8), Inches(6.7), Inches(0.3), Inches(0.3), M_CORAL)
        add_oval(s, Inches(0.6), Inches(6.75), Inches(0.34), Inches(0.34), M_PINK)
        add_text(s, Inches(1.0), Inches(2.05), Inches(11), Inches(0.4),
                 rp("eyebrow", "PORTFOLIO  ·  AGENTIC AI  ·  LLM SYSTEMS  ·  2026"),
                 font=H_FONT, size=12, bold=True, color=M_CORAL)
        add_text(s, Inches(1.0), Inches(2.5), Inches(12), Inches(1.5),
                 "Danial Kordmodanlou", font=H_FONT, size=60, bold=True, color=M_INK)
        add_rect(s, Inches(1.05), Inches(3.95), Inches(2.6), Inches(0.09), M_CYAN)
        add_text(s, Inches(1.05), Inches(4.2), Inches(11), Inches(0.6),
                 rp("title_role", "Building agentic AI systems  ·  grounded LLMs  ·  production at scale"),
                 font=B_FONT, size=19, color=M_GREY_D)
        add_text(s, Inches(1.05), Inches(5.1), Inches(11), Inches(0.4),
                 "Machine Learning Associate  ·  Vector Institute, Toronto",
                 font=B_FONT, size=14, color=M_INK)
        add_text(s, Inches(1.05), Inches(5.5), Inches(11), Inches(0.4),
                 "M.Sc. Computer Science  ·  York University  ·  Graduating June 2026",
                 font=B_FONT, size=13, italic=True, color=M_GREY_M)
        add_text(s, Inches(1.05), Inches(6.55), Inches(12), Inches(0.4),
                 "danielkordm@gmail.com   ·   linkedin.com/in/danial-kord   ·   github.com/Danial-Kord",
                 font=B_FONT, size=12, color=M_GREY_M)
        return

    fill_slide(s, NAVY)
    add_rect(s, Inches(0.6), Inches(0.9), Inches(0.12), Inches(5.7), AMBER)
    # decorative dots
    add_oval(s, Inches(11.6), Inches(0.7), Inches(0.5), Inches(0.5), AMBER)
    add_oval(s, Inches(11.85), Inches(0.95), Inches(0.5), Inches(0.5), INDIGO)
    # tag line
    add_text(s, Inches(1.0), Inches(0.95), Inches(11), Inches(0.4),
             rp("eyebrow", "PORTFOLIO  ·  AGENTIC AI  ·  LLM SYSTEMS  ·  2026"),
             font=H_FONT, size=12, bold=True, color=AMBER)
    # name
    add_text(s, Inches(1.0), Inches(1.65), Inches(12), Inches(1.6),
             "Danial Kordmodanlou",
             font=H_FONT, size=70, bold=True, color=IVORY)
    # role line
    add_text(s, Inches(1.0), Inches(3.4), Inches(11), Inches(0.7),
             rp("title_role", "Building agentic AI systems  ·  grounded LLMs  ·  production at scale"),
             font=B_FONT, size=22, color=GOLD)
    # divider
    add_rect(s, Inches(1.0), Inches(4.4), Inches(2.5), Inches(0.04), AMBER)
    # identity
    add_text(s, Inches(1.0), Inches(4.6), Inches(11), Inches(0.4),
             "Machine Learning Associate  ·  Vector Institute, Toronto",
             font=B_FONT, size=15, color=IVORY)
    add_text(s, Inches(1.0), Inches(5.05), Inches(11), Inches(0.4),
             "M.Sc. Computer Science  ·  York University",
             font=B_FONT, size=14, italic=True, color=SLATE)
    # contact strip
    add_text(s, Inches(1.0), Inches(6.2), Inches(12), Inches(0.4),
             "danielkordm@gmail.com   ·   linkedin.com/in/danial-kord   ·   github.com/Danial-Kord",
             font=B_FONT, size=12, color=SLATE)


# ====================================================================
# 2 — ABOUT  /  Why I'm here
# ====================================================================
def slide_about():
    s = new_slide(("fade",))
    page_bg(s)
    # left identity panel
    add_rect(s, 0, 0, Inches(5.0), SLIDE_H, PANEL)
    add_image(s, f"{MEDIA}/image39.jpg", Inches(0.9), Inches(1.4), w=Inches(3.2))
    add_rect(s, Inches(0.9), Inches(5.0), Inches(0.6), Inches(0.06), ACC1)
    add_text(s, Inches(0.9), Inches(5.1), Inches(4), Inches(0.4),
             "DANIAL KORDMODANLOU", font=H_FONT, size=13, bold=True, color=ON_DARK)
    add_text(s, Inches(0.9), Inches(5.45), Inches(4), Inches(0.4),
             "ML Associate — Vector Institute", font=B_FONT, size=11, color=PANEL_ACC)
    add_text(s, Inches(0.9), Inches(5.75), Inches(4), Inches(0.4),
             "M.Sc. Computer Science — York", font=B_FONT, size=11, color=MUTE_D)
    mark_break(s)  # — pause after left identity panel

    # right side
    rx = Inches(5.4)
    add_text(s, rx, Inches(0.7), Inches(7), Inches(0.5),
             "WHO I AM", font=H_FONT, size=12, bold=True, color=EYE)
    add_rect(s, rx, Inches(1.05), Inches(0.6), Inches(0.06), ACC1)

    add_text(s, rx, Inches(1.3), Inches(8), Inches(1.2),
             "Hi, I'm Danial.",
             font=H_FONT, size=44, bold=True, color=INK)
    add_text(s, rx, Inches(2.3), Inches(8), Inches(0.5),
             rp("about_lead", "I build [b]agentic AI systems[/b] grounded in real-world data."),
             font=B_FONT, size=17, italic=True, color=SUB)
    mark_break(s)  # — pause before bullets

    bullets = rp("about_bullets", [
        "Just completed a [b]Machine Learning Associate[/b] term at the [b]Vector Institute[/b], building a real-time VR firefighter training system: skeletal-telemetry deviation detection + an LLM coaching layer grounded in training manuals.",
        "Graduating with my [b]M.Sc. in Computer Science[/b] from [b]York University[/b] in June 2026 (GPA 4.0); presented VR depth-perception research at [b]ECVP 2025[/b], Mainz.",
        "Open-source: [b]DigiHuman[/b] (500⭐), a real-time 3D avatar animation pipeline.",
    ])
    add_bullets(s, rx, Inches(3.0), Inches(7.6), Inches(3.5),
                bullets, size=13, color=BODY, line_spacing=1.3,
                bullet_char="●")


# ====================================================================
# 3 — THE THREE PILLARS
# ====================================================================
def slide_pillars():
    s = new_slide(("zoom", "fade"))
    page_bg(s)
    marker(s, "FOCUS")
    add_text(s, Inches(0.7), Inches(1.05), Inches(11), Inches(0.9),
             "What I build", font=H_FONT, size=40, bold=True, color=INK)
    add_text(s, Inches(0.7), Inches(1.95), Inches(12), Inches(0.4),
             rp("pillars_sub", "Three pillars that cut through every project I take on."),
             font=B_FONT, size=15, italic=True, color=SUB)

    pillars = rp("pillars", [
        ("01",  "Agentic AI",
         "Multi-agent planning, tool-use orchestration, and per-claim verification.",
         ["LangChain", "LangGraph", "Multi-Agent", "n8n", "Cursor IDE"]),
        ("02",  "LLM Systems",
         "Grounded LLMs with RAG over hybrid retrieval, voice + vision modalities.",
         ["RAG", "FAISS", "Chroma", "Ollama", "GPT-4o", "Higgs Audio"]),
        ("03",  "System Design",
         "Microservices, CI/CD, containers and observability — built to ship.",
         ["Docker", "AWS", "Spring Boot", ".NET / Azure", "Streaming UI"]),
    ])
    _pcolors = [ACC1, ACC2, ACC1]
    sx = Inches(0.7); sy = Inches(2.7)
    cw = Inches(4.05); ch = Inches(3.95); gap = Inches(0.15)
    for i, (num, title, body, chips) in enumerate(pillars):
        color = _pcolors[i % 3]
        x = sx + i * (cw + gap)
        add_round_rect(s, x, sy, cw, ch, CARD, radius=0.05, line=cline(color))
        # left color band
        add_rect(s, x, sy, Inches(0.08), ch, color)
        add_text(s, x + Inches(0.35), sy + Inches(0.25), Inches(2), Inches(0.5),
                 num, font=H_FONT, size=22, bold=True, color=color)
        add_text(s, x + Inches(0.35), sy + Inches(0.85), cw - Inches(0.5), Inches(0.6),
                 title, font=H_FONT, size=22, bold=True, color=CHEAD)
        add_text(s, x + Inches(0.35), sy + Inches(1.55), cw - Inches(0.5), Inches(1.4),
                 body, font=B_FONT, size=12, color=CBODY,
                 line_spacing=1.3)
        # chips
        cy = sy + Inches(2.7)
        cx = x + Inches(0.35)
        max_x = x + cw - Inches(0.3)
        for c in chips:
            chip_w = Inches(0.28 + 0.085 * len(c))
            if cx + chip_w > max_x:
                cx = x + Inches(0.35); cy += Inches(0.4)
            add_round_rect(s, cx, cy, chip_w, Inches(0.32), CHIP_BG, radius=0.5)
            add_text(s, cx, cy + Inches(0.04), chip_w, Inches(0.28),
                     c, font=B_FONT, size=9, bold=True, color=CHIP_FG,
                     align=PP_ALIGN.CENTER)
            cx += chip_w + Inches(0.08)


# ====================================================================
# EDUCATION  (comprehensive · data-driven · paginated)
# ====================================================================
def build_education():
    entries = shown(EDUCATION)
    if not entries:
        return
    PER = 3
    pages = list(chunk(entries, PER))
    for pi, page in enumerate(pages):
        s = new_slide(("fade", "float_up"))
        title = "Education" + (" (cont.)" if pi else "")
        sub = "" if pi else "Where I trained — degrees, institutions and standing."
        section_header(s, "EDUCATION", title, sub)
        mark_break(s)

        x = Inches(0.7); w = Inches(11.93)
        ch = Inches(1.3); gap = Inches(0.18); sy = Inches(2.6)
        for i, e in enumerate(page):
            y = sy + i * (ch + gap)
            band = ACC1 if (pi * PER + i) % 2 == 0 else ACC2
            add_round_rect(s, x, y, w, ch, CARD, radius=0.06, line=cline(band))
            add_rect(s, x, y, Inches(0.1), ch, band)
            # left: degree / school / note
            add_text(s, x + Inches(0.4), y + Inches(0.16), Inches(8.0), Inches(0.45),
                     e["degree"], font=H_FONT, size=18, bold=True, color=CHEAD)
            add_text(s, x + Inches(0.4), y + Inches(0.62), Inches(8.2), Inches(0.35),
                     f'{e["school"]}  ·  {e["loc"]}', font=H_FONT, size=12.5,
                     bold=True, color=ACC2)
            add_text(s, x + Inches(0.4), y + Inches(0.94), Inches(8.2), Inches(0.32),
                     e["note"], font=B_FONT, size=10.5, italic=True, color=CBODY)
            # right: dates + GPA pills (right-aligned)
            right = x + w - Inches(0.3)
            pill(s, right - Inches(0.28 + 0.085 * len(e["dates"])),
                 y + Inches(0.2), e["dates"], fill=PILL_BG, fg=PILL_FG, size=10.5)
            if e.get("gpa"):
                gw = Inches(0.28 + 0.085 * len(e["gpa"]))
                pill(s, right - gw, y + Inches(0.66), e["gpa"],
                     fill=PILL2_BG, fg=PILL2_FG, size=10.5)
            mark_break(s)


# ====================================================================
# EXPERIENCE  (comprehensive · data-driven · paginated)
# ====================================================================
def build_experience():
    hide = rp("exp_hide", [])
    entries = [e for e in shown(EXPERIENCES) if e["org"] not in hide]
    if not entries:
        return
    PER = 3
    pages = list(chunk(entries, PER))
    for pi, page in enumerate(pages):
        s = new_slide(("fade", "wipe_up"))
        title = "Experience" + (" (cont.)" if pi else "")
        sub = "" if pi else rp("exp_sub",
                               "Roles where I shipped — research labs, startups and industry.")
        section_header(s, "EXPERIENCE", title, sub)
        mark_break(s)

        x = Inches(0.7); w = Inches(11.93)
        ch = Inches(1.4); gap = Inches(0.1); sy = Inches(2.5)
        for i, e in enumerate(page):
            y = sy + i * (ch + gap)
            band = ACC1 if (pi * PER + i) % 2 == 0 else ACC2
            add_round_rect(s, x, y, w, ch, CARD, radius=0.05, line=cline(band))
            add_rect(s, x, y, Inches(0.1), ch, band)
            # role + org line
            add_text(s, x + Inches(0.4), y + Inches(0.12), Inches(8.0), Inches(0.4),
                     e["role"], font=H_FONT, size=15, bold=True, color=CHEAD)
            add_text(s, x + Inches(0.4), y + Inches(0.5), Inches(8.0), Inches(0.32),
                     f'{e["org"]}  ·  {e["loc"]}', font=H_FONT, size=11,
                     bold=True, color=ACC2)
            # dates pill, right-aligned
            dlen = Inches(0.28 + 0.085 * len(e["dates"]))
            pill(s, x + w - dlen - Inches(0.3), y + Inches(0.16), e["dates"],
                 fill=PILL_BG, fg=PILL_FG, size=10)
            # condensed bullets (role-aware)
            ebullets = rp("exp_bullets", {}).get(e["org"], e["bullets"])
            add_bullets(s, x + Inches(0.4), y + Inches(0.8), w - Inches(0.8),
                        Inches(0.55), ebullets, size=9.5, color=CBODY,
                        line_spacing=1.1, bullet_char="▸", spacing_after=2)
            mark_break(s)


# ====================================================================
# DEEP-DIVE · VECTOR INSTITUTE  (current role headline)
# ====================================================================
def slide_vector():
    s = new_slide(("fade", "wipe_up"))
    page_bg(s)
    marker(s, "RECENT ROLE")
    add_text(s, Inches(0.7), Inches(1.0), Inches(8), Inches(0.4),
             "JANUARY 2026 — MAY 2026  ·  VECTOR INSTITUTE",
             font=H_FONT, size=12, bold=True, color=EYE)
    add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.9),
             "VR firefighter coaching — grounded LLMs in real time",
             font=H_FONT, size=32, bold=True, color=INK)
    add_text(s, Inches(0.7), Inches(2.25), Inches(12), Inches(0.5),
             "Partnered with [b]DXTR[/b] to detect skill deviations from headset telemetry and explain them in plain language.",
             font=B_FONT, size=14, italic=True, color=SUB)

    # two-column: Deviation Engine and Coaching Layer
    cw = Inches(5.95)
    sy = Inches(3.05)
    ch = Inches(3.7)

    # left card — Deviation Engine
    add_round_rect(s, Inches(0.7), sy, cw, ch, CARD, radius=0.06, line=cline(ACC1))
    add_rect(s, Inches(0.7), sy, Inches(0.1), ch, ACC1)
    add_text(s, Inches(0.95), sy + Inches(0.2), cw, Inches(0.4),
             "DEVIATION ENGINE", font=H_FONT, size=11, bold=True, color=ACC1)
    add_text(s, Inches(0.95), sy + Inches(0.55), cw, Inches(0.5),
             "Skeletal-telemetry analysis", font=H_FONT, size=20, bold=True, color=CHEAD)
    add_bullets(s, Inches(0.95), sy + Inches(1.2), cw - Inches(0.4), Inches(2.5),
                [
                    "Sliding windows of 3D skeletal telemetry",
                    "[b]MPJPE[/b] + quaternion angular distance vs. ground truth",
                    "Cross-platform [b]NTU-25 skeleton[/b] normalization",
                ],
                size=12, line_spacing=1.3, bullet_char="▸", color=CHEAD)

    # right card — Coaching Layer
    rx = Inches(6.85)
    add_round_rect(s, rx, sy, cw, ch, CARD, radius=0.06, line=cline(ACC2))
    add_rect(s, rx, sy, Inches(0.1), ch, ACC2)
    add_text(s, rx + Inches(0.25), sy + Inches(0.2), cw, Inches(0.4),
             "LLM COACHING LAYER", font=H_FONT, size=11, bold=True, color=ACC2)
    add_text(s, rx + Inches(0.25), sy + Inches(0.55), cw, Inches(0.5),
             "Grounded natural-language feedback", font=H_FONT, size=20, bold=True, color=CHEAD)
    add_bullets(s, rx + Inches(0.25), sy + Inches(1.2), cw - Inches(0.4), Inches(2.5),
                [
                    "[b]RAG[/b] over training-manual excerpts ([b]FAISS[/b] similarity)",
                    "[b]Ollama / LangChain[/b] for local & remote LLM orchestration",
                    "Structured evaluation of generative-AI training pipelines",
                ],
                size=12, line_spacing=1.3, bullet_char="▸", color=CHEAD)


# ====================================================================
# DEEP-DIVE · GUARDIAN  (NVIDIA Spark Hack — agentic emergency AI)
# ====================================================================
def slide_guardian():
    s = new_slide(("fade", "float_up"))
    page_bg(s)
    marker(s, "HACKATHON  ·  NVIDIA SPARK HACK")
    add_text(s, Inches(0.7), Inches(1.0), Inches(11), Inches(0.4),
             "NVIDIA SPARK HACK TORONTO  ·  2026",
             font=H_FONT, size=12, bold=True, color=EYE)
    add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(1.0),
             "Guardian — local-first emergency AI companion",
             font=H_FONT, size=32, bold=True, color=INK)
    add_text(s, Inches(0.7), Inches(2.3), Inches(12), Inches(0.5),
             "[b]Led the backend & multi-agent core[/b] (5-person team) — for people who live alone.",
             font=B_FONT, size=14, italic=True, color=SUB)

    # bullets left — how it works
    add_text(s, Inches(0.7), Inches(3.1), Inches(7), Inches(0.4),
             "HOW IT WORKS", font=H_FONT, size=11, bold=True, color=EYE)
    bullets = [
        "A [b]Wear OS watch[/b] detects falls & vitals anomalies; a deterministic fast-path routes emergencies to the Safety agent first.",
        "[b]FastAPI + LangGraph router[/b] over six specialist agents — only the Safety agent can dispatch 911.",
        "[b]CTAS-aligned risk scoring[/b] + idempotent Twilio tools in plain Python — the live demo never breaks on a flaky model.",
        "[b]Real-time voice[/b] (faster-whisper → Kokoro) + an SSE bus to a Next.js dashboard, Flutter app & Wear OS watch.",
    ]
    add_bullets(s, Inches(0.7), Inches(3.55), Inches(7.6), Inches(3.3),
                bullets, size=12, line_spacing=1.3, bullet_char="▸",
                color=BODY)

    # right column: hero stat + stack
    add_round_rect(s, Inches(8.6), Inches(3.1), Inches(4.2), Inches(3.7), CARD,
                   radius=0.06, line=cline(ACC1))
    add_rect(s, Inches(8.6), Inches(3.1), Inches(0.1), Inches(3.7), ACC1)
    add_text(s, Inches(8.85), Inches(3.25), Inches(3.8), Inches(0.5),
             "AUTONOMOUS", font=H_FONT, size=11, bold=True, color=ACC1)
    add_text(s, Inches(8.85), Inches(3.55), Inches(3.8), Inches(1.4),
             "911",
             font=H_FONT, size=88, bold=True, color=ACC1)
    add_text(s, Inches(8.85), Inches(4.95), Inches(3.7), Inches(0.5),
             "calls with full medical context — running on-device on an NVIDIA DGX Spark",
             font=B_FONT, size=11, color=CHEAD, line_spacing=1.2)
    add_text(s, Inches(8.85), Inches(5.55), Inches(3.8), Inches(0.4),
             "STACK", font=H_FONT, size=10, bold=True, color=ACC1)
    chip_row(s, Inches(8.85), Inches(5.95),
             ["FastAPI", "LangGraph", "Twilio", "Nemotron"],
             size=9, max_x=Inches(12.7))


# ====================================================================
# DEEP-DIVE · CASELOGIC  (3rd place hackathon)
# ====================================================================
def slide_caselogic():
    s = new_slide(("fade", "float_up"))
    page_bg(s)
    marker(s, "HACKATHON  ·  3RD PLACE")
    add_text(s, Inches(0.7), Inches(1.0), Inches(8), Inches(0.4),
             "TORONTO LEGAL-TECH HACKATHON  ·  2026",
             font=H_FONT, size=12, bold=True, color=EYE)
    add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(1.0),
             "CaseLogic — source-grounded legal research",
             font=H_FONT, size=32, bold=True, color=INK)
    add_text(s, Inches(0.7), Inches(2.3), Inches(12), Inches(0.5),
             "Built end-to-end in 24 hours. Every claim traced back to a paragraph.",
             font=B_FONT, size=14, italic=True, color=SUB)

    # bullets left
    add_text(s, Inches(0.7), Inches(3.1), Inches(7), Inches(0.4),
             "WHY IT WON", font=H_FONT, size=11, bold=True, color=EYE)
    bullets = [
        "[b]Hybrid retrieval[/b]: Chroma vector DB + SQLite FTS5 keyword search — vector recall plus literal-term precision for legal lookups.",
        "[b]Per-claim verification[/b]: every assertion grounded back to the originating source paragraph — no hallucinated citations.",
        "[b]Multi-agent planning workspace[/b]: orchestrated research, drafting, and verification agents in parallel.",
        "[b]Shipped end-to-end[/b] in 24h — UI + multi-agent backend, all deployed.",
    ]
    add_bullets(s, Inches(0.7), Inches(3.55), Inches(7.5), Inches(3.2),
                bullets, size=12, line_spacing=1.3, bullet_char="▸",
                color=BODY)

    # right column: badge + stack
    add_round_rect(s, Inches(8.6), Inches(3.1), Inches(4.2), Inches(3.7), CARD,
                   radius=0.06, line=cline(ACC1))
    add_rect(s, Inches(8.6), Inches(3.1), Inches(0.1), Inches(3.7), ACC1)

    # big "3rd"
    add_text(s, Inches(8.85), Inches(3.25), Inches(3.8), Inches(0.5),
             "RESULT", font=H_FONT, size=11, bold=True, color=ACC1)
    add_text(s, Inches(8.85), Inches(3.55), Inches(3.8), Inches(1.4),
             "3rd",
             font=H_FONT, size=88, bold=True, color=ACC1)
    add_text(s, Inches(8.85), Inches(5.0), Inches(3.8), Inches(0.4),
             "out of all Toronto Legal-Tech teams",
             font=B_FONT, size=11, color=CHEAD)

    add_text(s, Inches(8.85), Inches(5.55), Inches(3.8), Inches(0.4),
             "STACK", font=H_FONT, size=10, bold=True, color=ACC1)
    chip_row(s, Inches(8.85), Inches(5.95),
             ["Chroma", "SQLite FTS5", "LangChain", "Multi-Agent", "RAG"],
             size=9, max_x=Inches(12.7))


# ====================================================================
# DEEP-DIVE · SAFEZONE AI
# ====================================================================
def slide_safezone():
    s = new_slide(("fade", "fly_from_bot"))
    page_bg(s)
    marker(s, "HACKATHON  ·  SEMI-FINALIST")
    add_text(s, Inches(0.7), Inches(1.0), Inches(8), Inches(0.4),
             "BOSON AI × UofT MScAC HACKATHON  ·  2025",
             font=H_FONT, size=12, bold=True, color=EYE)
    add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(1.0),
             "SAFEZone AI — voice-first therapist agent",
             font=H_FONT, size=32, bold=True, color=INK)
    add_text(s, Inches(0.7), Inches(2.3), Inches(12), Inches(0.5),
             "A multi-stage agent pipeline that listens, reasons, and responds in a cloned voice with synced facial animation.",
             font=B_FONT, size=14, italic=True, color=SUB)

    # pipeline as steps
    steps = [
        ("ASR",        "Higgs Audio Understanding", "Transcript + emotional tone", ACC1),
        ("Reasoning",  "GPT-4o LLM",                "Onboarding-grounded reply", ACC2),
        ("TTS",        "Higgs Audio Generation",    "Cloned voice", ACC1),
        ("Avatar",     "Unity 3D therapist",        "Lip-sync + expressions", ACC2),
    ]
    sx = Inches(0.7); sy = Inches(3.15); cw = Inches(2.95); ch = Inches(3.4)
    gap = Inches(0.1)
    for i, (stage, name, desc, color) in enumerate(steps):
        x = sx + i * (cw + gap)
        add_round_rect(s, x, sy, cw, ch, CARD, radius=0.05, line=cline(color))
        add_rect(s, x, sy, Inches(0.08), ch, color)
        # step number
        add_text(s, x + Inches(0.3), sy + Inches(0.2), cw - Inches(0.4), Inches(0.4),
                 f"STEP {i+1}", font=H_FONT, size=10, bold=True, color=color)
        add_text(s, x + Inches(0.3), sy + Inches(0.55), cw - Inches(0.4), Inches(0.5),
                 stage, font=H_FONT, size=20, bold=True, color=CHEAD)
        add_rect(s, x + Inches(0.3), sy + Inches(1.1), Inches(0.4), Inches(0.04), color)
        add_text(s, x + Inches(0.3), sy + Inches(1.3), cw - Inches(0.4), Inches(0.6),
                 name, font=H_FONT, size=13, bold=True, color=ACC2)
        add_text(s, x + Inches(0.3), sy + Inches(1.85), cw - Inches(0.4), Inches(1.0),
                 desc, font=B_FONT, size=11, color=CHEAD, line_spacing=1.3)
        # arrow indicator between cards
        if i < len(steps) - 1:
            arrow_x = x + cw + Inches(-0.02)
            add_text(s, arrow_x, sy + Inches(1.5), Inches(0.15), Inches(0.4),
                     "▸", font=H_FONT, size=18, bold=True, color=ACC1,
                     align=PP_ALIGN.CENTER)


# ====================================================================
# DEEP-DIVE · SAFEZONE ARCHITECTURE (image)
# ====================================================================
def slide_safezone_arch():
    s = new_slide(("fade", "wipe_up"))
    page_bg(s, feature=True)
    add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
             "SAFEZONE AI  ·  ARCHITECTURE",
             font=H_FONT, size=11, bold=True, color=EYE)
    add_rect(s, Inches(0.7), Inches(0.85), Inches(0.4), Inches(0.04), ACC1)
    add_text(s, Inches(0.7), Inches(1.05), Inches(12), Inches(0.8),
             "End-to-end agent loop",
             font=H_FONT, size=30, bold=True, color=INK_F)
    add_text(s, Inches(0.7), Inches(1.85), Inches(12), Inches(0.4),
             "Voice  →  ASR + emotion  →  LLM grounded in APA PsycInfo  →  cloned-voice TTS  →  blend-shape lip-sync.",
             font=B_FONT, size=13, italic=True, color=SUB_F)
    mark_break(s)  # — pause before architecture diagram

    if DECO == "memphis":
        # The diagram art is dark-on-navy; on a white slide, seat it in an
        # intentional dark "exhibit" frame so it reads as deliberate.
        img_w = 10.6; ix = (13.333 - img_w) / 2.0; iy = 2.45
        img_h = img_w * 756.0 / 1847.0
        pad = 0.16
        add_round_rect(s, Inches(ix - pad), Inches(iy - pad),
                       Inches(img_w + 2 * pad), Inches(img_h + 2 * pad),
                       NAVY, radius=0.03)
        add_image(s, f"{MEDIA}/image26.png", Inches(ix), Inches(iy), w=Inches(img_w))
    else:
        add_image(s, f"{MEDIA}/image26.png", Inches(0.7), Inches(2.45), w=Inches(12))


# ====================================================================
# DEEP-DIVE · LATEX CV BUILDER  (System design + RAG)
# ====================================================================
def slide_latex_cv():
    """LaTeX CV Builder — animated pipeline exhibit (DigiHuman-style).
    One job description in -> a tailored CV, cover letter & fit score out.
    Mirrors the real repo: Obsidian vault + JD -> 5-stage cv_pipeline -> PDFs."""
    s = new_slide(("seq",))
    _no_footer_ids.add(id(s))            # full-bleed diagram: no footer
    fill_slide(s, M_WHITE)

    # --- decorative blobs (static page frame) ---
    add_oval(s, Inches(11.98), Inches(-0.80), Inches(2.05), Inches(2.05), M_PURPLE)
    add_oval(s, Inches(12.63), Inches(2.92), Inches(0.74), Inches(0.74), M_CYAN)
    add_oval(s, Inches(-0.55), Inches(6.40), Inches(1.45), Inches(1.45), M_YELLOW)
    add_oval(s, Inches(0.45), Inches(0.34), Inches(0.24), Inches(0.24), DH_CORAL)

    # --- title + subtitle (static) ---
    add_text(s, Inches(0.5), Inches(0.10), Inches(12.33), Inches(0.54),
             "LaTeX CV Builder", font=H_FONT, size=30, bold=True,
             color=DH_INK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(0.64), Inches(12.33), Inches(0.34),
             "One job description in   →   a tailored CV, cover letter & fit score out",
             font=B_FONT, size=12.5, italic=True, color=M_TEAL,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ---- nested helpers ----
    def card(x, y, w, h, fill, *, radius=0.10, line=None):
        return add_round_rect(s, Inches(x), Inches(y), Inches(w), Inches(h),
                              fill, radius=radius, line=line)

    def t(x, y, w, h, text, *, size=10, bold=False, color=DH_INK,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=None, italic=False,
          ls=None):
        add_text(s, Inches(x), Inches(y), Inches(w), Inches(h), text,
                 font=font or H_FONT, size=size, bold=bold, color=color,
                 align=align, anchor=anchor, italic=italic, line_spacing=ls)

    def badge(x, y, w, h, text, fill, size=8.5):
        add_round_rect(s, Inches(x), Inches(y), Inches(w), Inches(h), fill, radius=0.5)
        add_text(s, Inches(x), Inches(y), Inches(w), Inches(h), text, font=H_FONT,
                 size=size, bold=True, color=M_WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

    def arrow(x, y, w=0.52, h=0.56, left=False, down=False):
        c = add_chevron(s, Inches(x), Inches(y), Inches(w), Inches(h),
                        DH_ARROW, left=left)
        if down:
            c.rotation = 90
        return c

    TOP_Y, BOT_Y, TH = 1.62, 3.92, 1.88

    # ---- SOURCES lane (read every run) ----
    card(0.40, 1.06, 12.13, 0.46, M_CARD, radius=0.22, line=M_CYAN)
    t(0.62, 1.06, 2.0, 0.46, "READS EVERY RUN", size=8.5, bold=True,
      color=M_PURPLE, anchor=MSO_ANCHOR.MIDDLE)
    t(2.78, 1.06, 4.6, 0.46, "Obsidian vault  ·  experiences / projects / skills",
      size=9.5, bold=True, color=M_TEAL, anchor=MSO_ANCHOR.MIDDLE)
    t(7.55, 1.06, 2.3, 0.46, "canonical CV .tex", size=9.5, bold=True,
      color=DH_INK, anchor=MSO_ANCHOR.MIDDLE)
    t(10.0, 1.06, 2.4, 0.46, "referrals.db (SQLite)", size=9.5, bold=True,
      color=DH_INK, anchor=MSO_ANCHOR.MIDDLE)

    # ---- INPUT: Job Description (the trigger) ----
    card(0.40, TOP_Y, 1.85, TH, M_PURPLE)
    t(0.40, TOP_Y + 0.18, 1.85, 0.28, "INPUT", size=8.5, bold=True,
      color=M_WHITE, align=PP_ALIGN.CENTER)
    t(0.45, TOP_Y + 0.55, 1.75, 0.72, "Job Description", size=15, bold=True,
      color=M_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, ls=1.0)
    t(0.45, TOP_Y + 1.30, 1.75, 0.48, "paste · URL · Gmail queue", size=8.5,
      color=M_WHITE, align=PP_ALIGN.CENTER, ls=1.0)
    arrow(2.27, TOP_Y + 0.66)

    # ---- the 5-stage pipeline (LLM only in stages 1 & 3) ----
    def stage(x, y, num, title, module, desc, out, badge_txt, badge_fill,
              out_color=DH_CORAL):
        w = 2.85
        card(x, y, w, TH, DH_MINT)
        add_oval(s, Inches(x + 0.15), Inches(y + 0.16), Inches(0.44),
                 Inches(0.44), DH_CORAL)
        t(x + 0.15, y + 0.16, 0.44, 0.44, num, size=15, bold=True, color=M_WHITE,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        t(x + 0.68, y + 0.15, w - 0.84, 0.46, title, size=14, bold=True,
          color=DH_INK, anchor=MSO_ANCHOR.MIDDLE)
        t(x + 0.18, y + 0.72, w - 0.36, 0.28, module, size=9.5, bold=True,
          color=M_TEAL, font=B_FONT)
        t(x + 0.18, y + 1.01, w - 0.36, 0.52, desc, size=9, color=DH_INK, ls=1.04)
        t(x + 0.18, y + 1.54, w - 1.15, 0.30, out, size=10, bold=True,
          color=out_color, anchor=MSO_ANCHOR.MIDDLE)
        bw = 0.34 + 0.082 * len(badge_txt)
        badge(x + w - bw - 0.14, y + 1.55, bw, 0.30, badge_txt, badge_fill)

    # top row: 1 Parse JD -> 2 Select -> 3 Generate
    stage(2.88, TOP_Y, "1", "Parse JD", "job_parser.py",
          "LLM reads the JD into domain, role-fit, tags & constraints",
          "→ JobTargets", "LLM", DH_CORAL)
    arrow(5.78, TOP_Y + 0.66)
    stage(6.37, TOP_Y, "2", "Select", "candidate_selector.py",
          "Rule-based scoring of vault notes — anchors & de-dupe",
          "→ shortlist", "no LLM", M_TEAL)
    arrow(9.27, TOP_Y + 0.66)
    stage(9.86, TOP_Y, "3", "Generate", "bullet · cover · fit",
          "One LLM call writes bullets, a cover letter & a fit score",
          "→ bullets + letter", "LLM", DH_CORAL)
    # serpentine turn: stage 3 (top) -> stage 4 (bottom)
    arrow(11.05, 3.46, w=0.46, h=0.52, down=True)

    # bottom row (right -> left): 4 Render LaTeX -> 5 Compile
    stage(9.86, BOT_Y, "4", "Render LaTeX", "latex_formatter.py",
          "Fills résumé + per-entry .tex from the generation",
          "→ resume.tex", ".tex", M_PURPLE)
    arrow(9.27, BOT_Y + 0.66, left=True)
    stage(6.37, BOT_Y, "5", "Compile", "build_runner.py",
          "pdflatex builds the final PDF (skippable)",
          "→ CV.pdf", "pdflatex", M_TEAL)
    arrow(5.78, BOT_Y + 0.66, left=True)

    # ---- OUTPUTS panel ----
    card(0.40, BOT_Y, 5.30, TH, M_CARD, radius=0.10, line=M_GREEN)
    t(0.62, BOT_Y + 0.12, 3.0, 0.30, "OUTPUTS", size=10.5, bold=True,
      color=M_PURPLE)
    outs = [("CV.pdf", "tailored résumé", DH_CORAL),
            ("Cover letter .pdf", "matched tone & structure", M_TEAL),
            ("Fit score", "0–100 with strengths & gaps", M_GREEN),
            ("Archived", "JSON · SQLite · Obsidian note", M_PURPLE)]
    for i, (name, sub, acc) in enumerate(outs):
        ry = BOT_Y + 0.50 + i * 0.335
        add_oval(s, Inches(0.66), Inches(ry + 0.05), Inches(0.15),
                 Inches(0.15), acc)
        t(0.92, ry, 1.95, 0.30, name, size=10.5, bold=True, color=DH_INK,
          anchor=MSO_ANCHOR.MIDDLE)
        t(2.90, ry, 2.70, 0.30, sub, size=8.7, color=DH_INK,
          anchor=MSO_ANCHOR.MIDDLE)

    # ---- backend swap + UI (captions) ----
    card(1.95, 6.04, 9.45, 0.48, M_PURPLE, radius=0.24)
    t(1.95, 6.04, 9.45, 0.48,
      "LLM backend (swappable):   api_llm — LangChain → GPT-4o / Ollama"
      "      ·      claude_code — claude CLI",
      size=9.5, bold=True, color=M_WHITE, align=PP_ALIGN.CENTER,
      anchor=MSO_ANCHOR.MIDDLE)
    t(1.95, 6.60, 9.45, 0.32,
      "CV-UI · Next.js wraps the CLI — paste a JD and watch Step 1/5 … 5/5 stream live",
      size=10, bold=True, color=M_TEAL, align=PP_ALIGN.CENTER,
      anchor=MSO_ANCHOR.MIDDLE)


# ====================================================================
# DEEP-DIVE · DREAMFORGE  (AI tooling)
# ====================================================================
def slide_dreamforge():
    s = new_slide(("fade", "wipe_up"))
    page_bg(s)
    marker(s, "INDUSTRY  ·  AGENTIC TOOLING")
    add_text(s, Inches(0.7), Inches(1.0), Inches(8), Inches(0.4),
             "DECEMBER 2025 — FEBRUARY 2026  ·  DREAMFORGE",
             font=H_FONT, size=12, bold=True, color=EYE)
    add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(1.0),
             "Claude-assisted dev tools that fix themselves",
             font=H_FONT, size=30, bold=True, color=INK)
    add_text(s, Inches(0.7), Inches(2.3), Inches(12), Inches(0.5),
             "AI-driven engine for procedural generation — with autonomous build-health agents.",
             font=B_FONT, size=14, italic=True, color=SUB)

    # 3 callouts
    cards = [
        ("Auto-intercept",
         "Custom Claude tooling that watches CI/CD streams and intercepts thrown exceptions in real time.",
         ACC1),
        ("Auto-investigate",
         "Agents reproduce, isolate root causes across C# / Python / Docker layers, and gather evidence.",
         ACC2),
        ("Auto-resolve",
         "When confidence is high, the agent commits a permanent codebase fix straight into the repo.",
         ACC1),
    ]
    sx = Inches(0.7); sy = Inches(3.1); cw = Inches(4.15); ch = Inches(3.7); gap = Inches(0.1)
    for i, (title, desc, color) in enumerate(cards):
        x = sx + i * (cw + gap)
        add_round_rect(s, x, sy, cw, ch, CARD, radius=0.06, line=cline(color))
        add_rect(s, x, sy, Inches(0.1), ch, color)
        add_text(s, x + Inches(0.35), sy + Inches(0.3), cw, Inches(0.5),
                 f"0{i+1}", font=H_FONT, size=28, bold=True, color=color)
        add_text(s, x + Inches(0.35), sy + Inches(0.95), cw - Inches(0.5), Inches(0.5),
                 title, font=H_FONT, size=22, bold=True, color=CHEAD)
        add_rect(s, x + Inches(0.35), sy + Inches(1.55), Inches(0.4), Inches(0.04), color)
        add_text(s, x + Inches(0.35), sy + Inches(1.75), cw - Inches(0.6), Inches(1.7),
                 desc, font=B_FONT, size=12, color=CBODY, line_spacing=1.35)


# ====================================================================
# PROJECTS  (comprehensive · data-driven · paginated grid)
# ====================================================================
def build_projects():
    entries = shown(PROJECTS)
    if not entries:
        return
    # group by category order so colour bands cluster (role may reorder)
    order = rp("cat_order", CAT_ORDER)
    entries = sorted(entries, key=lambda p: order.index(p.get("cat", "ai"))
                     if p.get("cat", "ai") in order else len(order))
    PER = 12  # 4 cols x 3 rows
    pages = list(chunk(entries, PER))
    for pi, page in enumerate(pages):
        s = new_slide(("zoom", "fade"))
        title = "Projects" + (" (cont.)" if pi else "")
        sub = (rp("projects_sub",
                  "A snapshot of what I've shipped — across AI, vision, games and systems.")
               if pi == 0 else "")
        section_header(s, "PROJECTS", title, sub)

        # color-coded legend (every page)
        lx = Inches(0.7); ly = Inches(2.42)
        for key, label in CAT_LEGEND:
            add_oval(s, lx, ly + Inches(0.05), Inches(0.16), Inches(0.16), CATC[key])
            add_text(s, lx + Inches(0.25), ly, Inches(3.0), Inches(0.3),
                     label, font=H_FONT, size=9, bold=True, color=BODY)
            lx += Inches(3.1)
        mark_break(s)

        sx = Inches(0.7); sy = Inches(2.9)
        cw = Inches(3.0); ch = Inches(1.28)
        gap_x = Inches(0.13); gap_y = Inches(0.12)
        for i, p in enumerate(page):
            r = i // 4
            c = i % 4
            x = sx + c * (cw + gap_x)
            y = sy + r * (ch + gap_y)
            color = CATC.get(p.get("cat", "ai"), ACC1)
            add_round_rect(s, x, y, cw, ch, CARD, radius=0.06, line=cline(color))
            add_rect(s, x, y, Inches(0.08), ch, color)
            pdesc = rp("project_desc", {}).get(p["name"], p["desc"])
            add_text(s, x + Inches(0.22), y + Inches(0.16), cw - Inches(0.35), Inches(0.4),
                     p["name"], font=H_FONT, size=12, bold=True, color=color)
            add_text(s, x + Inches(0.22), y + Inches(0.56), cw - Inches(0.35), Inches(0.66),
                     pdesc, font=B_FONT, size=9.5, color=CBODY, line_spacing=1.2)
            # pause after each full row of 4 (except the last row on the page)
            if (i + 1) % 4 == 0 and (i + 1) < len(page):
                mark_break(s)


# ====================================================================
# STACK / TOOLKIT
# ====================================================================
def slide_stack():
    s = new_slide(("fade", "float_up"))
    page_bg(s, feature=True)
    add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
             "THE TOOLKIT",
             font=H_FONT, size=11, bold=True, color=EYE)
    add_rect(s, Inches(0.7), Inches(0.85), Inches(0.4), Inches(0.04), ACC1)
    add_text(s, Inches(0.7), Inches(1.05), Inches(12), Inches(0.9),
             "Stack I reach for", font=H_FONT, size=38, bold=True, color=INK_F)
    add_text(s, Inches(0.7), Inches(1.9), Inches(12), Inches(0.4),
             "Battle-tested across the projects above.",
             font=B_FONT, size=14, italic=True, color=SUB_F)

    cols = rp("stack", [
        ("Agentic AI",
         ["LangChain", "LangGraph", "Multi-Agent", "Ollama", "n8n", "Cursor IDE", "Claude API"]),
        ("LLM & Retrieval",
         ["RAG", "FAISS", "Chroma", "SQLite FTS5", "GPT-4o", "Higgs Audio", "Embeddings"]),
        ("System Design",
         ["Docker", "Kubernetes", "AWS (ECR)", "Spring Boot", ".NET / Azure",
          "GitHub Actions", "Next.js", "FastAPI", "MySQL", "Firebase"]),
    ])
    _scolors = [ACC1, ACC2, ACC1]
    sx = Inches(0.7); sy = Inches(2.8); cw = Inches(4.1); ch = Inches(3.7); gap = Inches(0.15)
    for i, (label, items) in enumerate(cols):
        color = _scolors[i % 3]
        x = sx + i * (cw + gap)
        add_round_rect(s, x, sy, cw, ch, CARD2, radius=0.05, line=cline(color))
        add_rect(s, x, sy, cw, Inches(0.55), STACKHDR)
        lbl_color = ON_DARK if DECO == "memphis" else color
        add_text(s, x + Inches(0.25), sy + Inches(0.1), cw, Inches(0.4),
                 label.upper(), font=H_FONT, size=14, bold=True, color=lbl_color,
                 anchor=MSO_ANCHOR.MIDDLE)
        cx = x + Inches(0.25); cy = sy + Inches(0.75)
        max_x = x + cw - Inches(0.2)
        for it in items:
            chip_w = Inches(0.28 + 0.085 * len(it))
            if cx + chip_w > max_x:
                cx = x + Inches(0.25); cy += Inches(0.45)
            add_round_rect(s, cx, cy, chip_w, Inches(0.34), STACKCHIP_BG, radius=0.4)
            add_text(s, cx, cy + Inches(0.05), chip_w, Inches(0.28),
                     it, font=B_FONT, size=10, bold=True, color=STACKCHIP_FG,
                     align=PP_ALIGN.CENTER)
            cx += chip_w + Inches(0.08)


# ====================================================================
# HONORS & AWARDS  (comprehensive · data-driven · paginated)
# ====================================================================
def build_honors():
    entries = shown(HONORS)
    if not entries:
        return
    PER = 6  # 3 cols x 2 rows
    pages = list(chunk(entries, PER))
    for pi, page in enumerate(pages):
        s = new_slide(("zoom", "fade"))
        title = "Honors & Awards" + (" (cont.)" if pi else "")
        sub = "" if pi else "Funded, ranked, and shipped — the full record."
        section_header(s, "RECOGNITION", title, sub)
        mark_break(s)

        cw = Inches(4.1); ch = Inches(2.05); sx = Inches(0.7); sy = Inches(2.65)
        gap_x = Inches(0.15); gap_y = Inches(0.18)
        for i, h in enumerate(page):
            r, c = i // 3, i % 3
            x = sx + c * (cw + gap_x); y = sy + r * (ch + gap_y)
            color = ACC1 if (pi * PER + i) % 2 == 0 else ACC2
            add_round_rect(s, x, y, cw, ch, CARD, radius=0.05, line=cline(color))
            add_rect(s, x, y, Inches(0.1), ch, color)
            add_text(s, x + Inches(0.3), y + Inches(0.15), cw - Inches(0.5), Inches(0.8),
                     h["big"], font=H_FONT, size=32, bold=True, color=color)
            add_text(s, x + Inches(0.3), y + Inches(1.0), cw - Inches(1.1), Inches(0.4),
                     h["title"], font=H_FONT, size=14, bold=True, color=CHEAD)
            add_text(s, x + Inches(0.3), y + Inches(1.45), cw - Inches(0.5), Inches(0.5),
                     h["desc"], font=B_FONT, size=10, color=CBODY, line_spacing=1.25)
            # year pill, top-right
            if h.get("year"):
                yl = Inches(0.28 + 0.085 * len(h["year"]))
                pill(s, x + cw - yl - Inches(0.25), y + Inches(0.22), h["year"],
                     fill=PILL_BG, fg=PILL_FG, size=10)
            mark_break(s)


# ====================================================================
# THANKS / CONTACT
# ====================================================================
def slide_thanks():
    s = new_slide(("fade", "float_up"))
    if DECO == "memphis":
        fill_slide(s, M_WHITE)
        add_oval(s, Inches(-1.3), Inches(-1.3), Inches(3.2), Inches(3.2), M_CYAN)
        add_oval(s, Inches(11.4), Inches(-1.0), Inches(2.8), Inches(2.8), M_CORAL)
        add_oval(s, Inches(12.3), Inches(6.6), Inches(2.0), Inches(2.0), M_YELLOW)
        add_oval(s, Inches(0.7), Inches(6.4), Inches(0.4), Inches(0.4), M_PURPLE)
        add_oval(s, Inches(2.0), Inches(1.0), Inches(0.3), Inches(0.3), M_PINK)
        add_text(s, Inches(1.0), Inches(1.45), Inches(11), Inches(0.5),
                 "Q & A", font=H_FONT, size=14, bold=True, color=M_CORAL)
        add_text(s, Inches(1.0), Inches(2.05), Inches(12), Inches(1.8),
                 "Thank you.", font=H_FONT, size=80, bold=True, color=M_INK)
        add_rect(s, Inches(1.05), Inches(3.9), Inches(2.6), Inches(0.09), M_CYAN)
        add_text(s, Inches(1.05), Inches(4.15), Inches(11), Inches(0.5),
                 "Happy to dig into any of these in more depth.",
                 font=B_FONT, size=18, italic=True, color=M_GREY_D)
        mark_break(s)
        contacts = [
            ("Email",    "danielkordm@gmail.com",       M_CORAL),
            ("LinkedIn", "linkedin.com/in/danial-kord", M_PURPLE),
            ("GitHub",   "github.com/Danial-Kord",      M_TEAL),
            ("Website",  "danial-kord.github.io",       M_CORAL),
        ]
        cw = Inches(3.0); ch = Inches(1.05); sx = Inches(1.0); sy = Inches(5.1)
        gap_x = Inches(0.1)
        for i, (label, value, color) in enumerate(contacts):
            x = sx + i * (cw + gap_x)
            add_round_rect(s, x, sy, cw, ch, M_CARD, radius=0.08, line=color)
            add_rect(s, x, sy, Inches(0.08), ch, color)
            add_text(s, x + Inches(0.25), sy + Inches(0.18), cw - Inches(0.4), Inches(0.3),
                     label.upper(), font=H_FONT, size=9, bold=True, color=color)
            add_text(s, x + Inches(0.25), sy + Inches(0.5), cw - Inches(0.4), Inches(0.45),
                     value, font=H_FONT, size=12, bold=True, color=M_INK)
        return

    fill_slide(s, NAVY)
    add_rect(s, Inches(0.6), Inches(0.9), Inches(0.12), Inches(5.7), AMBER)
    add_text(s, Inches(1.0), Inches(1.05), Inches(11), Inches(0.5),
             "Q & A", font=H_FONT, size=14, bold=True, color=AMBER)
    add_text(s, Inches(1.0), Inches(1.7), Inches(12), Inches(2.0),
             "Thank you.",
             font=H_FONT, size=88, bold=True, color=IVORY)
    add_text(s, Inches(1.0), Inches(3.85), Inches(11), Inches(0.5),
             "Happy to dig into any of these in more depth.",
             font=B_FONT, size=18, italic=True, color=GOLD)
    mark_break(s)  # — pause before contact cards

    # contact strip
    contacts = [
        ("Email",    "danielkordm@gmail.com",       AMBER),
        ("LinkedIn", "linkedin.com/in/danial-kord", GOLD),
        ("GitHub",   "github.com/Danial-Kord",      AMBER),
        ("Website",  "danial-kord.github.io",       GOLD),
    ]
    cw = Inches(3.0); ch = Inches(1.05); sx = Inches(1.0); sy = Inches(5.2)
    gap_x = Inches(0.1)
    for i, (label, value, color) in enumerate(contacts):
        x = sx + i * (cw + gap_x)
        add_round_rect(s, x, sy, cw, ch, INDIGO, radius=0.08)
        add_rect(s, x, sy, Inches(0.08), ch, color)
        add_text(s, x + Inches(0.25), sy + Inches(0.18), cw - Inches(0.4), Inches(0.3),
                 label.upper(), font=H_FONT, size=9, bold=True, color=color)
        add_text(s, x + Inches(0.25), sy + Inches(0.5), cw - Inches(0.4), Inches(0.45),
                 value, font=H_FONT, size=12, bold=True, color=IVORY)


# ====================================================================
# DIGIHUMAN PIPELINE  (recreated from the template's page 34 — animated)
# ====================================================================
def slide_digihuman_pipeline():
    s = new_slide(("seq",))          # ("seq",) → smooth auto fade build-up
    _no_footer_ids.add(id(s))        # full-bleed diagram: no footer
    fill_slide(s, M_WHITE)

    # decorative template-style blobs (behind content; fade in first)
    add_oval(s, Inches(11.95), Inches(-0.75), Inches(2.1), Inches(2.1), M_PURPLE)
    add_oval(s, Inches(12.55), Inches(3.15), Inches(0.85), Inches(0.85), DH_CORAL)
    add_oval(s, Inches(-0.55), Inches(6.35), Inches(1.5), Inches(1.5), M_CYAN)
    add_oval(s, Inches(0.40), Inches(0.25), Inches(0.26), Inches(0.26), M_YELLOW)

    add_text(s, Inches(0.5), Inches(0.14), Inches(12.33), Inches(0.7),
             "DigiHuman", font=H_FONT, size=34, bold=True, color=DH_INK,
             align=PP_ALIGN.CENTER)

    def tile(x, y, w, h):
        add_round_rect(s, Inches(x), Inches(y), Inches(w), Inches(h), DH_MINT, radius=0.09)

    def pill(x, y, w, h, text, size=12, r=0.5):
        add_round_rect(s, Inches(x), Inches(y), Inches(w), Inches(h), DH_CORAL, radius=r)
        add_text(s, Inches(x), Inches(y), Inches(w), Inches(h), text, font=H_FONT,
                 size=size, bold=True, color=M_WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

    def label(x, y, w, text):
        add_text(s, Inches(x), Inches(y), Inches(w), Inches(0.38), text, font=H_FONT,
                 size=13, bold=True, color=DH_INK, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

    # ----- TOP ROW: Video -> Server -----
    tile(0.93, 0.85, 2.30, 2.45)
    add_image(s, f"{DH_MEDIA}/dh_01.jpg", Inches(1.10), Inches(1.02), w=Inches(1.94))
    label(0.93, 2.86, 2.30, "Video")
    tile(3.75, 0.85, 2.30, 2.45)
    add_image(s, f"{DH_MEDIA}/dh_00.png", Inches(3.93), Inches(1.00), w=Inches(1.95))
    label(3.75, 2.86, 2.30, "Server")
    add_chevron(s, Inches(3.20), Inches(1.55), Inches(0.60), Inches(0.48), DH_ARROW)
    pill(2.50, 3.42, 1.95, 0.46, "Upload")

    # ----- Video Processing -> branch to the three landmark bars -----
    pill(5.71, 1.74, 2.34, 0.55, "Video Processing")
    add_rect(s, Inches(8.05), Inches(1.03), Inches(0.035), Inches(2.05), DH_ARROW)
    bars = [("Facial landmarks", 0.68, "dh_04.png", 10.72, 0.50, 0.89),
            ("Body Pose landmarks", 1.73, "dh_03.png", 10.62, 1.52, 0.94),
            ("Hand landmarks", 2.72, "dh_05.png", 10.61, 2.56, 0.89)]
    for txt, by, icon, ix, iy, isz in bars:
        add_rect(s, Inches(8.05), Inches(by + 0.33), Inches(0.30), Inches(0.035), DH_ARROW)
        add_round_rect(s, Inches(8.24), Inches(by), Inches(3.27), Inches(0.70), DH_MINT, radius=0.3)
        add_text(s, Inches(8.45), Inches(by), Inches(2.05), Inches(0.70), txt,
                 font=H_FONT, size=12, bold=True, color=DH_INK, anchor=MSO_ANCHOR.MIDDLE)
        add_image(s, f"{DH_MEDIA}/{icon}", Inches(ix), Inches(iy), w=Inches(isz))
    pill(8.95, 3.50, 1.95, 0.62, "Receive data from server", size=10, r=0.25)

    # ----- BOTTOM ROW (right -> left): Processing -> Smoothing -> Rendering -----
    tile(8.70, 3.86, 2.55, 2.55)
    add_image(s, f"{DH_MEDIA}/dh_02.png", Inches(9.10), Inches(4.05), w=Inches(1.80))
    pill(8.50, 6.42, 2.75, 0.95,
         "Processing 3D landmarks, calculating rotation & position of body joints",
         size=8.5, r=0.12)
    add_chevron(s, Inches(8.15), Inches(4.92), Inches(0.55), Inches(0.50), DH_ARROW, left=True)
    tile(5.53, 3.86, 2.60, 2.55)
    add_image(s, f"{DH_MEDIA}/dh_06.png", Inches(5.70), Inches(4.00), w=Inches(2.26))
    pill(5.10, 6.42, 3.20, 0.80,
         "Smoothing animation in frames, using signal filters (low-pass filter)",
         size=9, r=0.14)
    add_chevron(s, Inches(4.85), Inches(4.95), Inches(0.55), Inches(0.50), DH_ARROW, left=True)
    tile(2.26, 3.86, 2.55, 2.55)
    add_image(s, f"{DH_MEDIA}/dh_07.png", Inches(2.55), Inches(4.05), w=Inches(1.96))
    pill(2.00, 6.42, 3.00, 0.60, "Rendering the final animation", size=11, r=0.2)


# ====================================================================
# BUILD
# ====================================================================
def build(theme="midnight", role="general", out=None):
    apply_theme(theme)
    apply_role(role)
    _reset_deck()

    # Role profiles may hide slides that aren't relevant to the target role.
    show = {**SHOW_SLIDES, **(rp("show_slides") or {})}

    # ---- Ordered deck. Each section honours its toggle; the comprehensive
    #      sections additionally honour per-entry "show". ----
    slide_title()                                    # always
    if show["about"]:        slide_about()
    if show["pillars"]:      slide_pillars()
    if show["education"]:    build_education()         # all degrees
    if show["experience"]:   build_experience()        # all roles
    if show["vector"]:       slide_vector()            # deep-dive
    if show["latex_cv"]:     slide_latex_cv()          # deep-dive
    if show["dreamforge"]:   slide_dreamforge()        # deep-dive
    if show["guardian"]:     slide_guardian()          # deep-dive
    if show["caselogic"]:    slide_caselogic()         # deep-dive
    if show["safezone"]:     slide_safezone()          # deep-dive
    if show["safezone_arch"]: slide_safezone_arch()    # deep-dive
    if show["digihuman"]:    slide_digihuman_pipeline() # animated pipeline
    if show["projects"]:     build_projects()          # all projects
    if show["stack"]:        slide_stack()
    if show["honors"]:       build_honors()            # all honors
    slide_thanks()                                   # always

    # --- Animations + transitions (BEFORE footers, so footers stay static) ---
    # The page frame (background, decorations, marker, title, subtitle) is
    # static — present on slide enter. Each content box + its text fades in
    # together on a single click.
    for i, slide in enumerate(prs.slides):
        apply_fade_transition(slide, speed="med")
        apply_cascade_anim(slide, dur_ms=400)

    # --- Footers & page numbers: added AFTER animation → never animated.
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        if i == 0 or i == total - 1 or id(slide) in _no_footer_ids:
            continue
        add_footer(slide, i + 1, total)

    out = out or OUT_FILE.get(theme, "Daniel-CV-5min-Talk.pptx")
    sfx = rp("suffix")
    if sfx and out.endswith(".pptx"):
        out = f"{out[:-5]}-{sfx}.pptx"
    prs.save(out)
    print(f"[{theme} | {ROLE}] Saved {out} with {len(prs.slides)} slides")


if __name__ == "__main__":
    theme_arg = (sys.argv[1] if len(sys.argv) > 1 else "midnight").lower()
    role_arg = (sys.argv[2] if len(sys.argv) > 2 else "general").lower()
    themes = (["midnight", "memphis"] if theme_arg == "both"
              else [theme_arg] if theme_arg in ("midnight", "memphis") else None)
    if themes is None or role_arg not in ROLE_PROFILES:
        print("usage: python build_cv_talk.py [midnight|memphis|both] "
              f"[{'|'.join(ROLE_PROFILES)}]")
        sys.exit(1)
    for th in themes:
        build(th, role_arg)
